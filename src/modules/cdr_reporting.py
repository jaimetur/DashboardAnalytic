"""NetCheck CDR report preparation and template-backed PPT rendering."""

from __future__ import annotations

import csv
import gc
import io
import json
import re
import unicodedata
from collections import defaultdict
from dataclasses import dataclass, replace
from io import BytesIO
from itertools import product
from pathlib import Path
from typing import Callable, Iterable

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt

from src.utils.fonts import load_image_font


TEMPLATE_NAMES = {
    "nsa": "Template_CDR_analysis.pptx",
    "sa": "Template_CDR_analysis.pptx",
}
CDR_REPORT_VERSION = "2026-08-31-v11"
REPORTING_KINDS = {"data", "voice", "speech"}
COMMENT_HINTS = ("having ", "observed", "shows ", "similar performance", "worse ", "improvement", "degradation", "gap ")
CATALOG_HEADERS = ("Slide", "Slide Tittle", "Slide Subtittle", "Layout", "Chart Tittle", "CDR source", "KPI", "Chart type", "Filters", "Rows Aggregation", "Column Aggregation", "Legend", "Legend Position")
# Import the two immediately preceding schemas too, so existing templates remain
# usable after the aggregation columns were renamed and the legend was repositioned.
PREVIOUS_CATALOG_HEADERS = ("Slide", "Slide tittle", "Slide Subtittle", "Layout", "Chart Tittle", "CDR source", "KPI", "Chart type", "Legend", "Filters", "Grouping_Rows", "Grouping_Columns", "Legend Position")
OLDER_CATALOG_HEADERS = ("Slide", "Slide tittle", "Slide Subtittle", "Layout", "Chart Tittle", "CDR source", "KPI", "Chart type", "Legend", "Filters", "Grouping_Rows", "Grouping_Columns")
LEGACY_ROWS_COLUMNS_HEADERS = ("Slide", "Slide tittle", "Slide Subtittle", "Layout", "CDR source", "KPI", "Chart type", "Filters", "Grouping_Rows", "Grouping_Columns")
LEGACY_CATALOG_HEADERS = ("Slide", "Slide tittle", "Slide Subtittle", "Layout", "CDR source", "KPI", "Chart type", "Filters", "Grouping")
CATALOG_SOURCE_KINDS = {"cdr-data": "data", "cdr-voice": "voice", "cdr-speech": "speech"}
CHART_TYPES = {
    "100% stacked vertical bars", "count stacked horizontal bars", "cdf line", "scatter", "table",
    "distribution stacked vertical bars", "threshold stacked vertical bars", "average vertical bars", "median vertical bars",
}
STRUCTURAL_SLIDE_TYPES = {"title slide", "transition slide"}
PRESERVED_CHART_TYPES = {"not automated (preserve)"}
FILTER_OPERATORS = ("CONTAINS", "NOT CONTAINS", "IN", "NOT IN", ">=", "<=", "!=", "=", ">", "<")


def _catalogue_header_key(value: str) -> str:
    """Make harmless spelling/case/separator changes in imported headers equivalent."""
    return re.sub(r"[^a-z0-9]+", "", str(value).casefold())


CATALOG_HEADER_ALIASES = {
    "slide": "Slide",
    "slidetittle": "Slide Tittle",
    "slidetitle": "Slide Tittle",
    "slidesubtittle": "Slide Subtittle",
    "slidesubtitle": "Slide Subtittle",
    "layout": "Layout",
    "charttittle": "Chart Tittle",
    "charttitle": "Chart Tittle",
    "cdrsource": "CDR source",
    "kpi": "KPI",
    "charttype": "Chart type",
    "legend": "Legend",
    "legendposition": "Legend Position",
    "filter": "Filters",
    "filters": "Filters",
    "rowsaggregation": "Rows Aggregation",
    "rowaggregation": "Rows Aggregation",
    "groupingrows": "Rows Aggregation",
    "groupingrow": "Rows Aggregation",
    "columnaggregation": "Column Aggregation",
    "columnsaggregation": "Column Aggregation",
    "groupingcolumns": "Column Aggregation",
    "groupingcolumn": "Column Aggregation",
    "grouping": "Grouping",
}


def _canonical_catalog_headers(headers: Iterable[str]) -> tuple[str, ...]:
    """Normalize accepted historic header spellings before schema validation."""
    return tuple(
        CATALOG_HEADER_ALIASES.get(_catalogue_header_key(header), str(header))
        for header in headers
    )


def _default_catalogue_layout(technology: str, chart_count: int) -> str:
    """Choose the standard template layout when a legacy Slides Template omitted it."""
    if technology == "nsa":
        return {
            1: "Title and 1 column + Comments",
            2: "Title and 2 columns + Comments",
            3: "Title and 3 columns + Comments",
        }.get(chart_count, "Title and 2 columns and 2 rows + Comments right")
    return {
        1: "Title and 1 column",
        2: "Title and 2 columns",
        3: "Title and 3 columns",
    }.get(chart_count, "Title and 4 columns")


@dataclass(frozen=True)
class FilterCondition:
    column: str
    operator: str
    values: tuple[str, ...]


@dataclass(frozen=True)
class GroupingSpec:
    dimensions: tuple[str, ...]


@dataclass(frozen=True)
class CatalogEntry:
    slide: int
    slide_title: str
    slide_subtitle: str
    layout: str
    chart_title: str
    cdr_source: str
    kpi: str
    chart_type: str
    legend: str
    filters: str
    grouping_rows: str
    grouping_columns: str
    legend_position: str = "top"

    @property
    def source_kind(self) -> str | None:
        return CATALOG_SOURCE_KINDS.get(self.cdr_source.strip().casefold())

    @property
    def structural_type(self) -> str | None:
        value = self.chart_type.strip().casefold()
        return value if value in STRUCTURAL_SLIDE_TYPES else None


def parse_catalog_filters(value: str) -> tuple[FilterCondition, ...]:
    """Parse `Column OP value; ...` syntax without needing a particular CDR schema."""
    if not value.strip():
        return ()
    conditions: list[FilterCondition] = []
    for clause in (part.strip() for part in value.split(";") if part.strip()):
        match = re.fullmatch(r"(.+?)\s+(NOT\s+CONTAINS|NOT\s+IN|CONTAINS|IN|>=|<=|!=|=|>|<)\s+(.+)", clause, flags=re.I)
        if not match:
            # Compatibility for the initial supplied templates. New templates must use the syntax above.
            continue
        column, operator, raw_values = (part.strip() for part in match.groups())
        operator = re.sub(r"\s+", " ", operator).upper()
        # Existing quality-ratio rows describe both output states as "LQ < 1.6 vs ≥ 1.6".
        # That is chart metadata, not a source-row filter.
        if operator in {">", ">=", "<", "<="} and re.search(r"\bvs\b", raw_values, flags=re.I):
            continue
        if not column:
            raise ValueError(f"Invalid filter '{clause}': a column name is required.")
        if operator in {"IN", "NOT IN"}:
            if not raw_values.startswith("(") or not raw_values.endswith(")"):
                raise ValueError(f"Invalid filter '{clause}': IN values must use parentheses.")
            values = tuple(item.strip() for item in raw_values[1:-1].split(",") if item.strip())
        elif operator in {"CONTAINS", "NOT CONTAINS"}:
            # A contains condition accepts one term or a comma-separated list.
            # Parentheses are the canonical saved form for lists, but accept a
            # plain list as well so manually authored templates remain clear.
            list_values = raw_values[1:-1] if raw_values.startswith("(") and raw_values.endswith(")") else raw_values
            values = tuple(item.strip() for item in list_values.split(",") if item.strip())
        else:
            values = (raw_values,)
        if not values:
            raise ValueError(f"Invalid filter '{clause}': a value is required.")
        conditions.append(FilterCondition(column, operator, values))
    return tuple(conditions)


def parse_catalog_grouping(value: str) -> GroupingSpec:
    if not value.strip():
        return GroupingSpec(())
    dimensions = tuple(part.strip() for part in re.split(r"\s*(?:×|x)\s*", value, flags=re.I) if part.strip())
    if not dimensions:
        raise ValueError("Grouping must contain at least one dimension.")
    return GroupingSpec(dimensions)


def parse_legend_position(value: str) -> str:
    """Return the canonical legend position declared by a Slides Template row."""
    normalized = value.strip().casefold()
    if not normalized:
        return "top"
    aliases = {
        "top": "top", "above": "top", "arriba": "top",
        "bottom": "bottom", "below": "bottom", "abajo": "bottom",
        "left": "left", "izquierda": "left",
        "right": "right", "derecha": "right",
    }
    if normalized not in aliases:
        raise ValueError("Legend Position must be Top, Bottom, Left or Right.")
    return aliases[normalized]


def parse_catalog_csv(content: bytes | str, technology: str) -> list[CatalogEntry]:
    """Validate the editable report-template CSV and return its chart rows."""
    if technology not in TEMPLATE_NAMES:
        raise ValueError("Catalog technology must be NSA or SA.")
    if isinstance(content, bytes):
        try:
            text = content.decode("utf-8-sig")
        except UnicodeDecodeError as exc:
            raise ValueError("The report template must be a UTF-8 CSV file.") from exc
    else:
        text = content
    reader = csv.DictReader(io.StringIO(text))
    fieldnames = tuple(reader.fieldnames or ())
    accepted_schemas = {
        _canonical_catalog_headers(schema)
        for schema in (CATALOG_HEADERS, PREVIOUS_CATALOG_HEADERS, OLDER_CATALOG_HEADERS, LEGACY_ROWS_COLUMNS_HEADERS, LEGACY_CATALOG_HEADERS)
    }
    if _canonical_catalog_headers(fieldnames) not in accepted_schemas:
        raise ValueError("The report template must use exactly these columns: " + ", ".join(CATALOG_HEADERS))
    entries: list[CatalogEntry] = []
    for line_number, row in enumerate(reader, start=2):
        row = {
            CATALOG_HEADER_ALIASES.get(_catalogue_header_key(key), key): value
            for key, value in row.items()
        }
        try:
            slide = int((row.get("Slide") or "").strip())
        except ValueError as exc:
            raise ValueError(f"Catalog row {line_number} has an invalid Slide value.") from exc
        if slide < 1:
            raise ValueError(f"Catalog row {line_number} must use a positive slide number.")
        legacy_grouping = (row.get("Grouping") or "").strip()
        legacy_dimensions = parse_catalog_grouping(legacy_grouping).dimensions
        entry = CatalogEntry(
            slide=slide,
            slide_title=(row.get("Slide Tittle") or "").strip().replace("\\n", "\n"),
            slide_subtitle=(row.get("Slide Subtittle") or "").strip().replace("\\n", "\n"),
            layout=(row.get("Layout") or "").strip(),
            chart_title=(row.get("Chart Tittle") or "").strip().replace("\\n", "\n"),
            cdr_source=(row.get("CDR source") or "").strip(),
            kpi=(row.get("KPI") or "").strip(),
            chart_type=(row.get("Chart type") or "").strip(),
            legend=(row.get("Legend") or "").strip(),
            legend_position=parse_legend_position((row.get("Legend Position") or "").strip()),
            filters=(row.get("Filters") or "").strip(),
            grouping_rows=((row.get("Rows Aggregation") or row.get("Grouping_Rows") or "").strip() or " × ".join(legacy_dimensions[:1])),
            grouping_columns=((row.get("Column Aggregation") or row.get("Grouping_Columns") or "").strip() or " × ".join(legacy_dimensions[1:])),
        )
        if entry.structural_type:
            if not entry.slide_title:
                raise ValueError(f"Catalog row {line_number} requires Slide Tittle for a structural slide.")
            if not entry.layout:
                raise ValueError(f"Catalog row {line_number} requires Layout for a structural slide.")
            structural_chart_fields = (
                entry.chart_title, entry.cdr_source, entry.kpi, entry.legend,
                entry.filters, entry.grouping_rows, entry.grouping_columns,
            )
            if any(value.strip() for value in structural_chart_fields):
                raise ValueError(
                    f"Catalog row {line_number} is a {entry.chart_type} and cannot define chart, CDR, KPI, "
                    "legend, filter or grouping values."
                )
        if entry.source_kind and not entry.slide_title:
            raise ValueError(f"Catalog row {line_number} requires Slide Tittle for a CDR source.")
        if entry.source_kind and not entry.layout:
            raise ValueError(f"Catalog row {line_number} requires Layout for a CDR source.")
        if entry.cdr_source and entry.cdr_source.casefold() not in CATALOG_SOURCE_KINDS:
            raise ValueError(f"Catalog row {line_number} has unsupported CDR source '{entry.cdr_source}'.")
        if not entry.source_kind and not entry.structural_type and entry.chart_type.casefold() not in PRESERVED_CHART_TYPES:
            raise ValueError(
                f"Catalog row {line_number} must define a supported CDR chart, Title Slide or Transition Slide."
            )
        if entry.source_kind and (not entry.kpi or not entry.chart_type):
            raise ValueError(f"Catalog row {line_number} requires KPI and Chart type for a CDR source.")
        if entry.source_kind and entry.chart_type.casefold() not in CHART_TYPES:
            raise ValueError(f"Catalog row {line_number} has unsupported Chart type '{entry.chart_type}'.")
        if entry.source_kind and not (entry.grouping_rows or entry.grouping_columns):
            raise ValueError(f"Catalog row {line_number} requires Rows Aggregation or Column Aggregation for a CDR source.")
        try:
            parse_catalog_filters(entry.filters)
            parse_catalog_grouping(entry.grouping_rows)
            parse_catalog_grouping(entry.grouping_columns)
            parse_legend_position(entry.legend_position)
        except ValueError as exc:
            raise ValueError(f"Catalog row {line_number}: {exc}") from exc
        entries.append(entry)
    if not entries:
        raise ValueError("The report template does not contain any rows.")
    entries_by_slide: dict[int, list[CatalogEntry]] = defaultdict(list)
    for entry in entries:
        entries_by_slide[entry.slide].append(entry)
    for slide_number, slide_entries in entries_by_slide.items():
        structural_entries = [entry for entry in slide_entries if entry.structural_type]
        if structural_entries and len(slide_entries) != 1:
            raise ValueError(
                f"Slide {slide_number} cannot combine a structural Title/Transition row with chart rows."
            )
    return entries


def convert_catalog_csv(content: bytes | str, technology: str) -> bytes:
    """Migrate a compatible legacy CSV into the current editable Slides Templates schema.

    The importer deliberately accepts common title spelling variants and the former
    single ``Grouping`` column.  Missing newer presentation-only columns are left
    blank, while the normal validator still protects required report definitions.
    """
    if isinstance(content, bytes):
        try:
            text = content.decode("utf-8-sig")
        except UnicodeDecodeError as exc:
            raise ValueError("The report template must be a UTF-8 CSV file.") from exc
    else:
        text = content
    reader = csv.DictReader(io.StringIO(text))
    original_headers = tuple(reader.fieldnames or ())
    if not original_headers:
        raise ValueError("The report template does not contain a header row.")

    header_map: dict[str, str] = {}
    for original in original_headers:
        canonical = CATALOG_HEADER_ALIASES.get(_catalogue_header_key(original))
        if canonical and canonical not in header_map:
            header_map[canonical] = original
    if "Slide" not in header_map:
        raise ValueError("The uploaded CSV cannot be converted because it does not contain a Slide column.")

    converted_rows: list[dict[str, str]] = []
    for row in reader:
        converted = {header: "" for header in CATALOG_HEADERS}
        for canonical, original in header_map.items():
            if canonical == "Grouping":
                continue
            converted[canonical] = (row.get(original) or "").strip()
        legacy_grouping = (row.get(header_map.get("Grouping", "")) or "").strip()
        if legacy_grouping:
            dimensions = parse_catalog_grouping(legacy_grouping).dimensions
            if not converted["Rows Aggregation"]:
                converted["Rows Aggregation"] = " × ".join(dimensions[:1])
            if not converted["Column Aggregation"]:
                converted["Column Aggregation"] = " × ".join(dimensions[1:])
        converted_rows.append(converted)

    # Older templates did not contain a Layout column. Its suitable default is
    # determined by the number of automated charts represented by that slide.
    charts_per_slide: dict[str, int] = defaultdict(int)
    for row in converted_rows:
        if row["CDR source"].strip().casefold() in CATALOG_SOURCE_KINDS:
            charts_per_slide[row["Slide"].strip()] += 1
    for row in converted_rows:
        if (
            not row["Layout"].strip()
            and row["CDR source"].strip().casefold() in CATALOG_SOURCE_KINDS
        ):
            row["Layout"] = _default_catalogue_layout(technology, charts_per_slide[row["Slide"].strip()])

    # A slide-less template cannot preserve source slides. Migrate the former
    # preservation marker into an explicit structural slide contract instead.
    for row in converted_rows:
        if row["Chart type"].strip().casefold() not in PRESERVED_CHART_TYPES:
            continue
        is_title = row["Slide"].strip() == "1"
        row["Chart type"] = "Title Slide" if is_title else "Transition Slide"
        row["Layout"] = row["Layout"].strip() or ("Title Page" if is_title else "Title Only")

    output = io.StringIO(newline="")
    writer = csv.DictWriter(output, fieldnames=CATALOG_HEADERS, lineterminator="\n")
    writer.writeheader()
    writer.writerows(converted_rows)
    converted_content = output.getvalue().encode("utf-8")
    # Reuse the regular validator so a conversion never imports an incomplete chart
    # contract just because it used an older set of column headings.
    entries = parse_catalog_csv(converted_content, technology)
    return catalogue_csv(entries)


def load_catalog_csv(path: Path, technology: str) -> list[CatalogEntry]:
    return parse_catalog_csv(path.read_bytes(), technology)


def active_catalog_path(catalog_dir: Path, fallback_catalog: Path, technology: str) -> Path:
    """Return the built-in Slides Template kept in the technology library."""
    return fallback_catalog


def catalogue_csv(entries: list[CatalogEntry]) -> bytes:
    """Serialize the active template using the current editable CSV schema."""
    output = io.StringIO(newline="")
    writer = csv.DictWriter(output, fieldnames=CATALOG_HEADERS, lineterminator="\n")
    writer.writeheader()
    for entry in entries:
        writer.writerow({
            "Slide": entry.slide,
            "Slide Tittle": entry.slide_title.replace("\n", "\\n"),
            "Slide Subtittle": entry.slide_subtitle.replace("\n", "\\n"),
            "Layout": entry.layout,
            "Chart Tittle": entry.chart_title.replace("\n", "\\n"),
            "CDR source": entry.cdr_source,
            "KPI": entry.kpi,
            "Chart type": entry.chart_type,
            "Filters": entry.filters,
            "Rows Aggregation": entry.grouping_rows,
            "Column Aggregation": entry.grouping_columns,
            "Legend": entry.legend,
            "Legend Position": entry.legend_position.title(),
        })
    return output.getvalue().encode("utf-8")


# The PPT templates contain rasterised Tableau charts.  These rules are the
# source-of-truth replacement contract captured from every automated KPI slide:
# source CDR, visual grammar, KPI, and the template's test/session filters.
REPORT_CHART_SPECS = {
    "nsa": {
        8: {"kind": "status_100", "source": "voice", "sessions": ("volte", "multirab", "whatsapp")},
        9: {"kind": "failure_count", "source": "voice", "sessions": ("volte", "multirab", "whatsapp")},
        10: {"kind": "status_100", "source": "data", "tests": ("http", "youtube", "video", "brows")},
        11: {"kind": "status_100", "source": "data", "tests": ("fdfs",)},
        12: {"kind": "cdf_mean", "source": "voice", "metric": ("POLQA_LQ_Avg", "quality_score"), "sessions": ("volte", "multirab")},
        13: {"kind": "cdf_mean", "source": "speech", "metric": ("LQ", "quality_score"), "sessions": ("whatsapp",)},
        14: {"kind": "dual_quality_100", "source": "speech", "metric": ("LQ", "quality_score"), "sessions": ("whatsapp",), "threshold": 1.6,
             "secondary": {"source": "voice", "metric": ("POLQA_LQ_Avg", "quality_score"), "sessions": ("volte",)}},
        15: {"kind": "cdf_mean", "source": "voice", "metric": ("Call_Setup_Time", "setup_time_seconds"), "sessions": ("volte", "multirab")},
        16: {"kind": "cdf_bucket", "source": "data", "metric": ("FDTT_Sustainable_MDR", "Mean_Data_Rate", "throughput_mbps"), "tests": ("fdtt",), "directions": ("dl",)},
        17: {"kind": "cdf_bucket", "source": "data", "metric": ("FDTT_Sustainable_MDR", "Mean_Data_Rate", "throughput_mbps"), "tests": ("fdtt",), "directions": ("ul",)},
        18: {"kind": "cdf_pair", "source": "data", "metric": ("Mean_Data_Rate", "throughput_mbps"), "secondary_metric": ("Data_Test_Duration", "Transfer_Duration", "duration_seconds"), "tests": ("fdfs",), "directions": ("dl",)},
        19: {"kind": "cdf_pair", "source": "data", "metric": ("Mean_Data_Rate", "throughput_mbps"), "secondary_metric": ("Data_Test_Duration", "Transfer_Duration", "duration_seconds"), "tests": ("fdfs",), "directions": ("ul",)},
        20: {"kind": "cdf_pair", "source": "data", "metric": ("Interactivity_RTT_Median", "Interactivity_RTT_AVG", "latency_ms"), "secondary_metric": ("Interactivity_Packet_Error_Ratio", "Packet_Error_Ratio"), "tests": ("interactivity",)},
        21: {"kind": "cdf_mean", "source": "data", "metric": ("http_Browser_1MB_Reached_Duration", "http_Browser_Access_Duration", "setup_time_seconds"), "tests": ("brows", "http")},
    },
    "sa": {
        12: {"kind": "failure_count", "source": "voice", "sessions": ("call", "multirab", "whatsapp")},
        13: {"kind": "failure_count", "source": "voice", "operators": ("vodafone",), "sessions": ("call", "multirab", "whatsapp")},
        14: {"kind": "status_100", "source": "voice", "sessions": ("call", "multirab", "whatsapp")},
        15: {"kind": "quality_100", "source": "speech", "metric": ("LQ", "quality_score"), "threshold": 1.6},
        16: {"kind": "scatter", "source": "speech", "metric": ("LQ", "quality_score"), "x_metric": ("Playing_RSRP_NR_Avg", "NR_RSRP_Avg"), "operators": ("vodafone",), "sessions": ("whatsapp",)},
        17: {"kind": "quality_cdf", "source": "speech", "metric": ("LQ", "quality_score"), "sessions": ("whatsapp",), "operators": ("three", "3 uk", "3") , "threshold": 1.6},
        18: {"kind": "cdf_mean", "source": "voice", "metric": ("POLQA_LQ_Avg", "quality_score"), "sessions": ("call", "multirab", "whatsapp")},
        19: {"kind": "cdf_mean", "source": "voice", "metric": ("POLQA_LQ_Avg", "quality_score"), "sessions": ("call", "multirab", "whatsapp"), "city_scope": "london"},
        21: {"kind": "status_100", "source": "data", "tests": ("fdfs",)},
        22: {"kind": "cdf_mean", "source": "data", "metric": ("Mean_Data_Rate", "throughput_mbps"), "tests": ("fdfs",), "directions": ("dl",)},
        23: {"kind": "cdf_mean", "source": "data", "metric": ("Mean_Data_Rate", "throughput_mbps"), "tests": ("fdfs",), "directions": ("ul",)},
        24: {"kind": "cdf_mean", "source": "data", "metric": ("FDTT_Sustainable_MDR", "Mean_Data_Rate", "throughput_mbps"), "tests": ("fdtt",)},
        25: {"kind": "cdf_mean", "source": "data", "metric": ("Interactivity_RTT_Median", "Interactivity_RTT_AVG", "latency_ms"), "tests": ("interactivity",)},
        26: {"kind": "cdf_mean", "source": "data", "metric": ("http_Browser_1MB_Reached_Duration", "http_Browser_Access_Duration", "setup_time_seconds"), "tests": ("brows", "http")},
    },
}

OPERATOR_COLORS = {
    "vodafone": "#E15759", "vf": "#E15759",
    "three": "#F28E2B", "3": "#F28E2B",
    "o2": "#4E79A7", "telefonica": "#4E79A7",
    "ee": "#76B7B2",
}

OPERATOR_COLOUR_VARIANTS = {
    # Deliberately high-contrast shades: adjacent vendors must remain
    # distinguishable when they belong to the same operator family.
    "#E15759": ("#E15759", "#9B1D20", "#F58B8E", "#A03A6B"),
    "#F28E2B": ("#F28E2B", "#A84B09", "#FFC145", "#C96500"),
    "#4E79A7": ("#4E79A7", "#123B68", "#63A4E8", "#365A9B"),
    "#76B7B2": ("#76B7B2", "#176B70", "#3CC4BD", "#368A8F"),
}

NEUTRAL_SERIES_COLORS = ("#6F42C1", "#2E8B57", "#A0612A", "#C23B8B", "#607D8B", "#B8860B")

VENDOR_COLOUR_VARIANTS = {
    # A vendor remains recognisable in every chart. Different operators using
    # the same vendor receive contrasting shades from that vendor's family.
    "ericsson": ("#2E8B57", "#0D5A34", "#65B984", "#176B42"),
    "huawei": ("#E15759", "#A61E2B", "#F58B8E", "#C43D4D"),
    "samsung": ("#D9A514", "#9A7000", "#F2CC5C", "#C28D00"),
    "nsn": ("#4E79A7", "#123B68", "#63A4E8", "#365A9B"),
}


@dataclass(frozen=True)
class ReportSelection:
    data_id: int
    voice_id: int
    speech_id: int
    technology: str
    multivendor: bool
    vodafone_mapping_id: int | None = None
    three_mapping_id: int | None = None


def _normalise_operator(value: object) -> str:
    """Return the operator form used by the Vendor-mapping business formula."""
    text = str(value or "").strip()
    key = re.sub(r"[^a-z0-9]+", "", text.casefold())
    if key in {"3", "3uk", "three", "threeuk"}:
        return "3"
    if key in {"vodafone", "vodafoneuk", "vf", "vfuk"}:
        return "Vodafone UK"
    return text


def _normalise_report_operator(value: object) -> str:
    """Return the canonical report label for known historical UK aliases."""
    text = _normalise_operator(value)
    key = re.sub(
        r"[^a-z0-9]+",
        "",
        unicodedata.normalize("NFKD", text.casefold()).encode("ascii", "ignore").decode("ascii"),
    )
    if key in {"vodafone", "vodafoneuk", "vf", "vfuk"}:
        return "VF"
    if key.startswith("o2") or key in {"telefonica", "telefonicao2"}:
        return "O2"
    if key in {"ee", "eeuk", "everythingeverywhere"}:
        return "EE"
    return text


def normalise_report_operator_aliases(frame: pd.DataFrame) -> pd.DataFrame:
    """Canonicalise report-only operator labels without mutating stored CDRs.

    The normalisation applies to the physical ``Operator`` field used by
    template filters and to legacy operator-only values in ``report_vendor``.
    Mapped values such as ``Vodafone_Ericsson`` are intentionally left intact.
    """
    result = frame.copy()
    for column in result.columns:
        if _normalise_catalog_name(str(column)) == "operator":
            result[column] = result[column].map(_normalise_report_operator)
    if "report_vendor" in result.columns:
        result["report_vendor"] = result["report_vendor"].map(_normalise_report_operator)
    return result


def _split_global_cells(value: object) -> list[str]:
    """Return ordered Global CI values from NetCheck's timeline/list formatting."""
    if value is None or pd.isna(value):
        return []
    text = str(value).strip()
    if not text or text.lower() in {"nan", "none", "null"}:
        return []
    parts = re.split(r"\s*(?:->|;|,|\||\]\s*\[)\s*", text.strip("[] "))
    return [_canonical_cell_id(part) for part in parts if part.strip(" []'\"")]


def _canonical_cell_id(value: object) -> str:
    text = str(value or "").strip(" []'\"")
    if re.fullmatch(r"\d+\.0+", text):
        return text.split('.', 1)[0]
    return text


def vendor_from_cells(operator: object, cells: object, vendor_lookup: dict[str, str]) -> str:
    """Implement the business formula supplied for Vodafone and Three.

    Vodafone's Ericsson/null exceptions deliberately resolve to Mixed Vendor,
    exactly as specified by the reference formula.
    """
    normalized_operator = _normalise_operator(operator)
    if normalized_operator not in {"Vodafone UK", "3"}:
        return normalized_operator
    global_cells = _split_global_cells(cells)
    first = vendor_lookup.get(global_cells[0]) if global_cells else None
    last = vendor_lookup.get(global_cells[-1]) if global_cells else None
    if normalized_operator == "Vodafone UK":
        if first and first == last:
            return f"Vodafone_{first}"
        if (first == "Ericsson" and last != "Ericsson") or (last == "Ericsson" and first != "Ericsson"):
            return "Vodafone_Mixed Vendor"
        return "Vodafone_Other Vendor"
    if first and first == last:
        return f"3_{first}"
    return "3_Mixed Vendor"


def _first_existing(df: pd.DataFrame, candidates: Iterable[str]) -> str | None:
    # NetCheck exports are not entirely consistent between releases: a field can
    # be written as ``Cell_ID_A``, ``CELL ID A`` or ``cell-id-a``.  Treat only
    # spelling separators and case as insignificant, while retaining the
    # original column name for the caller.
    normalise = lambda value: re.sub(r"[^a-z0-9]+", "", str(value).casefold())
    lookup = {normalise(column): str(column) for column in df.columns}
    for candidate in candidates:
        actual = lookup.get(normalise(candidate))
        if actual:
            return actual
    return None


def classify_sessions(df: pd.DataFrame, technology: str) -> pd.DataFrame:
    main_rat_column = _first_existing(df, ["RAT", "RAT_A"])
    sample_rat_column = _first_existing(df, ["Sample_RAT_A"])
    normalized_rat_column = _first_existing(df, ["technology_primary"])
    rat_column = main_rat_column or sample_rat_column or normalized_rat_column
    call_mode_columns = list(dict.fromkeys(filter(None, (
        _first_existing(df, ["L1_Call_Mode_A", "L1_call_Mode_A"]),
        _first_existing(df, ["L2_Call_Mode_A", "L2_call_Mode_A"]),
    ))))
    if not rat_column and not call_mode_columns:
        raise ValueError("The selected CDR does not contain RAT or Call Mode fields required to separate NSA and SA sessions.")

    # WhatsApp's explicit RAT describes the packet session and remains
    # authoritative. Native and MultiRAB calls are classified by their call
    # mode first: a plain LTE RAT is common for valid VoLTE calls and must not
    # make those samples disappear from an NSA report.
    rat_values = (
        df[rat_column].fillna("").astype(str).str.strip()
        if rat_column else pd.Series("", index=df.index, dtype="string")
    )
    call_modes = pd.Series("", index=df.index, dtype="string")
    for column in call_mode_columns:
        call_modes = call_modes.str.cat(df[column].fillna("").astype(str), sep=" ")
    session_column = _first_existing(df, ["Session_Type", "session_type"])
    session_values = (
        df[session_column].fillna("").astype(str)
        if session_column else pd.Series("", index=df.index, dtype="string")
    )
    whatsapp = session_values.str.contains("whatsapp", case=False, na=False)
    recognised_call_mode = call_modes.str.contains(r"VoLTE|EPSFB|VoNR", case=False, na=False, regex=True)

    rat_nsa = rat_values.str.contains(r"EN[- ]?DC", case=False, na=False, regex=True)
    rat_sa = rat_values.str.contains(r"NR", case=False, na=False, regex=True) & ~rat_nsa
    mode_nsa = call_modes.str.contains(r"VoLTE|EPSFB", case=False, na=False, regex=True)
    mode_sa = call_modes.str.contains(r"VoNR", case=False, na=False, regex=True)
    conflicting_call_mode = mode_nsa & mode_sa
    mode_nsa &= ~conflicting_call_mode
    mode_sa &= ~conflicting_call_mode
    multirab_lte_fallback = (
        session_values.str.contains("multirab", case=False, na=False)
        & rat_values.str.contains(r"(?:^|/)LTE(?:/|$)", case=False, na=False, regex=True)
    )

    if technology == "nsa":
        mask = (whatsapp & rat_nsa) | (~whatsapp & (mode_nsa | (~recognised_call_mode & (rat_nsa | multirab_lte_fallback))))
    else:
        mask = (whatsapp & rat_sa) | (~whatsapp & (mode_sa | (~recognised_call_mode & rat_sa)))
    return df[mask].copy()


def _integer_cell_component(value: object) -> int | None:
    try:
        if value is None or pd.isna(value):
            return None
        return int(float(value))
    except (TypeError, ValueError):
        return None


def build_three_vendor_lookup(mapping: pd.DataFrame) -> dict[str, str]:
    cell_column = _first_existing(mapping, ["Cid__ECI", "CId___ECI"])
    vendor_column = _first_existing(mapping, ["Vendor", "OP/ Vendor", "OP_Vendor"])
    if not cell_column or not vendor_column:
        raise ValueError("The selected 3UK mapping must contain Cid__ECI and a Vendor column.")
    lookup: dict[str, str] = {}
    for cell, vendor in mapping[[cell_column, vendor_column]].dropna().itertuples(index=False):
        cell_text, vendor_text = _canonical_cell_id(cell), str(vendor).strip()
        if cell_text and vendor_text:
            lookup[cell_text] = vendor_text
    return lookup


def build_vodafone_vendor_lookup(mapping: pd.DataFrame) -> dict[str, str]:
    source_sheet = _first_existing(mapping, ["source_sheet"])
    four_g_mapping = mapping
    if source_sheet:
        four_g_mapping = mapping[mapping[source_sheet].astype(str).str.strip().str.casefold() == "4g"].copy()
    enodeb_column = _first_existing(four_g_mapping, ["eNodeB ID", "eNodeB_ID"])
    local_cell_column = _first_existing(four_g_mapping, ["Local Cell ID", "Local_Cell_ID"])
    vendor_column = _first_existing(four_g_mapping, ["OP/ Vendor", "OP_Vendor", "Vendor"])
    if four_g_mapping.empty or not enodeb_column or not local_cell_column or not vendor_column:
        raise ValueError("The selected VFUK mapping must contain the 4G sheet with eNodeB ID, Local Cell ID and OP/ Vendor columns.")
    lookup: dict[str, str] = {}
    for enodeb, local_cell, vendor in four_g_mapping[[enodeb_column, local_cell_column, vendor_column]].itertuples(index=False):
        enodeb_id = _integer_cell_component(enodeb)
        local_cell_id = _integer_cell_component(local_cell)
        vendor_text = str(vendor).strip()
        if enodeb_id is None or local_cell_id is None or not vendor_text or vendor_text.lower() == "nan":
            continue
        # Equivalent to Excel: HEX2DEC(DEC2HEX(eNodeB ID, 5) & DEC2HEX(Local Cell ID, 2)).
        if not 0 <= enodeb_id <= 0xFFFFF or not 0 <= local_cell_id <= 0xFF:
            continue
        lookup[str((enodeb_id << 8) | local_cell_id)] = vendor_text
    return lookup


def enrich_multivendor(df: pd.DataFrame, vodafone_mapping: pd.DataFrame, three_mapping: pd.DataFrame) -> pd.DataFrame:
    result = df.copy()
    operator_column = _first_existing(result, ["operator", "Operator"])
    cell_column = _first_existing(result, [
        "Cell_ID_A", "Cell_IDs_A", "Cell_ID", "Cell ID A", "Cell IDs A",
        "Global_Cell_ID_A", "Global_Cell_ID", "Global CI", "Global_CI",
        "GCID", "GCI", "CGI", "ECI", "Serving_Cell_ID",
    ])
    if not operator_column or not cell_column:
        raise ValueError("The selected CDR must contain Operator and a supported Cell ID field for multivendor reporting.")
    vodafone_lookup = build_vodafone_vendor_lookup(vodafone_mapping)
    three_lookup = build_three_vendor_lookup(three_mapping)
    result["report_vendor"] = [
        vendor_from_cells(
            operator,
            cells,
            vodafone_lookup if _normalise_operator(operator) == "Vodafone UK" else three_lookup,
        )
        for operator, cells in result[[operator_column, cell_column]].itertuples(index=False)
    ]
    return result


def assign_cdr_vendors(
    df: pd.DataFrame,
    vodafone_mapping: pd.DataFrame | None = None,
    three_mapping: pd.DataFrame | None = None,
) -> pd.DataFrame:
    """Assign the agreed multivendor value to the normalized CDR ``vendor`` field.

    This is the persistent Workspace counterpart of report-time enrichment.  The
    same first/last ``Cell_ID_A`` and operator formula is used, while allowing a
    CDR to be mapped with the VFUK and/or 3UK mapping files available in its
    Workspace.
    """
    result = df.copy()
    # CDR workbooks can already expose a source ``Vendor`` column.  Dataset
    # normalisation also creates a lower-case vendor field, which SQLite keeps
    # as ``vendor__2`` to avoid a case-insensitive name collision.  That source
    # field is not the calculated multivendor result, so remove every such
    # collision before writing the canonical, user-facing ``vendor`` column.
    vendor_collision_columns = [
        str(column) for column in result.columns
        if str(column).casefold() == 'vendor' or re.fullmatch(r'vendor__\d+', str(column).casefold())
    ]
    if vendor_collision_columns:
        result = result.drop(columns=vendor_collision_columns)
    operator_column = _first_existing(result, ["operator", "Operator"])
    cell_column = _first_existing(result, [
        "Cell_ID_A", "Cell_IDs_A", "Cell_ID", "Cell ID A", "Cell IDs A",
        "Global_Cell_ID_A", "Global_Cell_ID", "Global CI", "Global_CI",
        "GCID", "GCI", "CGI", "ECI", "Serving_Cell_ID",
    ])
    if not operator_column or not cell_column:
        raise ValueError(
            "The selected CDR must contain Operator and a Cell ID field. "
            "Supported names include Cell_ID_A, Cell_IDs_A, Cell_ID, Global CI, GCID, GCI, CGI or ECI."
        )

    vodafone_lookup = build_vodafone_vendor_lookup(vodafone_mapping) if vodafone_mapping is not None else {}
    three_lookup = build_three_vendor_lookup(three_mapping) if three_mapping is not None else {}
    assigned_vendors: list[object] = []
    report_groups: list[str] = []
    for operator, cells in result[[operator_column, cell_column]].itertuples(index=False):
        normalized_operator = _normalise_operator(operator)
        if normalized_operator == "Vodafone UK":
            if vodafone_mapping is None:
                assigned_vendors.append(pd.NA)
                report_groups.append(normalized_operator)
            else:
                mapped_value = vendor_from_cells(operator, cells, vodafone_lookup)
                assigned_vendors.append(mapped_value)
                report_groups.append(mapped_value)
        elif normalized_operator == "3":
            if three_mapping is None:
                assigned_vendors.append(pd.NA)
                report_groups.append(normalized_operator)
            else:
                mapped_value = vendor_from_cells(operator, cells, three_lookup)
                assigned_vendors.append(mapped_value)
                report_groups.append(mapped_value)
        else:
            # O2/EE are operators without a multivendor mapping.  They are not
            # vendors and must not be materialised in the Vendor field.  Their
            # operator remains available in the report-only comparison grouping,
            # exactly as the final ELSE branch of the supplied formula requires.
            assigned_vendors.append(pd.NA)
            report_groups.append(normalized_operator)
    result["vendor"] = assigned_vendors
    result["report_vendor"] = report_groups
    # Keep the calculated Vendor immediately after the source-sheet identifier
    # (or first when no worksheet identifier exists).  ``report_vendor`` is an
    # internal reporting comparison field and deliberately remains last.
    leading_columns = [column for column in ('source_sheet', 'vendor') if column in result.columns]
    remaining_columns = [column for column in result.columns if column not in {*leading_columns, 'report_vendor'}]
    return result.loc[:, [*leading_columns, *remaining_columns, 'report_vendor']]


def ensure_report_vendor_group(df: pd.DataFrame) -> pd.DataFrame:
    """Provide the formula-compatible comparison field for mapped CDRs.

    This also upgrades CDRs mapped before the report_vendor field was
    materialised: mapped Vodafone/3UK Vendor values are used where present,
    while O2/EE keep their operator as required by the formula's final ELSE.
    """
    result = df.copy()
    if "report_vendor" in result.columns:
        return result
    operator_column = _first_existing(result, ["operator", "Operator"])
    vendor_column = _first_existing(result, ["vendor", "Vendor"])
    if not operator_column:
        return result
    operators = result[operator_column].fillna("").astype(str).str.strip()
    if vendor_column:
        vendors = result[vendor_column].fillna("").astype(str).str.strip()
        result["report_vendor"] = vendors.where(vendors.ne(""), operators)
    else:
        result["report_vendor"] = operators
    return result


def _replace_word(value: str, source: str, replacement: str) -> str:
    """Replace a template word while retaining the source word's casing."""
    def replace_match(match: re.Match[str]) -> str:
        word = match.group(0)
        if word.isupper():
            return replacement.upper()
        if word[0].isupper():
            return replacement.capitalize()
        return replacement.lower()
    return re.sub(rf"\b{re.escape(source)}\b", replace_match, value, flags=re.I)


def _replace_operator_label(value: str, replacement: str) -> str:
    """Replace singular and plural display labels without touching filters."""
    return _replace_word(_replace_word(value, "Operators", f"{replacement}s"), "Operator", replacement)


def prepare_multivendor_catalog_entry(entry: CatalogEntry) -> CatalogEntry:
    """Apply the report-only multivendor wording and grouping interpretation.

    The stored template remains an operator-oriented definition.  For a
    multivendor run, only grouping dimensions, display legends and titles are
    transformed.  Filters deliberately remain untouched so an ``Operator``
    condition continues to filter the physical CDR Operator column.
    """
    def vendor_grouping(value: str) -> str:
        dimensions = parse_catalog_grouping(value).dimensions
        return " × ".join(
            "Vendor" if _normalise_catalog_name(dimension) == "operator" else dimension
            for dimension in dimensions
        )

    return replace(
        entry,
        slide_title=_replace_operator_label(entry.slide_title, "Vendor"),
        slide_subtitle=_replace_operator_label(entry.slide_subtitle, "Vendor"),
        chart_title=_replace_operator_label(entry.chart_title, "Vendor"),
        legend=_replace_operator_label(entry.legend, "Campaign"),
        grouping_rows=vendor_grouping(entry.grouping_rows),
        grouping_columns=vendor_grouping(entry.grouping_columns),
    )


def _font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    return load_image_font(size, bold)


def _column(frame: pd.DataFrame, candidates: Iterable[str]) -> str | None:
    return _first_existing(frame, candidates)


def _period_column(frame: pd.DataFrame) -> str | None:
    return _column(frame, ("Campaign", "period", "Period", "Quarter"))


def _group_column(frame: pd.DataFrame, multivendor: bool) -> str | None:
    return "report_vendor" if multivendor and "report_vendor" in frame.columns else _column(frame, ("Operator", "operator"))


def _normalise_catalog_name(value: str) -> str:
    return re.sub(r"[^a-z0-9]", "", value.casefold())


def _catalog_column(
    frame: pd.DataFrame,
    name: str,
    multivendor: bool,
    metric: str | None = None,
    bucket_edges: list[float] | None = None,
    operator_as_vendor: bool = True,
) -> str | None:
    """Resolve a template field name against a source column or supported semantic dimension."""
    normalized = _normalise_catalog_name(name)
    if normalized == "operator":
        return _group_column(frame, multivendor) if operator_as_vendor else _column(frame, ("Operator", "operator"))
    if normalized in {"vendor", "reportvendor"}:
        return _group_column(frame, multivendor) if multivendor else _column(frame, ("Vendor", "vendor"))
    if normalized == "callfamily":
        materialized = _column(frame, ("Call Family", "Call_Family", "call_family"))
        if materialized:
            return materialized
        session_column = _column(frame, ("Session_Type", "session_type"))
        if not session_column:
            return _column(frame, ("Call_Family", "call_family"))
        call_mode = _column(frame, ("L1_Call_Mode_A", "L1_Call_Mode_B", "Call_Mode", "call_mode"))
        session = frame[session_column].fillna("").astype(str)
        family = pd.Series("CALL", index=frame.index, dtype="string")
        family.loc[session.str.contains("multirab", case=False, na=False)] = "MultiRAB"
        family.loc[session.str.contains("whatsapp", case=False, na=False)] = "WhatsApp"
        family.loc[session.str.contains("volte", case=False, na=False)] = "VoLTE"
        family.loc[session.str.contains("vonr", case=False, na=False)] = "VoNR"
        if call_mode:
            modes = frame[call_mode].fillna("").astype(str)
            family.loc[(family == "CALL") & modes.str.contains("volte", case=False, na=False)] = "VoLTE"
            family.loc[(family == "CALL") & modes.str.contains("vonr", case=False, na=False)] = "VoNR"
        frame["__catalog_call_family"] = family
        return "__catalog_call_family"
    if normalized == "testfamily":
        materialized = _column(frame, ("Test Family", "Test_Family", "test_family"))
        if materialized:
            return materialized
        type_column = _column(frame, ("Type_of_Test", "Test_Type", "test_type"))
        name_column = _column(frame, ("Test_Name", "test_name"))
        if not type_column and not name_column:
            return None
        test_family = frame[type_column].fillna("").astype(str) if type_column else pd.Series("", index=frame.index, dtype="string")
        if name_column:
            test_names = frame[name_column].fillna("").astype(str)
            test_family.loc[test_names.str.contains("youtube", case=False, na=False)] = "YouTube"
            test_family.loc[test_names.str.contains("fdfs", case=False, na=False)] = "FDFS"
            test_family.loc[test_names.str.contains("fdtt", case=False, na=False)] = "FDTT"
        frame["__catalog_test_family"] = test_family
        return "__catalog_test_family"
    aliases = {
        "campaign": ("Campaign", "period", "Period", "Quarter"),
        "city": ("City", "city", "G_Level_1", "G_Level_2"),
        "failuretechnology": ("Failure_Technology",),
        "failurecategory": ("Failure_Category",),
        "n rband": ("NR_Band", "NR band", "Band_NR"),
        "nrband": ("NR_Band", "NR band", "Band_NR"),
        "lteband": ("LTE_Band", "4G_Band", "Band_LTE"),
        "radioband": ("NR_Band", "LTE_Band", "Band", "Radio_Band"),
    }
    if normalized in {"ratebucket", "valuebucket"}:
        if not metric:
            return None
        numeric = pd.to_numeric(frame[metric], errors="coerce")
        edges = bucket_edges or [1, 5, 10, 25, 50]
        labels = [f"<{edges[0]:g}"] + [f"{low:g}-{high:g}" for low, high in zip(edges, edges[1:])] + [f"{edges[-1]:g}+"]
        frame["__catalog_rate_bucket"] = pd.cut(numeric, bins=[float("-inf"), *edges, float("inf")], labels=labels, right=False).astype("string")
        return "__catalog_rate_bucket"
    candidate = _column(frame, aliases.get(normalized, (name,)))
    if candidate:
        return candidate
    for column in frame.columns:
        if _normalise_catalog_name(str(column)) == normalized:
            return str(column)
    return None


def _latest_campaign_value(series: pd.Series) -> str | None:
    """Return the latest year/quarter campaign independently of its text format."""
    values = [value for value in series.dropna().astype(str).str.strip().unique() if value]
    if not values:
        return None
    return max(values, key=_campaign_sort_key)


def _campaign_sort_key(value: object) -> tuple[int, int, str]:
    """Sort campaign values chronologically when their year/quarter is present."""
    text = str(value).strip()
    year_match = re.search(r"(?:19|20)\d{2}", text)
    quarter_match = re.search(r"(?:^|[^A-Z0-9])Q\s*([1-4])(?:[^0-9]|$)", text, flags=re.I)
    return (
        int(year_match.group(0)) if year_match else -1,
        int(quarter_match.group(1)) if quarter_match else -1,
        text.casefold(),
    )


def _campaign_display_value(value: object) -> str:
    """Reduce NetCheck campaign identifiers to a stable year/quarter label."""
    text = str(value).strip()
    year_match = re.search(r"(?:19|20)\d{2}", text)
    quarter_match = re.search(r"(?:^|[^A-Z0-9])Q\s*([1-4])(?:[^0-9]|$)", text, flags=re.I)
    if year_match and quarter_match:
        return f"{year_match.group(0)}-Q{quarter_match.group(1)}"
    return text


def _apply_catalog_filters(frame: pd.DataFrame, entry: CatalogEntry, multivendor: bool, metric: str | None) -> pd.DataFrame:
    result = frame.copy()
    for condition in parse_catalog_filters(entry.filters):
        if _normalise_catalog_name(condition.column) in {"threshold", "buckets"}:
            continue
        # Operator filters always refer to the source CDR Operator field. Only
        # grouping dimensions are promoted to Vendor for multivendor reports.
        column = _catalog_column(result, condition.column, multivendor, metric, operator_as_vendor=False)
        if not column:
            raise ValueError(f"Slide {entry.slide}: filter column '{condition.column}' does not exist in {entry.cdr_source}.")
        series = result[column]
        is_operator_filter = _normalise_catalog_name(condition.column) == "operator"
        comparison_series = series.map(_normalise_report_operator) if is_operator_filter else series
        if condition.operator in {">", ">=", "<", "<=", "=", "!="}:
            target = condition.values[0]
            if (
                _normalise_catalog_name(condition.column) == "campaign"
                and condition.operator in {"=", "!="}
                and _normalise_catalog_name(target) in {"latest", "latestcampaign"}
            ):
                latest_campaign = _latest_campaign_value(series)
                comparison = series.astype(str).eq(latest_campaign) if latest_campaign is not None else pd.Series(False, index=series.index)
                if condition.operator == "!=":
                    comparison = ~comparison
                result = result.loc[comparison].copy()
                continue
            if is_operator_filter:
                target = _normalise_report_operator(target)
            numeric = pd.to_numeric(series, errors="coerce")
            target_number = pd.to_numeric(pd.Series([target]), errors="coerce").iloc[0]
            if pd.notna(target_number):
                comparison = {">": numeric > target_number, ">=": numeric >= target_number, "<": numeric < target_number, "<=": numeric <= target_number, "=": numeric == target_number, "!=": numeric != target_number}[condition.operator]
            else:
                comparison = {"=": comparison_series.astype(str).str.casefold() == target.casefold(), "!=": comparison_series.astype(str).str.casefold() != target.casefold()}.get(condition.operator)
                if comparison is None:
                    raise ValueError(f"Slide {entry.slide}: '{condition.operator}' requires a numeric value for '{condition.column}'.")
        elif condition.operator in {"CONTAINS", "NOT CONTAINS"}:
            targets = tuple(_normalise_report_operator(item) if is_operator_filter else item for item in condition.values)
            comparison = pd.Series(False, index=series.index)
            for target in targets:
                comparison |= comparison_series.astype(str).str.contains(target, case=False, na=False, regex=False)
            if condition.operator == "NOT CONTAINS":
                comparison = ~comparison
        else:  # IN / NOT IN
            accepted = {
                (_normalise_report_operator(item) if is_operator_filter else item).casefold()
                for item in condition.values
            }
            comparison = comparison_series.astype(str).str.casefold().isin(accepted)
            if condition.operator == "NOT IN":
                comparison = ~comparison
        result = result.loc[comparison].copy()
    return result


def _catalog_threshold(entry: CatalogEntry) -> float:
    for condition in parse_catalog_filters(entry.filters):
        if _normalise_catalog_name(condition.column) == "threshold":
            try:
                return float(condition.values[0])
            except ValueError as exc:
                raise ValueError(f"Slide {entry.slide}: Threshold must be numeric.") from exc
    legacy = re.search(r"(?:<|<=)\s*([0-9]+(?:\.[0-9]+)?)\s+vs", entry.filters, flags=re.I)
    return float(legacy.group(1)) if legacy else 1.6


def _catalog_bucket_edges(entry: CatalogEntry) -> list[float] | None:
    for condition in parse_catalog_filters(entry.filters):
        if _normalise_catalog_name(condition.column) == "buckets":
            try:
                edges = [float(value.strip()) for value in condition.values[0].split(",")]
            except ValueError as exc:
                raise ValueError(f"Slide {entry.slide}: Buckets must be a comma-separated numeric list.") from exc
            if len(edges) < 1 or edges != sorted(set(edges)):
                raise ValueError(f"Slide {entry.slide}: Buckets must be unique ascending values.")
            return edges
    return None


def _apply_catalog_grouping(frame: pd.DataFrame, entry: CatalogEntry, multivendor: bool, metric: str | None) -> tuple[pd.DataFrame, str, str]:
    row_spec = parse_catalog_grouping(entry.grouping_rows)
    column_spec = parse_catalog_grouping(entry.grouping_columns)
    bucket_edges = _catalog_bucket_edges(entry)
    def resolve_dimensions(dimensions: tuple[str, ...], axis: str) -> list[str]:
        resolved: list[str] = []
        for dimension in dimensions:
            column = _catalog_column(frame, dimension, multivendor, metric, bucket_edges)
            if not column:
                raise ValueError(f"Slide {entry.slide}: {axis} grouping dimension '{dimension}' does not exist in {entry.cdr_source}.")
            resolved.append(column)
        return resolved

    row_columns = resolve_dimensions(row_spec.dimensions, "row")
    column_columns = resolve_dimensions(column_spec.dimensions, "column")
    explicit_dimension_values = {
        _normalise_catalog_name(condition.column): list(condition.values)
        for condition in parse_catalog_filters(entry.filters)
        if condition.operator == "IN"
    }
    configured_dimension_values: dict[str, list[str]] = {}
    # Preserve every resolved hierarchy level for renderers that need pane-like
    # rows and nested column headers. The flattened primary/series fields remain
    # available for chart grammars that intentionally use compact labels.
    def materialise_dimension(dimension: str, column: str, target: str) -> None:
        values = frame[column].fillna("(blank)").astype(str)
        normalized_dimension = _normalise_catalog_name(dimension)
        if normalized_dimension == "campaign":
            values = values.map(_campaign_display_value)
        frame[target] = values
        requested = explicit_dimension_values.get(normalized_dimension)
        if requested:
            if normalized_dimension == "campaign":
                requested = sorted(
                    dict.fromkeys(_campaign_display_value(value) for value in requested),
                    key=_campaign_sort_key,
                )
            elif normalized_dimension == "operator":
                requested = [_normalise_report_operator(value) for value in requested]
            configured_dimension_values[target] = list(dict.fromkeys(requested))

    row_display_columns: list[str] = []
    for index, (dimension, column) in enumerate(zip(row_spec.dimensions, row_columns, strict=True)):
        target = f"__catalog_row_{index}"
        materialise_dimension(dimension, column, target)
        row_display_columns.append(target)
    column_display_columns: list[str] = []
    for index, (dimension, column) in enumerate(zip(column_spec.dimensions, column_columns, strict=True)):
        target = f"__catalog_column_{index}"
        materialise_dimension(dimension, column, target)
        column_display_columns.append(target)
    # Rows form the category/table-row hierarchy. Columns form chart series and
    # table columns. Distribution charts reserve the final column level as the
    # stack/bucket breakdown and use any preceding column levels as the series.
    primary = "__catalog_primary"
    series = "__catalog_series"
    def materialise(columns: list[str], target: str) -> None:
        if not columns:
            frame[target] = "(all)"
        elif len(columns) == 1:
            frame[target] = frame[columns[0]].fillna("(blank)").astype(str)
        else:
            # On an empty frame pandas returns an empty DataFrame here rather
            # than an empty Series. Preserve a valid column so no-data chart
            # states do not interrupt the whole report.
            if frame.empty:
                frame[target] = pd.Series(index=frame.index, dtype="object")
            else:
                frame[target] = frame[columns].fillna("(blank)").astype(str).agg(" · ".join, axis=1)

    materialise(row_display_columns, primary)
    is_distribution = entry.chart_type.casefold() == "distribution stacked vertical bars"
    if is_distribution and column_display_columns:
        stack_column = column_display_columns[-1]
        series_columns = column_display_columns[:-1]
        materialise(series_columns, series)
        frame["__catalog_stack"] = frame[stack_column].fillna("(blank)").astype(str)
    else:
        materialise(column_display_columns, series)
    frame.attrs["catalogue_dimension_values"] = configured_dimension_values
    frame.attrs["catalogue_dimension_labels"] = {
        **{column: dimension for column, dimension in zip(row_display_columns, row_spec.dimensions, strict=True)},
        **{column: dimension for column, dimension in zip(column_display_columns, column_spec.dimensions, strict=True)},
        primary: row_spec.dimensions,
        series: column_spec.dimensions,
    }
    # Source tables are commonly appended newest campaign first. Whenever
    # Campaign participates in either aggregation axis, establish one shared
    # hierarchy order here so every chart grammar sees oldest -> newest while
    # retaining the first-seen order of all non-temporal dimensions.
    hierarchy = [
        *zip(row_display_columns, row_spec.dimensions, strict=True),
        *zip(column_display_columns, column_spec.dimensions, strict=True),
    ]
    campaign_columns = [
        column for column, dimension in hierarchy
        if _normalise_catalog_name(dimension) == "campaign"
    ]
    needs_campaign_sort = any(
        (observed := frame[column].drop_duplicates().tolist()) != sorted(observed, key=_campaign_sort_key)
        for column in campaign_columns
    ) if not frame.empty else False
    if needs_campaign_sort:
        sort_columns: list[str] = []
        for index, (column, dimension) in enumerate(hierarchy):
            values = frame[column].drop_duplicates().tolist()
            if _normalise_catalog_name(dimension) == "campaign":
                values = sorted(values, key=_campaign_sort_key)
            ranks = {value: rank for rank, value in enumerate(values)}
            sort_column = f"__catalog_sort_{index}"
            frame[sort_column] = frame[column].map(ranks)
            sort_columns.append(sort_column)
        preserved_attrs = frame.attrs.copy()
        frame = frame.sort_values(sort_columns, kind="stable").drop(columns=sort_columns)
        frame.attrs = preserved_attrs
    return frame, primary, series


def preview_catalog_chart_data(
    frame: pd.DataFrame,
    entry: CatalogEntry,
    *,
    limit: int = 200,
    offset: int = 0,
    column_filters: dict[str, tuple[str, ...]] | None = None,
) -> tuple[pd.DataFrame, dict[str, object]]:
    """Return the exact, post-filter rows supplied to one template chart.

    This deliberately shares the reporting filter and grouping resolution so
    the editor can expose derived dimensions (for example ``Rate Bucket``)
    without asking users to reproduce report logic by hand.
    """
    if not entry.source_kind:
        raise ValueError('Only chart rows with a CDR source can be previewed.')
    spec = _catalog_spec(entry)
    metric = _metric_column(frame, spec)
    filtered = _apply_catalog_filters(frame, entry, False, metric)
    grouped, primary, series = _apply_catalog_grouping(filtered, entry, False, metric)
    display_columns: list[tuple[str, str]] = []

    def include(source: str | None, label: str) -> None:
        if source and source in grouped.columns:
            display_columns.append((source, label))

    # Include the physical/effective CDR fields referenced by every part of
    # the chart definition, not only the synthetic labels used by renderers.
    bucket_edges = _catalog_bucket_edges(entry)
    for condition in parse_catalog_filters(entry.filters):
        if _normalise_catalog_name(condition.column) not in {'threshold', 'buckets'}:
            include(_catalog_column(grouped, condition.column, False, metric, bucket_edges, operator_as_vendor=False), f'Filter · {condition.column}')
    for axis, grouping in (('Rows Aggregation', entry.grouping_rows), ('Column Aggregation', entry.grouping_columns)):
        for dimension in parse_catalog_grouping(grouping).dimensions:
            include(_catalog_column(grouped, dimension, False, metric, bucket_edges), f'{axis} · {dimension}')
    include(metric, f'KPI · {metric}' if metric else 'KPI')
    include(primary, 'Resolved Rows Aggregation')
    include(series, 'Resolved Column Aggregation')
    include('__catalog_stack', 'Resolved Stack Aggregation')

    # Legend contains source dimensions selected by the template/editor.
    legend_dimensions = ', '.join(_legend_dimensions(entry.legend)) or '(automatic)'
    if display_columns:
        full_result = pd.concat(
            [grouped[source].rename(label) for source, label in display_columns],
            axis=1,
        ).copy()
    else:
        full_result = grouped.copy()
    full_result.insert(len(full_result.columns), 'Legend dimensions', legend_dimensions)
    filter_values = {
        str(column): sorted(
            {'' if pd.isna(value) else str(value) for value in full_result[column].tolist()},
            key=str.casefold,
        )
        for column in full_result.columns
    }
    for column, values in (column_filters or {}).items():
        if column not in full_result.columns or not values:
            continue
        accepted = {str(value) for value in values}
        full_result = full_result[full_result[column].map(lambda value: '' if pd.isna(value) else str(value)).isin(accepted)]
    # Callers that render a table still request small pages.  The template-editor
    # preview also uses this helper to build one temporary, server-side result so
    # later page changes do not need to recalculate the chart filters.
    page_size = max(1, min(int(limit), 100_000))
    page_offset = max(0, int(offset))
    result = full_result.iloc[page_offset:page_offset + page_size].copy()
    return result, {
        'source_rows': len(frame.index),
        'matched_rows': len(filtered.index),
        'shown_rows': len(result.index),
        'visible_rows': len(full_result.index),
        'page_offset': page_offset,
        'columns': list(result.columns),
        'filter_values': filter_values,
    }


def render_catalog_chart_preview(frame: pd.DataFrame, entry: CatalogEntry, *, multivendor: bool = False) -> bytes:
    """Render the same PNG chart used by a report for editor/report previews."""
    if not entry.source_kind:
        raise ValueError('Only chart rows with a CDR source can be previewed.')
    render_entry = prepare_multivendor_catalog_entry(entry) if multivendor else entry
    normalised_frame = normalise_report_operator_aliases(frame)
    return _chart_for_catalog_entry(render_entry, {render_entry.source_kind: normalised_frame}, multivendor).getvalue()


def is_empty_catalog_chart(image: bytes, entry: CatalogEntry) -> bool:
    """Identify the intentional no-samples placeholder emitted by the renderer."""
    title = entry.chart_title or entry.slide_title
    return image == _empty_chart(title).getvalue()


def _matches(frame: pd.DataFrame, column: str | None, tokens: tuple[str, ...] | None) -> pd.Series:
    if not column or not tokens:
        return pd.Series(True, index=frame.index)
    pattern = "|".join(re.escape(token) for token in tokens)
    return frame[column].astype(str).str.contains(pattern, case=False, na=False, regex=True)


def _source_for_spec(frames: dict[str, pd.DataFrame], spec: dict, multivendor: bool) -> tuple[pd.DataFrame, str | None, str | None]:
    frame = frames[spec["source"]].copy()
    session_column = _column(frame, ("Session_Type", "session_type", "Test_Name", "Test_Type"))
    test_column = _column(frame, ("Test_Name", "test_name", "Type_of_Test", "Test_Type"))
    direction_column = _column(frame, ("Direction", "direction", "Call_Direction"))
    operator_column = _group_column(frame, multivendor)
    mask = _matches(frame, session_column, spec.get("sessions"))
    mask &= _matches(frame, test_column, spec.get("tests"))
    mask &= _matches(frame, direction_column, spec.get("directions"))
    if spec.get("operators"):
        mask &= _matches(frame, operator_column, spec["operators"])
    if spec.get("city_scope"):
        city_column = _column(frame, ("city", "City", "G_Level_1", "G_Level_2"))
        mask &= _matches(frame, city_column, (spec["city_scope"],))
    return frame.loc[mask].copy(), operator_column, _period_column(frame)


def _metric_column(frame: pd.DataFrame, spec: dict) -> str | None:
    return _column(frame, spec.get("metric", ()))


def _operator_colour(label: object) -> str | None:
    """Return the stable palette colour for a recognised UK operator/vendor."""
    normalized = str(label).strip().casefold()
    if "vodafone" in normalized or re.search(r"(?:^|[^a-z0-9])vf(?:$|[^a-z0-9])", normalized):
        return OPERATOR_COLORS["vodafone"]
    if "three" in normalized or re.search(r"(?:^|[^a-z0-9])3(?:$|[^a-z0-9])", normalized):
        return OPERATOR_COLORS["three"]
    if "telefonica" in normalized or re.search(r"(?:^|[^a-z0-9])o2(?:$|[^a-z0-9])", normalized):
        return OPERATOR_COLORS["o2"]
    if re.search(r"(?:^|[^a-z0-9])ee(?:$|[^a-z0-9])", normalized):
        return OPERATOR_COLORS["ee"]
    return None


def _colour(label: object, index: int = 0) -> str:
    """Return a neutral series colour unless a renderer opts into a palette."""
    return NEUTRAL_SERIES_COLORS[index % len(NEUTRAL_SERIES_COLORS)]


def _hierarchy_group_colours(keys: list[tuple[object, ...]], level: int = 0) -> dict[str, str]:
    """Colour one hierarchy level consistently, with readable variants."""
    colours: dict[str, str] = {}
    offsets: dict[str, int] = {}
    for index, key in enumerate(keys):
        group = str(key[level]) if len(key) > level else ""
        if group in colours:
            continue
        base = _operator_colour(group)
        if base and base in OPERATOR_COLOUR_VARIANTS:
            offset = offsets.get(base, 0)
            colours[group] = OPERATOR_COLOUR_VARIANTS[base][offset % len(OPERATOR_COLOUR_VARIANTS[base])]
            offsets[base] = offset + 1
        else:
            colours[group] = _colour(group, index)
    return colours


def _dimension_roles(frame: pd.DataFrame, axis_columns: list[str]) -> list[set[str]]:
    """Identify the declared Operator/Vendor roles behind materialised axes."""
    definitions = frame.attrs.get("catalogue_dimension_labels", {})
    roles: list[set[str]] = []
    for column in axis_columns:
        definition = definitions.get(column, column)
        labels = definition if isinstance(definition, (tuple, list)) else (definition,)
        normalised = {_normalise_catalog_name(str(label)) for label in labels}
        roles.append({role for role in ("operator", "vendor") if any(role in label for label in normalised)})
    return roles


def _vendor_label(value: object) -> str:
    """Extract the vendor portion from an Operator_Vendor-style label."""
    parts = [part.strip() for part in re.split(r"[_·|/]", str(value)) if part.strip()]
    non_operators = [part for part in parts if _operator_colour(part) is None]
    return non_operators[-1] if non_operators else (parts[-1] if parts else str(value).strip())


def _vendor_colour_family(vendor: str) -> str | None:
    normalized = vendor.casefold()
    return next((family for family in VENDOR_COLOUR_VARIANTS if family in normalized), None)


def _series_colours(
    keys: list[tuple[object, ...]],
    axis_columns: list[str],
    frame: pd.DataFrame,
    *,
    line_chart: bool = False,
) -> dict[tuple[object, ...], str]:
    """Choose colours from declared dimensions rather than label coincidences.

    Vendor is the stable visual identity for bars, points, stacks and lines.
    Recognised vendors use their own colour family regardless of operator;
    repeated vendors across operators receive distinct shades in that family.
    """
    if not keys:
        return {}
    roles = _dimension_roles(frame, axis_columns)
    operator_levels = [index for index, role in enumerate(roles) if "operator" in role]
    vendor_levels = [index for index, role in enumerate(roles) if "vendor" in role]
    identity_levels = list(dict.fromkeys([*operator_levels, *vendor_levels]))
    operator_level = next(
        (index for index in identity_levels if any(_operator_colour(key[index]) for key in keys if len(key) > index)),
        None,
    )
    operator_for_key = {
        key: next(
            (str(key[index]) for index in identity_levels if len(key) > index and _operator_colour(key[index])),
            "",
        )
        for key in keys
    }
    operator_values = {colour for key in keys if (colour := _operator_colour(operator_for_key[key]))}

    if vendor_levels:
        vendor_level = vendor_levels[0]
        vendor_keys = [_vendor_label(key[vendor_level]) if len(key) > vendor_level else "" for key in keys]
        vendor_colours: dict[tuple[str, str], str] = {}
        family_offsets: dict[str, int] = {}
        neutral_index = 0
        for key, vendor in zip(keys, vendor_keys, strict=True):
            operator = operator_for_key[key].casefold()
            identity = (vendor.casefold(), operator)
            if identity in vendor_colours:
                continue
            family = _vendor_colour_family(vendor)
            if family:
                variants = VENDOR_COLOUR_VARIANTS[family]
                offset = family_offsets.get(family, 0)
                vendor_colours[identity] = variants[offset % len(variants)]
                family_offsets[family] = offset + 1
            else:
                vendor_colours[identity] = _colour(vendor, neutral_index)
                neutral_index += 1
        return {
            key: vendor_colours[(vendor.casefold(), operator_for_key[key].casefold())]
            for key, vendor in zip(keys, vendor_keys, strict=True)
        }

    if operator_level is not None and (len(operator_values) > 1 or line_chart):
        vendor_level = vendor_levels[0] if vendor_levels else operator_level
        palette_keys = [
            f"{operator_for_key[key]} · {_vendor_label(key[vendor_level])}" if len(key) > vendor_level else operator_for_key[key]
            for key in keys
        ]
        palette = _hierarchy_group_colours([(value,) for value in palette_keys])
        return {key: palette[palette_key] for key, palette_key in zip(keys, palette_keys, strict=True)}

    return {key: _colour("", index) for index, key in enumerate(keys)}


def _operator_hierarchy_level(keys: list[tuple[object, ...]], axis_columns: list[str]) -> int:
    """Find the column-aggregation level carrying operator/vendor identity."""
    candidates = [
        index for index, column in enumerate(axis_columns)
        if column.startswith("__catalog_column_")
    ]
    if not candidates:
        return 0
    # Templates may declare Campaign before Operator/Vendor. Prefer whichever
    # column level actually carries recognised operator identity, while keeping
    # the declared order as the deterministic fallback.
    return max(
        candidates,
        key=lambda index: (sum(_operator_colour(key[index]) is not None for key in keys if len(key) > index), -index),
    )


def _operator_vendor_key(*labels: object) -> str:
    """Choose the operator/vendor portion of a composite chart label."""
    for label in labels:
        for part in re.split(r"\s*·\s*", str(label)):
            if _operator_colour(part):
                return part.strip()
    return str(labels[0]).strip() if labels else ""


def _legend_dimensions(value: str) -> tuple[str, ...]:
    """Return the ordered CDR dimensions selected in the Legend field."""
    if "/" in value:
        return ()
    return tuple(dimension.strip() for dimension in value.split(",") if dimension.strip())


def _legend_labels(value: str) -> tuple[str, ...]:
    """Return legacy/manual legend captions separated with slashes."""
    if "/" not in value:
        return ()
    return tuple(label.strip() for label in value.split("/") if label.strip())


def _legend_key_caption(
    key: tuple[object, ...],
    axis_columns: list[str],
    frame: pd.DataFrame,
    dimensions: tuple[str, ...],
) -> str:
    """Build a caption from the selected dimensions represented by an axis key."""
    requested = {_normalise_catalog_name(value) for value in dimensions}
    labels = frame.attrs.get("catalogue_dimension_labels", {})
    selected_parts: list[str] = []
    for index, column in enumerate(axis_columns):
        declared = labels.get(column, column)
        declared_names = declared if isinstance(declared, tuple) else (declared,)
        if any(_normalise_catalog_name(str(name)) in requested for name in declared_names):
            selected_parts.append(str(key[index]))
    parts = selected_parts or [str(value) for value in key if str(value) != "(all)"]
    return " · ".join(parts) or "(all)"


def _legend_caption(labels: tuple[str, ...], index: int, fallback: object) -> str:
    return labels[index] if index < len(labels) else str(fallback)


def _draw_chart_legend(
    draw: ImageDraw.ImageDraw,
    items: list[tuple[str, str, int]],
    position: str,
    *,
    font_size: int = 15,
    line_markers: bool = False,
    side_x: int | None = None,
) -> None:
    """Draw a template legend in a row (top/bottom) or column (left/right)."""
    position = parse_legend_position(position)
    # Legends are part of the chart, not ancillary metadata. Enforce a
    # readable minimum because the 1600px PNG is normally scaled down in the
    # report and preview viewers.
    font_size = max(font_size, 17)
    marker_size = 22
    horizontal = position in {"top", "bottom"}
    if horizontal:
        # Dense CDFs can legitimately contain one curve per full hierarchy
        # combination.  Fit their compact line legend inside the 1600×900
        # canvas instead of letting the last rows fall below its lower edge.
        columns = 6 if line_markers else 5
        row_height = font_size + 12
        rows = max(1, (len(items) + columns - 1) // columns)
        start_x = 100
        start_y = 80 if position == "top" else 900 - (rows * row_height) - 8
        for index, (caption, colour, width) in enumerate(items):
            x = start_x + (index % columns) * (220 if line_markers else 275)
            y = start_y + (index // columns) * row_height
            if line_markers:
                draw.line((x, y + 11, x + 34, y + 11), fill=colour, width=max(width, 3))
            else:
                draw.rectangle((x, y, x + marker_size, y + marker_size), fill=colour)
            draw.text((x + (43 if line_markers else 32), y - 1), caption[:28], fill="#263B4A", font=_font(font_size, True))
        return
    x, start_y = (side_x if side_x is not None else (26 if position == "left" else 1380)), 112
    for index, (caption, colour, width) in enumerate(items):
        y = start_y + index * (font_size + 14)
        if line_markers:
            draw.line((x, y + 11, x + 34, y + 11), fill=colour, width=max(width, 3))
        else:
            draw.rectangle((x, y, x + marker_size, y + marker_size), fill=colour)
        draw.text((x + (43 if line_markers else 32), y - 1), caption[:24], fill="#263B4A", font=_font(font_size, True))


def _catalogue_display_label(category: object, series: object) -> str:
    """Keep row categories visible while using column groups as comparisons."""
    category_text, series_text = str(category), str(series)
    return category_text if series_text == "(all)" else f"{category_text} · {series_text}"


def _hierarchical_unique_keys(frame: pd.DataFrame, columns: list[str]) -> list[tuple[object, ...]]:
    """Keep every child grouping value adjacent to its parent.

    Source CDRs are commonly appended campaign by campaign.  Their natural row
    order can therefore be ``Vodafone/2025, O2/2025, Vodafone/2026, O2/2026``.
    That order would split the visible Vodafone group into two separate chart
    sections.  Retain the first-seen order of each level, but traverse the
    hierarchy depth-first so all campaign bars for one operator stay together.
    """
    keys = list(frame[columns].drop_duplicates().itertuples(index=False, name=None))
    if len(columns) < 2:
        return keys

    ordered: list[tuple[object, ...]] = []

    def visit(values: list[tuple[object, ...]], level: int) -> None:
        if level == len(columns) - 1:
            ordered.extend(values)
            return
        groups: dict[object, list[tuple[object, ...]]] = {}
        for key in values:
            groups.setdefault(key[level], []).append(key)
        for children in groups.values():
            visit(children, level + 1)

    visit(keys, 0)
    return ordered


def _hierarchical_complete_keys(frame: pd.DataFrame, columns: list[str]) -> list[tuple[object, ...]]:
    """Return the complete, first-seen Cartesian grid for hierarchy columns.

    A comparison axis must remain stable even where a source has no samples
    for one child value.  For example, every Operator/Vendor retains both
    Campaign columns even if it has no failures in one campaign.
    """
    if not columns:
        return [()]
    configured_values = frame.attrs.get("catalogue_dimension_values", {})
    values_by_level = [
        list(configured_values.get(column) or frame[column].drop_duplicates())
        for column in columns
    ]
    if any(not values for values in values_by_level):
        return []
    return [tuple(key) for key in product(*values_by_level)]


def _canvas(title: str) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    draw.text((32, 20), title, fill="#1D3345", font=_font(40, True))
    return image, draw


def _text_width(draw: ImageDraw.ImageDraw, value: str, font: ImageFont.ImageFont) -> int:
    """Return the rendered width of a label, including the actual font metrics."""
    return int(draw.textbbox((0, 0), value, font=font)[2])


def _draw_rotated_label(
    image: Image.Image,
    value: str,
    *,
    centre_x: float,
    bottom_y: float,
    fill: str,
    font: ImageFont.ImageFont,
) -> None:
    """Draw a 45° label centred on ``centre_x`` with its bottom at ``bottom_y``.

    Axis captions are deliberately rotated only when their measured text no
    longer fits the space assigned to a group.  Rendering to a transparent
    layer avoids clipping the anti-aliased glyphs on the white chart canvas.
    """
    probe = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    bbox = probe.textbbox((0, 0), value, font=font)
    label = Image.new("RGBA", (bbox[2] - bbox[0] + 8, bbox[3] - bbox[1] + 8), (0, 0, 0, 0))
    label_draw = ImageDraw.Draw(label)
    label_draw.text((4 - bbox[0], 4 - bbox[1]), value, fill=fill, font=font)
    rotated = label.rotate(45, expand=True, resample=Image.Resampling.BICUBIC)
    image.paste(rotated, (round(centre_x - rotated.width / 2), round(bottom_y - rotated.height)), rotated)


def _draw_vertical_label(
    image: Image.Image,
    value: str,
    *,
    centre_x: float,
    centre_y: float,
    fill: str,
    font: ImageFont.ImageFont,
) -> None:
    """Draw a y-axis caption centred vertically and rotated by 90 degrees."""
    probe = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    bbox = probe.textbbox((0, 0), value, font=font)
    label = Image.new("RGBA", (bbox[2] - bbox[0] + 8, bbox[3] - bbox[1] + 8), (0, 0, 0, 0))
    ImageDraw.Draw(label).text((4 - bbox[0], 4 - bbox[1]), value, fill=fill, font=font)
    rotated = label.rotate(90, expand=True, resample=Image.Resampling.BICUBIC)
    image.paste(rotated, (round(centre_x - rotated.width / 2), round(centre_y - rotated.height / 2)), rotated)


def _draw_dashed_vertical_line(
    draw: ImageDraw.ImageDraw,
    x: float,
    top: float,
    bottom: float,
    *,
    fill: str = "#AEBBC4",
    dash: int = 8,
    gap: int = 6,
) -> None:
    """Draw a light vertical group separator without obscuring bar values."""
    y = top
    while y < bottom:
        draw.line((x, y, x, min(y + dash, bottom)), fill=fill, width=1)
        y += dash + gap


def _draw_dashed_horizontal_line(
    draw: ImageDraw.ImageDraw,
    y: float,
    left: float,
    right: float,
    *,
    fill: str = "#AEBBC4",
    dash: int = 8,
    gap: int = 6,
) -> None:
    """Draw a light horizontal separator between child row values."""
    x = left
    while x < right:
        draw.line((x, y, min(x + dash, right), y), fill=fill, width=1)
        x += dash + gap


def _draw_top_column_group_separators(
    draw: ImageDraw.ImageDraw,
    keys: list[tuple[object, ...]],
    axis_columns: list[str],
    *,
    left: float,
    width: float,
    top: float,
    bottom: float,
) -> None:
    """Separate only adjacent values of the first column aggregation level.

    Row aggregations may precede the column hierarchy in ``keys``.  The first
    ``__catalog_column_*`` field is therefore the only level that defines a
    visual group boundary; lower column levels (for example Campaign) must
    remain together without their own vertical separators.
    """
    if len(keys) < 2:
        return
    top_column_level = next(
        (index for index, column in enumerate(axis_columns) if column.startswith("__catalog_column_")),
        None,
    )
    if top_column_level is None:
        return
    values = [str(key[top_column_level]) if len(key) > top_column_level else "" for key in keys]
    for index in range(1, len(values)):
        if values[index] != values[index - 1]:
            _draw_dashed_vertical_line(draw, left + index * width / len(keys), top, bottom)


def _empty_chart(title: str) -> BytesIO:
    image, draw = _canvas(title)
    draw.text((50, 440), "No valid samples for this KPI and technology filter", fill="#61727D", font=_font(24))
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0)
    return output


def _render_status_100(title: str, frame: pd.DataFrame, group: str | None, period: str | None, quality: bool = False, threshold: float = 1.6, metric: str | None = None, legend_labels: tuple[str, ...] = (), legend_position: str = "top") -> BytesIO:
    if frame.empty or not group or not period:
        return _empty_chart(title)
    image, draw = _canvas(title)
    state_column = metric if quality else _column(frame, ("Call_Status", "Test_Result", "status"))
    if not state_column:
        return _empty_chart(title)
    states = ("< 1.6", "≥ 1.6") if quality else ("Completed", "Dropped", "Failed")
    colours = ("#E15759", "#59A14F") if quality else ("#4E79A7", "#F28E2B", "#E15759")
    hierarchy_columns = sorted(
        [column for column in frame.columns if column.startswith("__catalog_row_") or column.startswith("__catalog_column_")],
        key=lambda column: (0 if column.startswith("__catalog_row_") else 1, int(column.rsplit("_", 1)[1])),
    )
    data = frame[[group, period, state_column, *hierarchy_columns]].copy()
    if quality:
        numeric_state = pd.to_numeric(data[state_column], errors="coerce")
        # Missing and non-numeric KPI values are not samples and therefore
        # must not affect either segment of a percentage chart.
        data = data.loc[numeric_state.notna()].copy()
        numeric_state = numeric_state.loc[data.index]
        data["state"] = numeric_state.map(lambda value: "< 1.6" if value < threshold else "≥ 1.6")
    else:
        raw = data[state_column].astype("string").str.strip().str.casefold()
        # SQLite/pandas sources can represent absent values as NULL, NaN,
        # pandas <NA>, or an empty string. Exclude all of them before assigning
        # status categories so they cannot be counted as Failed by fallback.
        valid_state = raw.notna() & raw.ne("") & ~raw.isin({"nan", "<na>"})
        data = data.loc[valid_state].copy()
        raw = raw.loc[data.index]
        data["state"] = raw.map(lambda value: "Completed" if any(item in value for item in ("complete", "success", "pass", "ok")) else "Dropped" if "drop" in value else "Failed")
    data = data.dropna(subset=[group, period])
    row_hierarchy = [column for column in hierarchy_columns if column.startswith("__catalog_row_")]
    column_hierarchy = [column for column in hierarchy_columns if column.startswith("__catalog_column_")]
    if column_hierarchy:
        return _render_status_100_hierarchy(title, data, row_hierarchy, column_hierarchy, states, colours, legend_labels, legend_position)
    if len(row_hierarchy) > 1:
        return _render_status_100_hierarchy(title, data, [], row_hierarchy, states, colours, legend_labels, legend_position)
    combos = [(str(g), str(p)) for g, p in data[[group, period]].drop_duplicates().itertuples(index=False)]
    if not combos:
        return _empty_chart(title)
    chart_left, chart_top, chart_width, chart_height = 145, 115, 1300, 640
    bar_width = max(20, min(72, chart_width // max(len(combos) * 2, 1)))
    for i, (g, p) in enumerate(combos):
        subset = data[(data[group].astype(str) == g) & (data[period].astype(str) == p)]
        total = max(len(subset), 1); x = chart_left + i * (chart_width / len(combos)) + 12
        running = 0
        for state, colour in zip(states, colours, strict=True):
            value = len(subset[subset["state"] == state]) / total
            height = value * chart_height
            y = chart_top + chart_height - running - height
            draw.rectangle((x, y, x + bar_width, y + height), fill=colour)
            if value >= .08: draw.text((x + 2, y + height / 2 - 8), f"{value:.0%}", fill="white", font=_font(16, True))
            running += height
        label = _catalogue_display_label(g, p)[:24]
        label_font = _font(18, True)
        if _text_width(draw, label, label_font) > chart_width / len(combos) - 8:
            _draw_rotated_label(image, label, centre_x=x + bar_width / 2, bottom_y=chart_top + chart_height + 75, fill="#5A6B78", font=label_font)
        else:
            draw.text((x - 4, chart_top + chart_height + 8), label, fill="#5A6B78", font=label_font)
    for y in range(0, 101, 20):
        value_y = chart_top + chart_height - (y / 100 * chart_height)
        draw.line((chart_left - 20, value_y, chart_left + chart_width, value_y), fill="#E4E9ED", width=1)
        draw.text((64, value_y - 10), f"{y}%", fill="#4E6271", font=_font(18, True))
    _draw_chart_legend(draw, [(_legend_caption(legend_labels, index, state), colour, 2) for index, (state, colour) in enumerate(zip(states, colours, strict=True))], legend_position, font_size=16)
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0); return output


def _render_status_100_hierarchy(
    title: str,
    data: pd.DataFrame,
    row_hierarchy: list[str],
    column_hierarchy: list[str],
    states: tuple[str, ...],
    colours: tuple[str, ...],
    legend_labels: tuple[str, ...] = (),
    legend_position: str = "top",
) -> BytesIO:
    """Render template row groups as panes and column groups as nested headers."""
    row_keys = _hierarchical_unique_keys(data, row_hierarchy) if row_hierarchy else [()]
    column_keys = _hierarchical_unique_keys(data, column_hierarchy)
    if not row_keys or not column_keys:
        return _empty_chart(title)

    image, draw = _canvas(title)
    # Reserve a dedicated header band below the chart title.  Rotated vendor /
    # operator captions can be tall, so they must never share the title area.
    row_labels = [" · ".join(str(value) for value in row_key) for row_key in row_keys]
    row_label_font = _font(18, True)
    widest_row_label = max((_text_width(draw, label, row_label_font) for label in row_labels), default=0)
    # Leave a measured gutter for row labels, percentage ticks and a visual
    # gap before the plot.  Fixed gutters caused long FDFS/test-family labels
    # to cross the y axis on multi-pane charts.
    chart_left = max(205, min(540, 24 + widest_row_label + 92))
    chart_top, chart_right, chart_height = 245, 1395, 510
    chart_width = chart_right - chart_left
    row_height = chart_height / len(row_keys)
    column_width = chart_width / len(column_keys)
    bar_width = max(18, min(86, column_width * 0.68))

    # The first column level is the upper header (Operator in the template
    # contract); lower levels, such as Campaign, are shown below every column.
    outer_values = [str(key[0]) for key in column_keys]
    outer_groups: list[tuple[int, int, str]] = []
    start = 0
    while start < len(column_keys):
        end = start + 1
        while end < len(column_keys) and outer_values[end] == outer_values[start]:
            end += 1
        outer_groups.append((start, end, outer_values[start]))
        start = end
    header_font = _font(16, True)
    rotate_outer_headers = any(
        _text_width(draw, caption[:22], header_font) + 14 > (end - start) * column_width
        for start, end, caption in outer_groups
    )
    lower_captions = [" · ".join(str(value) for value in key[1:]) or str(key[0]) for key in column_keys]
    rotate_axis_captions = rotate_outer_headers or any(
        _text_width(draw, caption[:20], _font(14)) + 8 > column_width
        for caption in lower_captions
    )
    for start, end, caption in outer_groups:
        centre = chart_left + ((start + end) / 2) * column_width
        caption = caption[:22]
        if rotate_axis_captions:
            _draw_rotated_label(image, caption, centre_x=centre, bottom_y=chart_top - 16, fill="#566A78", font=header_font)
        else:
            draw.text((centre - min(len(caption) * 4, 70), chart_top - 46), caption, fill="#566A78", font=header_font)
        draw.line((chart_left + start * column_width, chart_top - 14, chart_left + end * column_width, chart_top - 14), fill="#CDD7DE", width=1)

    _draw_top_column_group_separators(
        draw,
        column_keys,
        column_hierarchy,
        left=chart_left,
        width=chart_width,
        top=chart_top,
        bottom=chart_top + chart_height,
    )

    for row_index, row_key in enumerate(row_keys):
        pane_top = chart_top + row_index * row_height
        pane_bottom = pane_top + row_height
        row_label = row_labels[row_index]
        if row_label:
            draw.text((24, pane_top + row_height / 2 - 10), row_label, fill="#566A78", font=row_label_font)
        draw.line((24, pane_bottom, chart_left + chart_width, pane_bottom), fill="#D7DEE3", width=1)
        ticks = (0, 50, 100) if row_index == len(row_keys) - 1 else (50, 100)
        for tick in ticks:
            tick_y = pane_bottom - tick / 100 * row_height
            draw.line((chart_left, tick_y, chart_left + chart_width, tick_y), fill="#E8ECEF", width=1)
            draw.text((chart_left - 50, tick_y - 9), f"{tick}%", fill="#566A78", font=_font(15, True))

        row_mask = pd.Series(True, index=data.index)
        for field, value in zip(row_hierarchy, row_key, strict=True):
            row_mask &= data[field].astype(str).eq(str(value))
        for column_index, column_key in enumerate(column_keys):
            mask = row_mask.copy()
            for field, value in zip(column_hierarchy, column_key, strict=True):
                mask &= data[field].astype(str).eq(str(value))
            subset = data.loc[mask]
            if subset.empty:
                # A grouping grid intentionally retains every column so its
                # headers remain comparable between rows.  Mark absent source
                # combinations explicitly instead of leaving misleading blank
                # space (for example, a vendor with no MultiRAB samples).
                centre_x = chart_left + (column_index + .5) * column_width
                draw.text((centre_x - 5, pane_top + row_height / 2 - 7), "—", fill="#B5C0C8", font=_font(14))
                continue
            x = chart_left + column_index * column_width + (column_width - bar_width) / 2
            total = len(subset)
            running = 0.0
            for state, colour in zip(states, colours, strict=True):
                ratio = float(subset["state"].eq(state).sum()) / total
                segment_height = ratio * row_height
                y = pane_bottom - running - segment_height
                draw.rectangle((x, y, x + bar_width, y + segment_height), fill=colour)
                if ratio >= 0.08:
                    draw.text((x + 3, y + segment_height / 2 - 9), f"{ratio:.0%}", fill="white", font=_font(17, True), stroke_width=1, stroke_fill="#42515C")
                running += segment_height

    for column_index, lower_caption in enumerate(lower_captions):
        centre = chart_left + (column_index + 0.5) * column_width
        if rotate_axis_captions:
            _draw_rotated_label(image, lower_caption[:20], centre_x=centre, bottom_y=chart_top + chart_height + 90, fill="#4E6271", font=_font(16, True))
        else:
            draw.text((centre - min(len(lower_caption) * 4, 70), chart_top + chart_height + 10), lower_caption[:20], fill="#4E6271", font=_font(16, True))

    _draw_chart_legend(draw, [(_legend_caption(legend_labels, index, state), colour, 2) for index, (state, colour) in enumerate(zip(states, colours, strict=True))], legend_position)
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0)
    return output


def _render_failure_count(title: str, frame: pd.DataFrame, group: str | None, period: str | None, legend_labels: tuple[str, ...] = (), legend_position: str = "top") -> BytesIO:
    if frame.empty or not group:
        return _empty_chart(title)
    status = _column(frame, ("Call_Status", "Test_Result", "status"))
    if not status: return _empty_chart(title)
    failed = frame[frame[status].astype(str).str.contains("failed|drop|cutoff", case=False, na=False)].copy()
    failed["__catalog_failure_state"] = failed[status].astype(str).map(
        lambda value: "Dropped" if "drop" in value.casefold() else "Failed"
    )
    row_hierarchy = sorted(
        [column for column in frame.columns if column.startswith("__catalog_row_")],
        key=lambda column: int(column.rsplit("_", 1)[1]),
    )
    column_hierarchy = sorted(
        [column for column in frame.columns if column.startswith("__catalog_column_")],
        key=lambda column: int(column.rsplit("_", 1)[1]),
    )
    if column_hierarchy:
        return _render_failure_count_hierarchy(
            title, failed, row_hierarchy, column_hierarchy, legend_labels, legend_position, comparison_frame=frame,
        )
    if len(row_hierarchy) > 1:
        return _render_failure_count_hierarchy(title, failed, [], row_hierarchy, legend_labels, legend_position, comparison_frame=frame)
    has_series = bool(period) and not frame[period].fillna("(all)").astype(str).eq("(all)").all()
    fields = [group, period] if has_series else [group]
    counts = failed.groupby([*fields, "__catalog_failure_state"], dropna=False).size().unstack(fill_value=0)
    comparison_keys = _hierarchical_complete_keys(frame, fields)
    if len(fields) == 1:
        comparison_index = pd.Index([key[0] for key in comparison_keys], name=fields[0])
    else:
        comparison_index = pd.MultiIndex.from_tuples(comparison_keys, names=fields)
    counts = counts.reindex(comparison_index, fill_value=0)
    counts = counts.head(16)
    image, draw = _canvas(title); maximum = max(int(counts.sum(axis=1).max()), 1)
    colours = {"Failed": "#E15759", "Dropped": "#F28E2B"}
    for index, (labels, values) in enumerate(counts.iterrows()):
        labels = labels if isinstance(labels, tuple) else (labels,)
        y = 120 + index * 42; x = 390
        draw.text((28, y + 4), " · ".join(str(value) for value in labels)[:42], fill="#263B4A", font=_font(17, True))
        for state in ("Failed", "Dropped"):
            count = int(values.get(state, 0)); width = int(980 * count / maximum)
            if width:
                draw.rectangle((x, y, x + width, y + 25), fill=colours[state])
                if width > 26: draw.text((x + 5, y + 3), str(count), fill="white", font=_font(16, True), stroke_width=1, stroke_fill="#42515C")
            x += width
    _draw_chart_legend(draw, [(_legend_caption(legend_labels, index, state), colours[state], 2) for index, state in enumerate(("Failed", "Dropped"))], legend_position, font_size=13)
    draw.text((390, 820), "# of failed / dropped sessions", fill="#4E6271", font=_font(19, True))
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0); return output


def _render_failure_count_hierarchy(
    title: str,
    failed: pd.DataFrame,
    row_hierarchy: list[str],
    column_hierarchy: list[str],
    legend_labels: tuple[str, ...] = (),
    legend_position: str = "top",
    comparison_frame: pd.DataFrame | None = None,
) -> BytesIO:
    """Render failure counts with template rows and columns as separate axes."""
    comparison = comparison_frame if comparison_frame is not None else failed
    # Keep the complete aggregation grids even when a combination has no
    # failure rows. This makes zero-count horizontal bars explicit rather than
    # silently removing a row/column from the comparison.
    row_keys = _hierarchical_complete_keys(comparison, row_hierarchy) if row_hierarchy else [()]
    column_keys = _hierarchical_complete_keys(comparison, column_hierarchy)
    if not row_keys or not column_keys:
        return _empty_chart(title)

    counts = failed.groupby([*row_hierarchy, *column_hierarchy, "__catalog_failure_state"], dropna=False).size()
    if counts.empty:
        maximum = 1
    else:
        maximum = max(int(counts.groupby(level=list(range(len(row_hierarchy) + len(column_hierarchy)))).sum().max()), 1)
    image, draw = _canvas(title)
    # See the status renderer above: hierarchy captions use the space between
    # the title and the plot, not the title itself.
    chart_left, chart_top, chart_height = 285, 245, 510
    # A right-side legend needs its own canvas lane. Without reserving it, the
    # diagonal outer column captions extend into the legend area on dense
    # hierarchy charts (for example Operator × Campaign failure matrices).
    chart_width = 980 if legend_position == "right" else 1250
    outer_separator_top = chart_top - 64
    row_height = chart_height / len(row_keys)
    column_width = chart_width / len(column_keys)
    colours = {"Failed": "#E15759", "Dropped": "#F28E2B"}

    outer_values = [str(key[0]) for key in column_keys]
    outer_groups: list[tuple[int, int, str]] = []
    start = 0
    while start < len(column_keys):
        end = start + 1
        while end < len(column_keys) and outer_values[end] == outer_values[start]:
            end += 1
        outer_groups.append((start, end, outer_values[start]))
        start = end
    header_font = _font(18, True)
    rotate_outer_headers = any(
        _text_width(draw, caption[:20], header_font) + 14 > (end - start) * column_width
        for start, end, caption in outer_groups
    )
    for start, end, caption in outer_groups:
        centre = chart_left + ((start + end) / 2) * column_width
        caption = caption[:20]
        if rotate_outer_headers:
            _draw_rotated_label(image, caption, centre_x=centre, bottom_y=chart_top - 27, fill="#566A78", font=header_font)
        else:
            draw.text((centre - min(len(caption) * 4, 64), chart_top - 58), caption, fill="#566A78", font=header_font)
        draw.line((chart_left + start * column_width, chart_top - 24, chart_left + end * column_width, chart_top - 24), fill="#C8D2D9", width=1)

    for column_index, column_key in enumerate(column_keys):
        lower_caption = " · ".join(str(value) for value in column_key[1:]) or str(column_key[0])
        centre = chart_left + (column_index + 0.5) * column_width
        draw.text((centre - min(len(lower_caption) * 4, 68), chart_top - 23), lower_caption[:18], fill="#4E6271", font=_font(15, True))
        cell_left = chart_left + column_index * column_width
        if column_index and column_key[0] == column_keys[column_index - 1][0]:
            # Campaigns/child values under one Operator/Vendor are related,
            # but need a lighter dashed division to remain readable.
            _draw_dashed_vertical_line(draw, cell_left, chart_top - 24, chart_top + chart_height + 25)
        else:
            # A new first-level aggregation value begins a new solid group.
            draw.line((cell_left, outer_separator_top, cell_left, chart_top + chart_height + 25), fill="#AEBBC4", width=2)
        draw.text((cell_left + 3, chart_top + chart_height + 7), "0", fill="#566A78", font=_font(13, True))
        draw.text((cell_left + column_width - 25, chart_top + chart_height + 7), str(maximum), fill="#566A78", font=_font(13, True))

    # Render each row hierarchy level in its own label column. Repeated outer
    # values are merged visually so Call Family remains distinct from G Level 4.
    label_width = max((chart_left - 28) / len(row_hierarchy), 65) if row_hierarchy else 0
    for level, field in enumerate(row_hierarchy):
        values = [str(key[level]) for key in row_keys]
        start = 0
        while start < len(row_keys):
            end = start + 1
            while end < len(row_keys) and row_keys[end][:level + 1] == row_keys[start][:level + 1]:
                end += 1
            centre_y = chart_top + ((start + end) / 2) * row_height
            x = 20 + level * label_width
            draw.text((x, centre_y - 9), values[start][:22], fill="#405765", font=_font(14, True))
            start = end

    for row_index, row_key in enumerate(row_keys):
        row_top = chart_top + row_index * row_height
        row_bottom = row_top + row_height
        next_row_key = row_keys[row_index + 1] if row_index + 1 < len(row_keys) else None
        if next_row_key is not None and len(row_hierarchy) > 1 and row_key[0] == next_row_key[0]:
            # Child values of the same first row dimension (for example the
            # cities inside one Call Family) use the same dashed hierarchy
            # convention as Campaign columns inside one Operator.
            _draw_dashed_horizontal_line(draw, row_bottom, 20, chart_left + chart_width)
        else:
            draw.line((20, row_bottom, chart_left + chart_width, row_bottom), fill="#AEBBC4", width=2)
        for column_index, column_key in enumerate(column_keys):
            key_prefix = (*row_key, *column_key)
            state_counts = {
                state: int(counts.get((*key_prefix, state), 0))
                for state in ("Failed", "Dropped")
            }
            cell_left = chart_left + column_index * column_width
            available_width = max(column_width - 10, 1)
            x = cell_left + 4
            # Dense failure charts can contain dozens of city rows.  Reserve
            # enough height for a real count label rather than rendering a
            # clipped white glyph against the bar edge.
            bar_height = max(12, min(22, row_height * 0.84))
            y = row_top + (row_height - bar_height) / 2
            count_font = _font(12, True)
            outside_counts: list[str] = []
            for state in ("Failed", "Dropped"):
                count = state_counts[state]
                segment_width = available_width * count / maximum
                if segment_width:
                    draw.rectangle((x, y, x + segment_width, y + bar_height), fill=colours[state])
                    label = str(count)
                    label_width = _text_width(draw, label, count_font)
                    if segment_width >= label_width + 10:
                        label_box = draw.textbbox((0, 0), label, font=count_font)
                        label_height = label_box[3] - label_box[1]
                        draw.text(
                            (x + (segment_width - label_width) / 2, y + (bar_height - label_height) / 2 - label_box[1]),
                            label, fill="white", font=count_font, stroke_width=1, stroke_fill="#42515C",
                        )
                    else:
                        outside_counts.append(label)
                x += segment_width
            if outside_counts:
                # Keep labels visible even when an individual stacked segment
                # is too narrow. The ordered values still follow Failed,
                # Dropped as documented by the legend.
                draw.text((x + 3, y + 1), " / ".join(outside_counts), fill="#34495A", font=count_font)

    draw.line((chart_left + chart_width, outer_separator_top, chart_left + chart_width, chart_top + chart_height + 25), fill="#AEBBC4", width=2)
    _draw_chart_legend(
        draw,
        [(_legend_caption(legend_labels, index, state), colours[state], 2) for index, state in enumerate(("Failed", "Dropped"))],
        legend_position,
        font_size=13,
        side_x=int(chart_left + chart_width + 24) if legend_position == "right" else None,
    )
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0)
    return output


def _chart_axis_hierarchy(frame: pd.DataFrame, *, distribution: bool = False) -> list[str]:
    """Return visible template hierarchy levels in their declared order."""
    rows = sorted(
        (column for column in frame.columns if column.startswith("__catalog_row_")),
        key=lambda column: int(column.rsplit("_", 1)[1]),
    )
    columns = sorted(
        (column for column in frame.columns if column.startswith("__catalog_column_")),
        key=lambda column: int(column.rsplit("_", 1)[1]),
    )
    # Distribution charts use their final column level as the stack, never as
    # an x-axis category.
    return [*rows, *(columns[:-1] if distribution and columns else columns)]


def _hierarchy_spans(keys: list[tuple[object, ...]], level: int) -> list[tuple[int, int]]:
    """Return contiguous spans sharing the same hierarchy prefix at *level*."""
    spans: list[tuple[int, int]] = []
    start = 0
    while start < len(keys):
        end = start + 1
        while end < len(keys) and keys[end][:level + 1] == keys[start][:level + 1]:
            end += 1
        spans.append((start, end))
        start = end
    return spans


def _draw_hierarchical_axis_labels(
    image: Image.Image,
    draw: ImageDraw.ImageDraw,
    keys: list[tuple[object, ...]],
    left: float,
    width: float,
    top: float,
    bottom: float,
) -> None:
    """Draw each grouping level, keeping child labels under its parent group."""
    if not keys:
        return
    levels = len(keys[0])
    item_width = width / len(keys)
    # One chart must use one axis-caption orientation. Mixing horizontal and
    # diagonal captions makes adjacent columns look misaligned and obscures
    # the grouping hierarchy.
    rotate_all_labels = any(
        _text_width(draw, str(keys[start][level])[:20], _font(18, True)) + 12 > (end - start) * item_width
        for level in range(max(levels - 1, 0))
        for start, end in _hierarchy_spans(keys, level)
    ) or any(
        _text_width(draw, str(key[-1])[:18], _font(16, True)) + 8 > item_width
        for key in keys
    )
    for level in range(max(levels - 1, 0)):
        for start, end in _hierarchy_spans(keys, level):
            centre = left + ((start + end) / 2) * item_width
            text = str(keys[start][level])[:20]
            y = top - 30 * (levels - level)
            font = _font(18, True)
            if rotate_all_labels:
                _draw_rotated_label(image, text, centre_x=centre, bottom_y=y + 24, fill="#566A78", font=font)
            else:
                draw.text((centre - min(len(text) * 4.5, 86), y), text, fill="#566A78", font=font)
            draw.line((left + start * item_width, y + 24, left + end * item_width, y + 24), fill="#CDD7DE", width=1)
    for index, key in enumerate(keys):
        text = str(key[-1])[:18]
        centre = left + (index + .5) * item_width
        font = _font(16, True)
        if rotate_all_labels:
            _draw_rotated_label(image, text, centre_x=centre, bottom_y=bottom + 78, fill="#62727E", font=font)
        else:
            draw.text((centre - min(len(text) * 4, 68), bottom + 11), text, fill="#62727E", font=font)


def _render_stacked_distribution(title: str, frame: pd.DataFrame, group: str | None, series: str | None, stack: str, legend_labels: tuple[str, ...] = (), legend_position: str = "top") -> BytesIO:
    if frame.empty or not group or not series or stack not in frame.columns:
        return _empty_chart(title)
    axis_columns = _chart_axis_hierarchy(frame, distribution=True) or [group, series]
    data = frame[[*axis_columns, stack]].dropna()
    data.attrs = frame.attrs.copy()
    combinations = _hierarchical_unique_keys(data, axis_columns)
    buckets = list(data[stack].drop_duplicates())
    if not combinations or not buckets:
        return _empty_chart(title)
    image, draw = _canvas(title); left, top, width, height = 125, 260, 1260, 475
    bar_width = max(20, min(70, width // max(len(combinations) * 2, 1)))
    bucket_colours = _series_colours([(bucket,) for bucket in buckets], [stack], data)
    for index, key in enumerate(combinations):
        subset = data
        for column, value in zip(axis_columns, key, strict=True):
            subset = subset[subset[column].astype(str) == str(value)]
        total = max(len(subset), 1); x = left + index * (width / len(combinations)) + 10; running = 0
        for bucket_index, bucket in enumerate(buckets):
            value = len(subset[subset[stack] == bucket]) / total; segment = value * height; y = top + height - running - segment
            draw.rectangle((x, y, x + bar_width, y + segment), fill=bucket_colours.get((bucket,), _colour(bucket, bucket_index)))
            running += segment
    _draw_top_column_group_separators(
        draw,
        combinations,
        axis_columns,
        left=left,
        width=width,
        top=top,
        bottom=top + height,
    )
    _draw_hierarchical_axis_labels(image, draw, combinations, left, width, top, top + height)
    _draw_chart_legend(draw, [(_legend_caption(legend_labels, index, bucket), bucket_colours.get((bucket,), _colour(bucket, index)), 2) for index, bucket in enumerate(buckets[:8])], legend_position)
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0); return output


def _combine_charts(title: str, charts: list[BytesIO]) -> BytesIO:
    """Place the two chart grammars used together by the supplied templates."""
    usable = [Image.open(chart).convert("RGB") for chart in charts]
    if not usable:
        return _empty_chart(title)
    image, _ = _canvas(title)
    width = image.width // len(usable)
    for index, chart in enumerate(usable):
        chart.thumbnail((width - 18, image.height - 85))
        x = index * width + (width - chart.width) // 2
        image.paste(chart, (x, 80 + (image.height - 80 - chart.height) // 2))
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0)
    return output


def _cdf_terminal_x_maximum(
    series_values: list[list[float]],
    low: float,
    fallback: float,
    minimum_separation: float = 0.015,
) -> float:
    """Trim only a converged CDF tail that is already effectively complete."""
    if len(series_values) < 2:
        return fallback
    candidates = sorted({value for values in series_values for value in values if low <= value <= fallback})
    for value in candidates:
        levels = sorted(sum(point <= value for point in values) / len(values) for values in series_values)
        completed_levels = [level for level in levels if level > 0.98]
        # Never crop meaningful CDF data. A tail is eligible only once at
        # least three curves have *exceeded* 98%, and those completed curves
        # have themselves converged too closely to distinguish. A coincident
        # pair alone must never truncate other still-separated CDF curves.
        completed_curves_converged = (
            len(completed_levels) >= 3
            and max(completed_levels) - min(completed_levels) < minimum_separation
        )
        if completed_curves_converged:
            return value
    return fallback


def _render_cdf_line(title: str, frame: pd.DataFrame, group: str | None, period: str | None, metric: str | None, legend_labels: tuple[str, ...] = (), legend_position: str = "top") -> BytesIO:
    if frame.empty or not group or not metric: return _empty_chart(title)
    campaign_column = _period_column(frame)
    # A CDF series is defined by the complete aggregation hierarchy, not by
    # the flattened primary/series pair.  This preserves every combination
    # when dimensions are split between Rows Aggregation and Column
    # Aggregation, or when several dimensions live on either one.
    hierarchy_columns = _chart_axis_hierarchy(frame)
    grouping_columns = hierarchy_columns or [
        *([group] if group else []),
        *([period] if period and period != group else []),
    ]
    columns = list(dict.fromkeys([*grouping_columns, metric, *([campaign_column] if campaign_column else [])]))
    data = frame[columns].copy()
    data.attrs = frame.attrs.copy()
    if campaign_column:
        data["__cdf_campaign"] = data[campaign_column].fillna("(blank)").astype(str).map(_campaign_display_value)
    data[metric] = pd.to_numeric(data[metric], errors="coerce"); data = data.dropna()
    if data.empty: return _empty_chart(title)
    combinations = _hierarchical_unique_keys(data, grouping_columns)
    series_data: list[tuple[tuple[str, ...], pd.DataFrame, list[float]]] = []
    for combination in combinations:
        mask = pd.Series(True, index=data.index)
        for column, value in zip(grouping_columns, combination, strict=True):
            mask &= data[column].astype(str).eq(str(value))
        subset = data.loc[mask]
        values = subset[metric].sort_values().tolist()
        if values:
            series_data.append((combination, subset, values))
    if not series_data:
        return _empty_chart(title)
    low = float(data[metric].min())
    observed_high = float(data[metric].max())
    series_values = [values for _, _, values in series_data]
    high = _cdf_terminal_x_maximum(series_values, low, observed_high)
    high = high if high > low else low + 1
    image, draw = _canvas(title); left, top, width, height = 100, 135, 1320, 590
    # Colour policy is driven by the template's declared dimensions: operator
    # families are used only for genuine multi-operator comparisons.
    comparison_colours = _series_colours(combinations, grouping_columns, data, line_chart=True)
    # A campaign is the temporal comparison within an operator/vendor.  Keep
    # that relationship visible even in monochrome printouts by making newer
    # campaigns progressively thicker than their earlier counterparts.
    latest_campaign = None
    campaign_count = 0
    if campaign_column and "__cdf_campaign" in data:
        campaigns = sorted(data["__cdf_campaign"].dropna().astype(str).unique(), key=_campaign_sort_key)
        campaign_count = len(campaigns)
        if len(campaigns) > 1:
            latest_campaign = campaigns[-1]
    legend_items: list[tuple[str, str, int]] = []
    for index, (combination, subset, values) in enumerate(series_data):
        visible_values = [value for value in values if value <= high]
        if not visible_values:
            continue
        label = _legend_key_caption(combination, grouping_columns, data, legend_labels)
        points = [(left + (value - low) / (high - low) * width, top + height - ((n + 1) / len(values)) * height) for n, value in enumerate(visible_values)]
        line_campaigns = subset["__cdf_campaign"].astype(str).unique() if latest_campaign else ()
        # A single-campaign chart has no historical curve to de-emphasise, so
        # use the same readable weight as the newest curve in a comparison.
        line_width = 4 if campaign_count <= 1 or (latest_campaign and latest_campaign in line_campaigns) else 1
        colour = comparison_colours.get(combination, _colour(label, index))
        draw.line(points, fill=colour, width=line_width)
        legend_items.append((label, colour, line_width))
    for tick in range(0, 101, 20):
        y = top + height - tick / 100 * height; draw.line((left, y, left + width, y), fill="#E4E9ED", width=1); draw.text((16, y - 10), f"{tick}%", fill="#4E6271", font=_font(18, True))
    for tick in range(0, 6):
        value = low + (high - low) * tick / 5
        x = left + width * tick / 5
        draw.line((x, top + height, x, top + height + 7), fill="#62727E", width=1)
        draw.text((x - 18, top + height + 7), f"{value:.1f}", fill="#4E6271", font=_font(16, True))
    draw.text((left + width / 2 - 70, top + height + 31), metric.replace("_", " "), fill="#405765", font=_font(20, True))
    _draw_chart_legend(draw, legend_items, legend_position, font_size=11, line_markers=True)
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0); return output


def _render_scatter(title: str, frame: pd.DataFrame, group: str | None, metric: str | None, x_metric: str | None, legend_labels: tuple[str, ...] = (), legend_position: str = "top") -> BytesIO:
    if frame.empty or not group or not metric or not x_metric: return _empty_chart(title)
    data = frame[[group, metric, x_metric]].copy(); data.attrs = frame.attrs.copy(); data[metric] = pd.to_numeric(data[metric], errors="coerce"); data[x_metric] = pd.to_numeric(data[x_metric], errors="coerce"); data = data.dropna()
    if data.empty: return _empty_chart(title)
    image, draw = _canvas(title); left, top, width, height = 130, 120, 1220, 600
    x_low, x_high = float(data[x_metric].min()), float(data[x_metric].max()); y_low, y_high = float(data[metric].min()), float(data[metric].max()); x_high = x_high if x_high > x_low else x_low + 1; y_high = y_high if y_high > y_low else y_low + 1
    group_keys = [(str(label),) for label in data[group].drop_duplicates().tolist()]
    group_colours = _series_colours(group_keys, [group], data)
    legend_items: list[tuple[str, str, int]] = []
    for index, (label, subset) in enumerate(data.groupby(group, sort=False)):
        for x_value, y_value in subset[[x_metric, metric]].itertuples(index=False):
            x = left + (x_value - x_low) / (x_high - x_low) * width; y = top + height - (y_value - y_low) / (y_high - y_low) * height
            draw.ellipse((x - 4, y - 4, x + 4, y + 4), fill=group_colours.get((str(label),), _colour(label, index)))
        legend_items.append((_legend_caption(legend_labels, index, label), group_colours.get((str(label),), _colour(label, index)), 2))
    for tick in range(0, 6):
        x = left + width * tick / 5; y = top + height - height * tick / 5
        x_value = x_low + (x_high - x_low) * tick / 5; y_value = y_low + (y_high - y_low) * tick / 5
        draw.line((x, top + height, x, top + height + 7), fill="#62727E", width=1)
        draw.line((left - 7, y, left, y), fill="#62727E", width=1)
        draw.text((x - 18, top + height + 7), f"{x_value:.0f}", fill="#4E6271", font=_font(16, True))
        draw.text((68, y - 9), f"{y_value:.1f}", fill="#4E6271", font=_font(16, True))
    draw.text((left + width / 2 - 120, top + height + 30), x_metric.replace("_", " "), fill="#405765", font=_font(20, True)); draw.text((40, 90), metric.replace("_", " "), fill="#405765", font=_font(20, True))
    _draw_chart_legend(draw, legend_items, legend_position, font_size=14)
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0); return output


def _render_mean_column(
    title: str,
    frame: pd.DataFrame,
    group: str | None,
    series: str | None,
    metric: str | None,
    aggregation: str = "mean",
    legend_dimensions: tuple[str, ...] = (),
    legend_position: str = "top",
) -> BytesIO:
    if frame.empty or not group or not metric:
        return _empty_chart(title)
    axis_columns = _chart_axis_hierarchy(frame) or ([group, series] if series and series != group else [group])
    data = frame[[*axis_columns, metric]].copy()
    data.attrs = frame.attrs.copy()
    data[metric] = pd.to_numeric(data[metric], errors="coerce")
    aggregate = data.dropna().groupby(axis_columns, dropna=False, sort=False)[metric]
    means = aggregate.median() if aggregation == "median" else aggregate.mean()
    if means.empty:
        return _empty_chart(title)
    image, draw = _canvas(title)
    # Reserve a true y-axis lane for the KPI label and a generous lower band
    # for a uniformly rotated hierarchy of column captions.
    left, top, chart_width, baseline, maximum = 155, 280, 1165, 680, max(float(means.max()), 1.0)
    bar_width = min(150, max(30, chart_width / max(len(means) * 1.7, 1)))
    keys = [label if isinstance(label, tuple) else (label,) for label in means.index]
    group_colours = _series_colours(keys, axis_columns, data)
    for index, (label, value) in enumerate(means.items()):
        height = (baseline - top) * float(value) / maximum
        x = left + (index + .5) * chart_width / len(means) - bar_width / 2
        draw.rectangle((x, baseline - height, x + bar_width, baseline), fill=group_colours.get(keys[index], _colour(label, index)))
        draw.text((x, baseline - height - 31), f"{float(value):.2f}", fill="#263B4A", font=_font(20, True))
    _draw_top_column_group_separators(
        draw,
        keys,
        axis_columns,
        left=left,
        width=chart_width,
        top=top,
        bottom=baseline,
    )
    _draw_hierarchical_axis_labels(image, draw, keys, left, chart_width, top, baseline)
    _draw_vertical_label(
        image, metric.replace("_", " "), centre_x=48, centre_y=(top + baseline) / 2,
        fill="#405765", font=_font(21, True),
    )
    if legend_dimensions:
        legend_items: list[tuple[str, str, int]] = []
        seen_captions: set[str] = set()
        for index, key in enumerate(keys):
            caption = _legend_key_caption(key, axis_columns, data, legend_dimensions)
            if caption in seen_captions:
                continue
            seen_captions.add(caption)
            legend_items.append((caption, group_colours.get(key, _colour(key, index)), 2))
        _draw_chart_legend(draw, legend_items, legend_position)
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0)
    return output


def _render_table(title: str, frame: pd.DataFrame, group: str | None, series: str | None, metric: str | None) -> BytesIO:
    """Render a compact mean-value table grouped exactly as declared in the template."""
    if frame.empty or not group or not metric:
        return _empty_chart(title)
    data = frame[[group, series, metric]].copy() if series else frame[[group, metric]].copy()
    data[metric] = pd.to_numeric(data[metric], errors="coerce")
    data = data.dropna(subset=[metric])
    if data.empty:
        return _empty_chart(title)
    has_series = bool(series) and not data[series].fillna("(all)").astype(str).eq("(all)").all()
    if has_series:
        table = data.pivot_table(index=group, columns=series, values=metric, aggfunc="mean", sort=False)
    else:
        table = data.groupby(group, sort=False)[metric].mean().to_frame("Value")
    image, draw = _canvas(title)
    headers = [str(table.index.name or "Category")] + [str(value) for value in table.columns]
    rows = [(str(index), *["" if pd.isna(value) else f"{float(value):.2f}" for value in values]) for index, values in table.head(18).iterrows()]
    col_width = min(310, 1450 // max(len(headers), 1)); row_height = 34; left, top = 55, 115
    for col, header in enumerate(headers):
        x = left + col * col_width
        draw.rectangle((x, top, x + col_width, top + row_height), fill="#23384A")
        draw.text((x + 8, top + 8), header[:28], fill="white", font=_font(14, True))
    for row_index, row in enumerate(rows):
        y = top + (row_index + 1) * row_height
        fill = "#F4F7F9" if row_index % 2 == 0 else "#FFFFFF"
        for col, value in enumerate(row):
            x = left + col * col_width
            draw.rectangle((x, y, x + col_width, y + row_height), fill=fill, outline="#D9E1E6")
            draw.text((x + 8, y + 8), str(value)[:28], fill="#34495A", font=_font(13))
    output = BytesIO(); image.save(output, format="PNG"); output.seek(0)
    return output


def _catalog_tokens(filters: str, keyword: str) -> tuple[str, ...]:
    text = filters.casefold()
    if keyword == "session":
        return tuple(token for token in ("volte", "multirab", "whatsapp", "classic", "call") if token in text)
    if keyword == "test":
        return tuple(token for token in ("fdfs", "fdtt", "interactivity", "brows", "http", "youtube", "video") if token in text)
    if keyword == "direction":
        return tuple(token for token in ("dl", "ul") if re.search(rf"\b{token}\b", text))
    return ()


def _catalog_spec(entry: CatalogEntry) -> dict:
    chart_type = entry.chart_type.casefold()
    metric_parts = tuple(part.strip(" `") for part in re.split(r"\s+vs\s+", entry.kpi, flags=re.I) if part.strip())
    spec: dict = {"source": entry.source_kind, "metric": metric_parts[:1] or (entry.kpi,)}
    if "scatter" in chart_type:
        spec["kind"] = "scatter"
        spec["x_metric"] = metric_parts[1:] or ("Playing_RSRP_NR_Avg", "NR_RSRP_Avg")
    elif chart_type == "count stacked horizontal bars":
        spec["kind"] = "failure_count"
    elif "100%" in chart_type or chart_type == "threshold stacked vertical bars":
        spec["kind"] = "quality_100" if chart_type == "threshold stacked vertical bars" or any(token in entry.kpi.casefold() for token in ("lq", "polqa")) else "status_100"
        if spec["kind"] == "quality_100":
            spec["threshold"] = _catalog_threshold(entry)
    else:
        spec["kind"] = "cdf_mean"
    return spec


def _chart_for_catalog_entry(entry: CatalogEntry, frames: dict[str, pd.DataFrame], multivendor: bool) -> BytesIO:
    spec = _catalog_spec(entry)
    chart_title = entry.chart_title or entry.slide_title
    legend_dimensions = _legend_dimensions(entry.legend)
    legend_labels = _legend_labels(entry.legend)
    legend_position = parse_legend_position(entry.legend_position)
    frame, group, period = _source_for_spec(frames, spec, multivendor)
    metric = _metric_column(frame, spec)
    try:
        frame = _apply_catalog_filters(frame, entry, multivendor, metric)
        frame, group, period = _apply_catalog_grouping(frame, entry, multivendor, metric)
    except ValueError:
        # A partial CDR upload should leave only the affected chart empty, not fail the report.
        return _empty_chart(chart_title)
    chart_type = entry.chart_type.casefold()
    if chart_type != "distribution stacked vertical bars" and "__catalog_stack" in frame.columns:
        frame[period] = frame[period].astype(str) + " · " + frame["__catalog_stack"].astype(str)
    if spec["kind"] == "status_100":
        return _render_status_100(chart_title, frame, group, period, legend_labels=legend_labels, legend_position=legend_position)
    if spec["kind"] == "quality_100":
        return _render_status_100(chart_title, frame, group, period, True, spec.get("threshold", 1.6), metric, legend_labels, legend_position)
    if spec["kind"] == "failure_count":
        return _render_failure_count(chart_title, frame, group, period, legend_labels, legend_position)
    if chart_type == "distribution stacked vertical bars":
        return _render_stacked_distribution(chart_title, frame, group, period, "__catalog_stack", legend_labels, legend_position)
    # Non-stacked visuals have one visual series per row/column combination. A
    # rows-only chart remains a single category series rather than becoming the
    # misleading ``Operator · Operator`` label used by the earlier renderer.
    frame["__catalog_label"] = [
        _catalogue_display_label(category, series)
        for category, series in frame[[group, period]].fillna("(blank)").itertuples(index=False, name=None)
    ]
    if spec["kind"] == "scatter":
        return _render_scatter(chart_title, frame, "__catalog_label", metric, _column(frame, spec.get("x_metric", ())), (), legend_position)
    if chart_type == "table":
        return _render_table(chart_title, frame, group, period, metric)
    if "vertical bars" in chart_type:
        return _render_mean_column(
            chart_title, frame, group, period, metric,
            aggregation="median" if chart_type == "median vertical bars" else "mean",
            legend_dimensions=legend_dimensions, legend_position=legend_position,
        )
    return _render_cdf_line(chart_title, frame, group, period, metric, legend_dimensions, legend_position)


def _chart_for_spec(title: str, frames: dict[str, pd.DataFrame], spec: dict, multivendor: bool) -> BytesIO:
    frame, group, period = _source_for_spec(frames, spec, multivendor); metric = _metric_column(frame, spec)
    if spec["kind"] == "status_100": return _render_status_100(title, frame, group, period)
    if spec["kind"] == "quality_100": return _render_status_100(title, frame, group, period, True, spec.get("threshold", 1.6), metric)
    if spec["kind"] == "failure_count": return _render_failure_count(title, frame, group, period)
    if spec["kind"] == "scatter": return _render_scatter(title, frame, group, metric, _column(frame, spec.get("x_metric", ())))
    if spec["kind"] == "dual_quality_100":
        secondary = {"kind": "quality_100", "threshold": spec.get("threshold", 1.6), **spec["secondary"]}
        other_frame, other_group, other_period = _source_for_spec(frames, secondary, multivendor)
        other_metric = _metric_column(other_frame, secondary)
        return _combine_charts(title, [
            _render_status_100("WhatsApp", frame, group, period, True, spec.get("threshold", 1.6), metric),
            _render_status_100("VoLTE", other_frame, other_group, other_period, True, secondary["threshold"], other_metric),
        ])
    if spec["kind"] == "quality_cdf":
        return _combine_charts(title, [
            _render_status_100("POLQA <1.6 rate", frame, group, period, True, spec.get("threshold", 1.6), metric),
            _render_cdf_line("POLQA AVG MOS CDF", frame, group, period, metric),
        ])
    if spec["kind"] == "cdf_pair":
        secondary_metric = _column(frame, spec.get("secondary_metric", ()))
        return _combine_charts(title, [
            _render_cdf_line(metric.replace("_", " ") if metric else title, frame, group, period, metric),
            _render_cdf_line(secondary_metric.replace("_", " ") if secondary_metric else title, frame, group, period, secondary_metric),
        ])
    if spec["kind"] == "cdf_bucket":
        # The template combines a throughput CDF with a low-rate distribution.
        # The averaged columns are retained as the numerical distribution summary.
        return _render_cdf_line(title, frame, group, period, metric)
    return _render_cdf_line(title, frame, group, period, metric)


def _clear_commentary(slide) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 10:
            shape.text_frame.clear()
            continue
        text = shape.text.strip().lower()
        if len(text) > 45 and any(hint in text for hint in COMMENT_HINTS):
            shape.text_frame.clear()


def _set_slide_header(slide, title: str, subtitle: str) -> None:
    title_shape = next(
        (
            shape for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and getattr(shape, "is_placeholder", False)
            and shape.placeholder_format.type in {1, 3}
        ),
        None,
    )
    if title_shape is None:
        title_shape = next(
            (shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.top < Inches(1.2)),
            None,
        )
    if title_shape is not None and (title or subtitle):
        # The template title placeholder is the single header surface. Keep the
        # template subtitle inside it as a second paragraph rather than adding
        # a separate text box below the placeholder.
        text_frame = title_shape.text_frame
        text_frame.clear()
        title_paragraph = text_frame.paragraphs[0]
        title_paragraph.text = title
        if subtitle:
            subtitle_paragraph = text_frame.add_paragraph()
            subtitle_paragraph.text = subtitle
            subtitle_paragraph.font.size = Pt(16)
            subtitle_paragraph.font.color.rgb = RGBColor(36, 90, 150)
    existing_subtitle = next((shape for shape in slide.shapes if shape.name == "catalogue-subtitle"), None)
    if existing_subtitle is not None:
        existing_subtitle._element.getparent().remove(existing_subtitle._element)


def _set_structural_slide_text(slide, title: str, subtitle: str) -> None:
    """Populate a title/transition layout without creating extra text boxes."""
    title_shape = next(
        (
            shape for shape in slide.placeholders
            if shape.placeholder_format.type in {1, 3}
        ),
        None,
    )
    subtitle_shape = next(
        (
            shape for shape in slide.placeholders
            if shape.placeholder_format.type == 4
        ),
        None,
    )
    if title_shape is not None:
        title_shape.text = title
        title_shape.text_frame.word_wrap = True
        title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if subtitle_shape is not None:
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.word_wrap = True
        subtitle_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif subtitle:
        _set_slide_header(slide, title, subtitle)

    # Structural slides keep only their title/subtitle placeholders. Branding
    # and decorations inherited from the master/layout remain untouched.
    retained_elements = {
        shape._element for shape in (title_shape, subtitle_shape) if shape is not None
    }
    for shape in list(slide.placeholders):
        if shape._element in retained_elements:
            continue
        shape._element.getparent().remove(shape._element)


def _chart_frames(slide) -> list[tuple[int, int, int, int]]:
    """Remove example chart images/groups and return their occupied areas.

    Logos and small decorative images remain untouched. The new report chart is
    placed in the exact bounding area of the removed template chart(s).
    """
    removable_types = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.GROUP}
    frames: list[tuple[int, int, int, int]] = []
    for shape in list(slide.shapes):
        if shape.shape_type not in removable_types:
            continue
        area = shape.width * shape.height
        if shape.top < Inches(0.9) or area < Inches(1.5) * Inches(1.5):
            continue
        frames.append((shape.left, shape.top, shape.width, shape.height))
        shape._element.getparent().remove(shape._element)
    return frames


def _remove_template_chart_placeholders(slide) -> None:
    """Remove inherited chart pictures stored *inside* the template placeholders.

    The supplied templates encode their sample Tableau exports as picture
    placeholders rather than regular picture shapes.  Removing only regular
    images therefore left the old chart under the new one.  Index 0 is the
    master title and index 10 is the deliberately blank analyst-comments area.
    """
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue
        if shape.placeholder_format.idx in {0, 10}:
            continue
        shape._element.getparent().remove(shape._element)


def _combined_frame(frames: list[tuple[int, int, int, int]]) -> tuple[int, int, int, int] | None:
    if not frames:
        return None
    left = min(frame[0] for frame in frames)
    top = min(frame[1] for frame in frames)
    right = max(frame[0] + frame[2] for frame in frames)
    bottom = max(frame[1] + frame[3] for frame in frames)
    return left, top, right - left, bottom - top


def _named_slide_layout(presentation: Presentation, layout_name: str):
    """Resolve a template layout name against the template slide master."""
    expected = layout_name.strip().casefold()
    for layout in presentation.slide_layouts:
        if layout.name.strip().casefold() == expected:
            return layout
    return None


def _remove_all_slides(presentation: Presentation) -> None:
    """Leave the source deck as a master/layout-only presentation."""
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _apply_catalogue_layout(presentation: Presentation, slide, layout_name: str):
    """Assign the requested layout and align retained structural placeholders."""
    layout = _named_slide_layout(presentation, layout_name)
    if layout is None:
        return None
    layout_relationship_ids = [
        relationship_id
        for relationship_id, relationship in slide.part.rels.items()
        if relationship.reltype == RT.SLIDE_LAYOUT
    ]
    for relationship_id in layout_relationship_ids:
        slide.part.rels.pop(relationship_id)
    slide.part.rels.get_or_add(RT.SLIDE_LAYOUT, layout.part)

    layout_placeholders = {
        placeholder.placeholder_format.idx: placeholder
        for placeholder in layout.placeholders
    }
    for shape in slide.placeholders:
        target = layout_placeholders.get(shape.placeholder_format.idx)
        if target is None or shape.placeholder_format.idx not in {0, 10}:
            continue
        shape.left, shape.top = target.left, target.top
        shape.width, shape.height = target.width, target.height
    return layout


def _layout_chart_frames(layout) -> list[tuple[int, int, int, int]]:
    """Read chart placeholder frames from the selected template layout."""
    if layout is None:
        return []
    frames = [
        (shape.left, shape.top, shape.width, shape.height)
        for shape in layout.placeholders
        if shape.placeholder_format.type == 7 and shape.placeholder_format.idx != 10
    ]
    # PowerPoint layouts commonly differ by one or two EMU between placeholders
    # that are visually on the same row. A strict (top, left) sort therefore
    # placed the right-hand chart before the left-hand chart. Build visual rows
    # with a small vertical tolerance, then order each row from left to right.
    visual_rows: list[list[tuple[int, int, int, int]]] = []
    row_tolerance = int(Inches(0.08))
    for frame in sorted(frames, key=lambda item: item[1]):
        matching_row = next(
            (row for row in visual_rows if abs(frame[1] - min(item[1] for item in row)) <= row_tolerance),
            None,
        )
        if matching_row is None:
            visual_rows.append([frame])
        else:
            matching_row.append(frame)
    visual_rows.sort(key=lambda row: min(frame[1] for frame in row))
    return [frame for row in visual_rows for frame in sorted(row, key=lambda item: item[0])]


def render_cdr_report(destination: Path, template: Path, frames: dict[str, pd.DataFrame] | None, technology: str,
                      multivendor: bool, catalog: list[CatalogEntry] | None = None,
                      chart_output_dir: Path | None = None,
                      frame_loader: Callable[[str], pd.DataFrame] | None = None,
                      on_chart_rendered: Callable[[CatalogEntry, int, bool], None] | None = None) -> Path:
    if not template.exists():
        raise FileNotFoundError(f"Reporting template not found: {template.name}")
    if not catalog:
        raise ValueError("A Slides Template is required to generate the report.")
    # A report may concatenate campaigns that were exported using different
    # operator spellings.  Apply aliases only to these in-memory report frames
    # before filtering and grouping; Workspace datasets remain source-faithful.
    if frames is None and frame_loader is None:
        raise ValueError('Report rendering requires CDR frames or a frame loader.')
    # Keep only the source currently needed when a loader is supplied.  This
    # avoids retaining Data, Voice and Speech frames simultaneously for large
    # reports on resource-constrained servers.
    cached_frames = {source: normalise_report_operator_aliases(frame) for source, frame in (frames or {}).items()}
    active_source: str | None = None
    def frame_for(source: str) -> pd.DataFrame:
        nonlocal active_source
        # Keep at most one lazy-loaded CDR frame resident.  The catalogue
        # order still determines slide/chart order, while switching sources
        # releases the previous large DataFrame before loading the next one.
        if frame_loader is not None and source != active_source:
            cached_frames.clear()
            gc.collect()
        if source not in cached_frames:
            if frame_loader is None:
                raise ValueError(f'No CDR frame is available for {source}.')
            cached_frames[source] = normalise_report_operator_aliases(frame_loader(source))
        active_source = source
        return cached_frames[source]
    presentation = Presentation(template)
    _remove_all_slides(presentation)
    rendered_charts: list[dict[str, object]] = []

    catalogue_slides: dict[int, list[CatalogEntry]] = defaultdict(list)
    render_catalog = [prepare_multivendor_catalog_entry(entry) if multivendor else entry for entry in catalog]
    for entry in render_catalog:
        catalogue_slides[entry.slide].append(entry)

    for number in sorted(catalogue_slides):
        slide_entries = catalogue_slides[number]
        header = slide_entries[0]
        structural = header.structural_type
        if not structural and header.chart_type.strip().casefold() in PRESERVED_CHART_TYPES:
                structural = "title slide" if number == min(catalogue_slides) else "transition slide"

        if structural:
            layout_name = header.layout or ("Title Page" if structural == "title slide" else "Title Only")
            layout = _named_slide_layout(presentation, layout_name)
            if layout is None:
                raise ValueError(f"Slide {number}: layout '{layout_name}' does not exist in the selected template.")
            slide = presentation.slides.add_slide(layout)
            _set_structural_slide_text(slide, header.slide_title, header.slide_subtitle)
            continue

        chart_entries = [entry for entry in slide_entries if entry.source_kind]
        if not chart_entries:
            raise ValueError(
                f"Slide {number} does not define an automated chart, a Title Slide or a Transition Slide."
            )
        layouts = {entry.layout for entry in chart_entries}
        if len(layouts) != 1:
            raise ValueError(f"Slide {number} uses more than one Layout in the active template.")
        layout_name = layouts.pop()
        layout = _named_slide_layout(presentation, layout_name)
        if layout is None:
            raise ValueError(f"Slide {number}: layout '{layout_name}' does not exist in the selected template.")
        placement_frames = _layout_chart_frames(layout)
        if len(placement_frames) < len(chart_entries):
            raise ValueError(
                f"Slide {number}: layout '{layout_name}' has {len(placement_frames)} chart placeholders, "
                f"but the template defines {len(chart_entries)} charts."
            )
        slide = presentation.slides.add_slide(layout)
        _set_slide_header(slide, header.slide_title, header.slide_subtitle)
        _clear_commentary(slide)
        _remove_template_chart_placeholders(slide)
        for chart_index, (entry, placement) in enumerate(zip(chart_entries, placement_frames[:len(chart_entries)], strict=True), start=1):
            # The PowerPoint must consume the exact same renderer used by the
            # editor and Report Charts.  Keeping this call shared prevents an
            # unnoticed drift in normalisation, multivendor preparation or
            # catalogue filtering between preview and export.
            source_frame = frame_for(entry.source_kind)
            chart_bytes = render_catalog_chart_preview(source_frame, entry, multivendor=multivendor)
            # Rebuild the source frame before retrying an unexpected empty
            # chart. Retrying the same already-loaded frame cannot recover a
            # worker that was under memory pressure while materialising it.
            for _attempt in range(2):
                if not is_empty_catalog_chart(chart_bytes, entry) or frame_loader is None:
                    break
                del chart_bytes
                cached_frames.pop(entry.source_kind, None)
                gc.collect()
                source_frame = frame_for(entry.source_kind)
                chart_bytes = render_catalog_chart_preview(source_frame, entry, multivendor=multivendor)
            if on_chart_rendered:
                on_chart_rendered(entry, len(source_frame.index), is_empty_catalog_chart(chart_bytes, entry))
            if chart_output_dir is not None:
                chart_output_dir.mkdir(parents=True, exist_ok=True)
                file_name = f'slide-{number:03d}-chart-{chart_index:02d}.png'
                (chart_output_dir / file_name).write_bytes(chart_bytes)
                rendered_charts.append({
                    'slide': number,
                    'title': entry.chart_title or entry.slide_title or f'Slide {number}',
                    'source': entry.cdr_source,
                    'chart_type': entry.chart_type,
                    'file': file_name,
                })
            slide.shapes.add_picture(BytesIO(chart_bytes), *placement)
    cached_frames.clear()
    gc.collect()
    presentation.save(destination)
    if chart_output_dir is not None:
        (chart_output_dir / 'manifest.json').write_text(json.dumps({'charts': rendered_charts}), encoding='utf-8')
    return destination
