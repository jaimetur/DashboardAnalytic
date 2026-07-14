from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

POWERPOINT_EXPORT_VERSION = "2026-07-14-v6"
PPT_WIDTH_IN = 13.333
PPT_HEIGHT_IN = 7.5
SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 720
MAX_EXPORT_CDF_POINTS = 320

BG = "#F4F0E7"
PANEL = "#FFFDF8"
PANEL_ALT = "#F6FBFA"
TEXT = "#102131"
MUTED = "#61727D"
TEAL = "#0B7A75"
ORANGE = "#DD653E"
BLUE = "#245A96"
LINE = "#D9E4E8"
DARK_BG = "#143048"
SERIES_COLORS = ["#0B7A75", "#DD653E", "#245A96", "#B84D3A", "#6D46A8", "#228A5D", "#C78B1D", "#4D6A88"]


def export_word_report(destination: Path, analysis: dict[str, Any]) -> Path:
    document = Document()
    document.add_heading("Dashboard Analytic Report", level=0)
    document.add_paragraph(f"Metric analysed: {analysis['selected_metric']}")
    document.add_paragraph(f"Filtered rows: {analysis['kpis']['rows']}")

    document.add_heading("Key KPIs", level=1)
    for key, value in analysis["kpis"].items():
        document.add_paragraph(f"{key}: {value}")

    document.add_heading("Scorecard", level=1)
    for row in analysis["scorecard"]:
        document.add_paragraph(f"{row['label']}: {row['value']}")

    document.save(destination)
    return destination


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _load_font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/SFNS.ttf",
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def _format_label(value: str) -> str:
    normalized = str(value or "").strip()
    if not normalized:
        return "n/a"
    return normalized.replace("_", " ").title()


def _format_value(value: Any) -> str:
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return f"{value:,.4f}".rstrip("0").rstrip(".")
    return str(value)


def _add_full_bg(slide, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(PPT_WIDTH_IN), Inches(PPT_HEIGHT_IN))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()


def _add_panel(slide, left: float, top: float, width: float, height: float, fill: str = PANEL, line: str = LINE, radius: int = 18):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.08
    return shape


def _add_textbox(slide, left: float, top: float, width: float, height: float, text: str, *, size: int = 14, bold: bool = False,
                 color: str = TEXT, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _add_multiline_textbox(slide, left: float, top: float, width: float, line_height: float, lines: list[str], *,
                           size: int = 10, bold: bool = False, color: str = MUTED) -> None:
    if not lines:
        return
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(line_height * len(lines)))
    frame = textbox.text_frame
    frame.clear()
    first = True
    for line in lines:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.LEFT
        if paragraph.runs:
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
            run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        first = False


def _add_badge(slide, left: float, top: float, width: float, text: str, *, fill: str, color: str = "#FFFFFF") -> None:
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.36))
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(fill)
    badge.line.fill.background()
    _add_textbox(slide, left + 0.08, top + 0.03, width - 0.16, 0.22, text, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)


def _filters_summary(filters: dict[str, Any]) -> str:
    fragments: list[str] = []
    for key in ["market", "period"]:
        values = filters.get(key) or []
        if values:
            fragments.append(f"{_format_label(key)}: {', '.join(str(item) for item in values)}")
    for key, values in (filters.get("extra_filters") or {}).items():
        if not values or values == ["__none__"]:
            continue
        if not isinstance(values, list):
            values = [values]
        fragments.append(f"{_format_label(key)}: {', '.join(str(item) for item in values)}")
    if filters.get("date_from"):
        fragments.append(f"Date From: {filters['date_from']}")
    if filters.get("date_to"):
        fragments.append(f"Date To: {filters['date_to']}")
    return " | ".join(fragments) if fragments else "No filters selected"


def _filters_lines(filters_text: str) -> list[str]:
    lines = [line.strip() for line in str(filters_text or "").split(" | ") if line.strip()]
    return lines or ["No filters selected"]


def _filters_summary_lines(filters_text: str, max_lines: int = 2) -> list[str]:
    parts = _filters_lines(filters_text)
    if len(parts) <= 1:
        return parts
    if len(parts) <= max_lines:
        return parts
    first_line_count = (len(parts) + 1) // 2
    return [
        " | ".join(parts[:first_line_count]),
        " | ".join(parts[first_line_count:]),
    ][:max_lines]


def _structured_filters_summary_lines(filters_text: str) -> list[str]:
    parts = _filters_lines(filters_text)
    if not parts:
        return ["No filters selected"]
    buckets = [
        ["Market", "Period", "Operator", "Vendor"],
        ["Region", "City"],
        ["Session Type", "Test Name", "Direction"],
    ]
    lines: list[str] = []
    used: set[str] = set()
    for bucket in buckets:
        bucket_values = [
            part for part in parts
            if any(part.startswith(f"{label}:") for label in bucket)
        ]
        if bucket_values:
            lines.append(" | ".join(bucket_values))
            used.update(bucket_values)
    remaining = [part for part in parts if part not in used]
    if remaining:
        lines.extend(remaining)
    return lines[:3]


def _draw_wrapped_text(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], text: str, font: ImageFont.ImageFont,
                       fill: str = TEXT, line_spacing: int = 8) -> None:
    x0, y0, x1, y1 = box
    max_width = max(50, x1 - x0)
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if draw.textlength(candidate, font=font) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    y = y0
    for line in lines:
        draw.text((x0, y), line, font=font, fill=fill)
        y += font.size + line_spacing
        if y > y1:
            break


def _empty_chart_image(title: str, message: str) -> BytesIO:
    image = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), BG)
    draw = ImageDraw.Draw(image)
    title_font = _load_font(26, bold=True)
    body_font = _load_font(18)
    draw.rounded_rectangle((24, 24, SLIDE_WIDTH - 24, SLIDE_HEIGHT - 24), radius=20, fill=PANEL, outline=LINE, width=2)
    draw.text((44, 34), title, font=title_font, fill=TEXT)
    _draw_wrapped_text(draw, (44, 120, SLIDE_WIDTH - 44, SLIDE_HEIGHT - 44), message, body_font, fill=MUTED)
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


def _downsample_series(labels: list[float], series: list[float], limit: int = MAX_EXPORT_CDF_POINTS) -> tuple[list[float], list[float]]:
    if len(labels) <= limit or len(series) <= limit:
        return labels, series
    sampled_labels: list[float] = []
    sampled_series: list[float] = []
    last_index = len(labels) - 1
    for position in range(limit):
        index = round((position / max(limit - 1, 1)) * last_index)
        sampled_labels.append(labels[index])
        sampled_series.append(series[index])
    return sampled_labels, sampled_series


def _draw_bar_chart(chart: dict[str, Any]) -> BytesIO:
    labels = [str(value) for value in chart.get("labels", [])]
    values = [float(value) for value in chart.get("series", [])]
    if not labels or not values:
        return _empty_chart_image("Group Benchmark", "The selected comparison does not contain grouped values for this metric.")

    image = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), BG)
    draw = ImageDraw.Draw(image)
    title_font = _load_font(22, bold=True)
    axis_font = _load_font(15)
    value_font = _load_font(14, bold=True)
    draw.rounded_rectangle((20, 20, SLIDE_WIDTH - 20, SLIDE_HEIGHT - 20), radius=18, fill=PANEL, outline=LINE, width=2)
    draw.text((36, 28), "Group Benchmark", font=title_font, fill=TEXT)

    left, top, right, bottom = 70, 90, SLIDE_WIDTH - 40, SLIDE_HEIGHT - 70
    draw.line((left, top, left, bottom), fill=MUTED, width=2)
    draw.line((left, bottom, right, bottom), fill=MUTED, width=2)

    compact_labels = labels[:8]
    compact_values = values[:8]
    max_value = max(compact_values) if compact_values else 1.0
    max_value = max(max_value, 1.0)
    count = max(1, len(compact_values))
    gap = 16
    usable_width = right - left - gap * (count + 1)
    bar_width = max(36, usable_width // count)
    chart_height = bottom - top - 24

    for index, (label, value) in enumerate(zip(compact_labels, compact_values, strict=False)):
        x0 = left + gap + index * (bar_width + gap)
        x1 = x0 + bar_width
        bar_height = int((value / max_value) * chart_height)
        y0 = bottom - bar_height
        draw.rounded_rectangle((x0, y0, x1, bottom), radius=10, fill=ORANGE)
        value_text = _format_value(value)
        text_width = draw.textlength(value_text, font=value_font)
        label_y = y0 + 8 if bar_height > 28 else max(top + 8, y0 - 20)
        label_color = "#FFFFFF" if bar_height > 28 else TEXT
        draw.text((x0 + (bar_width - text_width) / 2, label_y), value_text, font=value_font, fill=label_color)
        short_label = label if len(label) <= 14 else f"{label[:12]}.."
        text_width = draw.textlength(short_label, font=axis_font)
        draw.text((x0 + (bar_width - text_width) / 2, bottom + 10), short_label, font=axis_font, fill=MUTED)

    buffer = BytesIO()
    image.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


def _draw_line_chart(chart: dict[str, Any]) -> BytesIO:
    if chart.get("series_collection"):
        series_collection = []
        for item in chart.get("series_collection", []):
            if not item.get("labels") or not item.get("series"):
                continue
            sampled_labels, sampled_series = _downsample_series(
                [float(value) for value in item.get("labels", [])],
                [float(value) for value in item.get("series", [])],
            )
            series_collection.append({
                "name": str(item.get("name") or "Series"),
                "labels": sampled_labels,
                "series": sampled_series,
            })
    else:
        labels, series = _downsample_series(
            [float(value) for value in chart.get("labels", [])],
            [float(value) for value in chart.get("series", [])],
        )
        series_collection = [{"name": "CDF", "labels": labels, "series": series}] if labels and series else []

    if not series_collection:
        return _empty_chart_image("CDF Curve", "The selected CDF comparison does not contain numeric values for this metric.")

    image = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), BG)
    draw = ImageDraw.Draw(image)
    title_font = _load_font(22, bold=True)
    legend_font = _load_font(14)
    draw.rounded_rectangle((20, 20, SLIDE_WIDTH - 20, SLIDE_HEIGHT - 20), radius=18, fill=PANEL, outline=LINE, width=2)
    draw.text((36, 28), "CDF Curve", font=title_font, fill=TEXT)

    left, top, right, bottom = 70, 90, SLIDE_WIDTH - 40, SLIDE_HEIGHT - 110
    draw.line((left, top, left, bottom), fill=MUTED, width=2)
    draw.line((left, bottom, right, bottom), fill=MUTED, width=2)

    all_x = [value for s in series_collection for value in s["labels"]]
    all_y = [value for s in series_collection for value in s["series"]]
    min_x, max_x = min(all_x or [0.0]), max(all_x or [1.0])
    min_y, max_y = min(all_y or [0.0]), max(all_y or [1.0])
    if max_x == min_x:
        max_x = min_x + 1.0
    if max_y == min_y:
        max_y = min_y + 1.0

    for index, s in enumerate(series_collection[:8]):
        color = SERIES_COLORS[index % len(SERIES_COLORS)]
        points: list[tuple[float, float]] = []
        for x_value, y_value in zip(s["labels"], s["series"], strict=False):
            px = left + ((x_value - min_x) / (max_x - min_x)) * (right - left)
            py = bottom - ((y_value - min_y) / (max_y - min_y)) * (bottom - top)
            points.append((px, py))
        if len(points) >= 2:
            draw.line(points, fill=color, width=4)
        for point in points[::max(1, len(points) // 10 or 1)]:
            draw.ellipse((point[0] - 3, point[1] - 3, point[0] + 3, point[1] + 3), fill=color)

    legend_x = 48
    legend_y = SLIDE_HEIGHT - 64
    for index, s in enumerate(series_collection[:8]):
        color = SERIES_COLORS[index % len(SERIES_COLORS)]
        draw.rounded_rectangle((legend_x, legend_y, legend_x + 16, legend_y + 16), radius=4, fill=color)
        draw.text((legend_x + 24, legend_y - 2), str(s["name"])[:18], font=legend_font, fill=MUTED)
        legend_x += 140
        if legend_x > SLIDE_WIDTH - 180:
            legend_x = 48
            legend_y += 22

    buffer = BytesIO()
    image.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


def _add_metric_cards(slide, analyses: list[dict[str, Any]]) -> None:
    cols = 3
    gap_x = 0.2
    gap_y = 0.18
    card_w = 4.0
    card_h = 2.45
    start_x = 0.55
    start_y = 1.64
    for index, analysis_item in enumerate(analyses[:6]):
        result = analysis_item["result"]
        metric_name = result.get("selected_metric") or analysis_item.get("metric") or "Metric"
        row = index // cols
        col = index % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        _add_panel(slide, x, y, card_w, card_h, fill="#163247", line="#254961")
        _add_textbox(slide, x + 0.16, y + 0.12, 2.4, 0.24, metric_name, size=14, bold=True, color="#FFFFFF")
        _add_badge(slide, x + card_w - 1.28, y + 0.12, 1.0, f"{result.get('metric_kpis', {}).get('samples', 0)}", fill=TEAL)
        metrics = result.get("metric_kpis") or {}
        summary = [
            ("Mean", metrics.get("mean_metric", 0)),
            ("Avg", metrics.get("avg_metric", 0)),
            ("P10", metrics.get("p10_metric", 0)),
            ("P90", metrics.get("p90_metric", 0)),
            ("Min", metrics.get("min_metric", 0)),
            ("Max", metrics.get("max_metric", 0)),
        ]
        inner_cols = 3
        stat_w = 1.12
        stat_h = 0.66
        stat_gap_x = 0.12
        stat_gap_y = 0.1
        start_stat_x = x + 0.16
        start_stat_y = y + 0.62
        for metric_index, (label, value) in enumerate(summary):
            inner_row = metric_index // inner_cols
            inner_col = metric_index % inner_cols
            sx = start_stat_x + inner_col * (stat_w + stat_gap_x)
            sy = start_stat_y + inner_row * (stat_h + stat_gap_y)
            _add_panel(slide, sx, sy, stat_w, stat_h, fill="#21435B", line="#2B536B")
            _add_textbox(slide, sx + 0.08, sy + 0.08, stat_w - 0.16, 0.12, label, size=8, bold=True, color="#AFC3D2")
            _add_textbox(slide, sx + 0.08, sy + 0.28, stat_w - 0.16, 0.16, _format_value(value), size=11, bold=True, color="#FFFFFF")


def _add_kpi_grid(slide, items: list[tuple[str, Any]], *, left: float, top: float, width: float, columns: int, card_height: float = 0.9,
                  fill: str = PANEL_ALT, accent: str = TEAL) -> None:
    if not items:
        return
    gap_x = 0.16
    gap_y = 0.14
    card_width = (width - gap_x * (columns - 1)) / columns
    for index, (label, value) in enumerate(items):
        row = index // columns
        col = index % columns
        x = left + col * (card_width + gap_x)
        y = top + row * (card_height + gap_y)
        _add_panel(slide, x, y, card_width, card_height, fill=fill, line=LINE)
        _add_textbox(slide, x + 0.14, y + 0.12, card_width - 0.28, 0.18, _format_label(label), size=9, bold=True, color=MUTED)
        _add_textbox(slide, x + 0.14, y + 0.44, card_width - 0.28, 0.24, _format_value(value), size=17, bold=True, color=TEXT)


def _add_metric_kpi_strip(slide, metric_kpis: dict[str, Any], *, left: float, top: float, width: float) -> None:
    summary = [
        ("Mean", metric_kpis.get("mean_metric", 0)),
        ("Avg", metric_kpis.get("avg_metric", 0)),
        ("P10", metric_kpis.get("p10_metric", 0)),
        ("P90", metric_kpis.get("p90_metric", 0)),
        ("Min", metric_kpis.get("min_metric", 0)),
        ("Max", metric_kpis.get("max_metric", 0)),
    ]
    cols = len(summary)
    gap = 0.08
    cell_width = (width - gap * (cols - 1)) / cols
    for index, (label, value) in enumerate(summary):
        x = left + index * (cell_width + gap)
        _add_panel(slide, x, top, cell_width, 0.34, fill="#FFF9F3", line=LINE)
        _add_textbox(slide, x + 0.05, top + 0.04, cell_width - 0.1, 0.08, label, size=6, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        _add_textbox(slide, x + 0.05, top + 0.15, cell_width - 0.1, 0.1, _format_value(value), size=9, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def _add_grouped_scorecard_table(slide, groups: list[dict[str, Any]], left: float, top: float, width: float, height: float) -> None:
    if not groups:
        _add_textbox(slide, left, top, width, 0.3, "No percentile groups available for this metric.", size=12, color=MUTED)
        return
    table = slide.shapes.add_table(len(groups) + 1, 6, Inches(left), Inches(top), Inches(width), Inches(height)).table
    headers = ["Group", "P10", "P25", "P50", "P75", "P90"]
    for index, header in enumerate(headers):
        cell = table.cell(0, index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb("E8F4F3")
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.color.rgb = _rgb(TEXT)
    for row_index, group in enumerate(groups[:8], start=1):
        values_by_label = {str(item.get("label")): item.get("value") for item in group.get("items") or []}
        row_values = [
            str(group.get("group") or "Overall"),
            _format_value(values_by_label.get("P10", "-")),
            _format_value(values_by_label.get("P25", "-")),
            _format_value(values_by_label.get("P50", "-")),
            _format_value(values_by_label.get("P75", "-")),
            _format_value(values_by_label.get("P90", "-")),
        ]
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("FFFFFF" if row_index % 2 else "F8FBFB")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = _rgb(TEXT)


def _add_data_table(slide, rows: list[dict[str, Any]], left: float, top: float, width: float, height: float) -> None:
    if not rows:
        _add_textbox(slide, left, top, width, 0.3, "No table rows are available for the selected filters.", size=12, color=MUTED)
        return
    headers = list(rows[0].keys())[:8]
    sliced_rows = rows[:10]
    compact_height = min(height, 0.42 + (len(sliced_rows) + 1) * 0.32)
    table = slide.shapes.add_table(len(sliced_rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(compact_height)).table
    for index, header in enumerate(headers):
        cell = table.cell(0, index)
        cell.text = _format_label(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb("EAF1F8")
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(9)
                run.font.color.rgb = _rgb(TEXT)
    for row_index, row in enumerate(sliced_rows, start=1):
        for col_index, header in enumerate(headers):
            cell = table.cell(row_index, col_index)
            cell.text = _format_value(row.get(header, ""))
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("FFFFFF" if row_index % 2 else "F9FBFD")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.color.rgb = _rgb(TEXT)


def _build_powerpoint_payload(report: dict[str, Any]) -> dict[str, Any]:
    analyses = report.get("analyses") or []
    primary = analyses[0]["result"] if analyses else None
    return {
        "dataset_name": report.get("dataset_name") or "Dataset",
        "dataset_type": report.get("dataset_type") or "Other",
        "filters_text": report.get("filters_text") or "No filters selected",
        "selected_metrics": report.get("selected_metrics") or [],
        "global_kpis": (primary or {}).get("global_kpis") or {},
        "analyses": analyses,
    }


def _init_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(PPT_WIDTH_IN)
    presentation.slide_height = Inches(PPT_HEIGHT_IN)
    return presentation


def export_powerpoint_report(destination: Path, report: dict[str, Any]) -> Path:
    payload = _build_powerpoint_payload(report)
    presentation = _init_presentation()

    cover = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_full_bg(cover, DARK_BG)
    _add_textbox(cover, 0.7, 0.8, 7.8, 0.4, "Dashboard Analytic", size=28, bold=True, color="#FFFFFF")
    _add_textbox(cover, 0.7, 1.35, 8.5, 0.35, payload["dataset_name"], size=22, bold=True, color="#D9EEF3")
    _add_badge(cover, 0.7, 2.05, 1.4, payload["dataset_type"], fill=TEAL)
    for line_index, filter_line in enumerate(_filters_lines(payload["filters_text"])[:8]):
        _add_textbox(cover, 0.7, 2.55 + line_index * 0.28, 7.6, 0.22, filter_line, size=12, color="#E4EEF1")
    _add_panel(cover, 8.75, 0.75, 3.85, 5.95, fill="#1C4665", line="#2D607B")
    _add_textbox(cover, 9.0, 1.05, 3.2, 0.28, "Export Contents", size=14, bold=True, color="#FFFFFF")
    contents = [
        "Executive Dashboard",
        "Global Metrics",
        "Metric KPI cards",
        "Visual Analytics per metric",
        "Processed Metrics table",
    ]
    for index, label in enumerate(contents):
        _add_badge(cover, 9.0, 1.55 + index * 0.72, 3.0, label, fill=BLUE)

    global_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_full_bg(global_slide, BG)
    _add_textbox(global_slide, 0.55, 0.35, 5.5, 0.28, "Executive Dashboard", size=22, bold=True)
    _add_badge(global_slide, 9.95, 0.34, 2.75, payload["dataset_name"][:28], fill=TEAL)
    _add_multiline_textbox(global_slide, 0.55, 0.72, 12.1, 0.18, _structured_filters_summary_lines(payload["filters_text"]), size=9, color=MUTED)
    global_items = [(key, value) for key, value in payload["global_kpis"].items() if key not in {"date_from", "date_to"}]
    _add_kpi_grid(global_slide, global_items[:16], left=0.55, top=1.58, width=12.2, columns=4, card_height=1.0)

    metric_card_pages = [payload["analyses"][index:index + 6] for index in range(0, len(payload["analyses"]), 6)] or [[]]
    for page_index, metric_page in enumerate(metric_card_pages, start=1):
        metric_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _add_full_bg(metric_slide, BG)
        _add_textbox(metric_slide, 0.55, 0.35, 5.5, 0.28, "Executive Dashboard", size=22, bold=True)
        title = "Selected Metric Cards" if len(metric_card_pages) == 1 else f"Selected Metric Cards · Page {page_index}"
        _add_textbox(metric_slide, 0.55, 0.72, 6.0, 0.24, title, size=12, bold=True, color=ORANGE)
        _add_multiline_textbox(metric_slide, 0.55, 0.96, 12.1, 0.18, _structured_filters_summary_lines(payload["filters_text"]), size=9, color=MUTED)
        _add_metric_cards(metric_slide, metric_page)

    for analysis_item in payload["analyses"]:
        result = analysis_item["result"]
        metric_name = result.get("selected_metric") or analysis_item.get("metric") or "Metric"
        visual = presentation.slides.add_slide(presentation.slide_layouts[6])
        _add_full_bg(visual, BG)
        _add_textbox(visual, 0.55, 0.35, 6.0, 0.28, f"Visual Analytics · {metric_name}", size=22, bold=True)
        _add_multiline_textbox(visual, 0.55, 0.72, 12.1, 0.18, _structured_filters_summary_lines(payload["filters_text"]), size=9, color=MUTED)
        cdf_image = _draw_line_chart(result.get("cdf_chart") or {})
        comparison_image = _draw_bar_chart(result.get("comparison_chart") or {})
        visual.shapes.add_picture(cdf_image, Inches(0.55), Inches(1.48), width=Inches(6.0), height=Inches(2.8))
        visual.shapes.add_picture(comparison_image, Inches(6.78), Inches(1.48), width=Inches(6.0), height=Inches(2.8))
        _add_metric_kpi_strip(visual, result.get("metric_kpis") or {}, left=0.55, top=4.38, width=12.2)
        _add_textbox(visual, 0.55, 4.82, 3.4, 0.18, f"Grouped Percentiles · { _format_label(result.get('filters', {}).get('aggregation') or 'all') }", size=9, bold=True, color=BLUE)
        _add_grouped_scorecard_table(visual, (result.get("scorecard_groups") or [])[:8], left=0.55, top=5.04, width=12.2, height=0.84)

    primary_result = payload["analyses"][0]["result"] if payload["analyses"] else {}
    table_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_full_bg(table_slide, BG)
    _add_textbox(table_slide, 0.55, 0.35, 5.5, 0.28, "Processed Metrics", size=22, bold=True)
    _add_multiline_textbox(table_slide, 0.55, 0.72, 12.1, 0.18, _structured_filters_summary_lines(payload["filters_text"]), size=9, color=MUTED)
    _add_data_table(table_slide, primary_result.get("table_rows") or [], left=0.55, top=1.56, width=12.2, height=5.28)

    presentation.save(destination)
    return destination
