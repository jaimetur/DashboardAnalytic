from __future__ import annotations

import json
import time
from io import BytesIO
from pathlib import Path
from urllib.parse import quote
import warnings
import zipfile

import pandas as pd
from fastapi import BackgroundTasks

from src.modules.auth import hash_password
from src.version import __release_date__

def login(client) -> None:
    response = client.post("/login", data={"username": "admin", "password": "admin123"}, follow_redirects=False)
    assert response.status_code == 303


def login_super(client) -> None:
    response = client.post("/login", data={"username": "super", "password": "super123"}, follow_redirects=False)
    assert response.status_code == 303


def test_login_page_loads(client) -> None:
    response = client.get("/login")
    assert response.status_code == 200
    assert "Log in" in response.text
    assert "Dashboard Analytic" in response.text
    assert __release_date__ in response.text
    assert "Default Access:" in response.text
    assert "<strong class=\"login-default-role login-default-role-super-admin\">Role: super-admin</strong>" in response.text
    assert "<strong class=\"login-default-role login-default-role-admin\">Role: admin</strong>" in response.text
    assert "<strong class=\"login-default-role login-default-role-user\">Role: user</strong>" in response.text
    assert 'class="login-workspace-field">Workspace' in response.text
    assert 'data-login-password-toggle' in response.text
    assert '<span class="login-password-editor">' in response.text
    assert 'Default Workspace' in response.text


def test_new_environment_creates_the_three_bootstrap_roles(client) -> None:
    import src.DashboardAnalytic as app_module

    roles = {row['username']: row['role'] for row in app_module.repository.list_users()}
    assert roles['super'] == 'super-admin'
    assert roles['admin'] == 'admin'
    assert roles['demo'] == 'user'


def test_bootstrap_users_are_not_recreated_after_the_first_start(client) -> None:
    import src.DashboardAnalytic as app_module

    with app_module.repository.global_connection() as conn:
        conn.execute("DELETE FROM users WHERE username = ?", ('super',))
        conn.execute("UPDATE users SET username = ? WHERE username = ?", ('renamed-admin', 'admin'))

    app_module.repository.initialize()

    assert app_module.repository.get_user('super') is None
    assert app_module.repository.get_user('admin') is None
    assert app_module.repository.get_user('renamed-admin') is not None
    assert app_module.repository.get_user('demo') is not None


def test_login_page_hides_missing_default_access_accounts(client) -> None:
    import src.DashboardAnalytic as app_module

    with app_module.repository.global_connection() as conn:
        conn.execute("DELETE FROM users WHERE username = ?", ("demo",))

    response = client.get("/login")
    assert response.status_code == 200
    assert "Default Access:" in response.text
    assert "admin / admin123" in response.text
    assert "demo / demo123" not in response.text


def test_login_page_hides_default_access_when_password_differs_from_default(client) -> None:
    import src.DashboardAnalytic as app_module

    with app_module.repository.global_connection() as conn:
        conn.execute(
            "UPDATE users SET password_hash = ? WHERE username = ?",
            (hash_password("changed-password"), "demo"),
        )

    response = client.get("/login")
    assert response.status_code == 200
    assert "Default Access:" in response.text
    assert "admin / admin123" in response.text
    assert "demo / demo123" not in response.text


def test_login_page_hides_default_access_section_when_no_default_users_exist(client) -> None:
    import src.DashboardAnalytic as app_module

    with app_module.repository.global_connection() as conn:
        conn.execute("DELETE FROM users WHERE username IN (?, ?, ?)", ("super", "admin", "demo"))

    response = client.get("/login")
    assert response.status_code == 200
    assert "Default Access:" not in response.text
    assert "admin / admin123" not in response.text
    assert "demo / demo123" not in response.text


def test_login_username_is_case_insensitive_and_rejects_case_duplicates(client) -> None:
    import src.DashboardAnalytic as app_module

    admin = next(row for row in app_module.repository.list_users() if row['username'] == 'admin')
    app_module.repository.set_user_workspace_access(int(admin['id']), ['default'])
    mixed_case_login = client.post(
        '/login',
        data={'username': 'AdMiN', 'password': 'admin123', 'workspace_id': 'default'},
        follow_redirects=False,
    )
    assert mixed_case_login.status_code == 303

    assert app_module.repository.get_user('ADMIN').username == 'admin'
    demo = next(row for row in app_module.repository.list_users() if row['username'] == 'demo')
    app_module.repository.set_user_workspace_access(int(demo['id']), ['default'])
    assert app_module.repository.user_has_workspace_access('DeMo', 'default')

    duplicate = client.post(
        '/admin/users',
        data={'username': 'ADMIN', 'password': 'other-password', 'role': 'user'},
    )
    assert duplicate.status_code == 400
    assert 'already exists' in duplicate.text


def test_login_warns_when_valid_user_cannot_access_selected_workspace(client) -> None:
    response = client.post(
        '/login',
        data={'username': 'demo', 'password': 'demo123', 'workspace_id': 'default'},
    )
    assert response.status_code == 403
    assert 'You do not have access to that workspace.' in response.text
    assert 'class="alert alert-warning" role="alert"' in response.text


def test_only_super_admin_has_implicit_workspace_access(client) -> None:
    import src.DashboardAnalytic as app_module

    denied_admin = client.post(
        '/login',
        data={'username': 'admin', 'password': 'admin123', 'workspace_id': 'default'},
    )
    assert denied_admin.status_code == 403
    assert 'You do not have access to that workspace.' in denied_admin.text

    admin = next(row for row in app_module.repository.list_users() if row['username'] == 'admin')
    app_module.repository.set_user_workspace_access(int(admin['id']), ['default'])
    granted_admin = client.post(
        '/login',
        data={'username': 'admin', 'password': 'admin123', 'workspace_id': 'default'},
        follow_redirects=False,
    )
    assert granted_admin.status_code == 303

    app_module.repository.set_user_workspace_access(int(admin['id']), [])
    super_admin = client.post(
        '/login',
        data={'username': 'super', 'password': 'super123', 'workspace_id': 'default'},
        follow_redirects=False,
    )
    assert super_admin.status_code == 303


def test_admin_import_export_packages_detect_configuration_and_workspaces(client) -> None:
    import src.DashboardAnalytic as app_module

    login_super(client)
    admin_response = client.get('/admin')
    assert admin_response.status_code == 200
    assert 'Import/Export' in admin_response.text
    assert 'Config</option>' in admin_response.text
    assert 'Config + Slides Templates' in admin_response.text
    assert 'Workspace: Default Workspace' in admin_response.text

    config_response = client.get('/admin/import-export/export?export_target=config')
    assert config_response.status_code == 200
    with zipfile.ZipFile(BytesIO(config_response.content)) as archive:
        manifest = json.loads(archive.read('manifest.json'))
        assert 'config/workspace-registry.db' not in archive.namelist()
    assert manifest == {
        'format': 'dashboard-analytic-export',
        'includes_slides_templates': False,
        'kind': 'config',
        'version': 1,
    }

    workspace_response = client.get('/admin/import-export/export?export_target=workspace:default')
    assert workspace_response.status_code == 200
    inspection_response = client.post(
        '/admin/import-export/inspect',
        files={'package': ('default-workspace.zip', BytesIO(workspace_response.content), 'application/zip')},
    )
    assert inspection_response.json() == {
        'kind': 'workspace',
        'includes_slides_templates': False,
        'workspace_collisions': ['Default Workspace'],
    }
    close_response = client.post('/workspace/close', data={'workspace_id': 'default'}, follow_redirects=False)
    assert close_response.status_code == 303
    imported_response = client.post(
        '/admin/import-export/import',
        data={'confirmed_import': 'true'},
        files={'package': ('default-workspace.zip', BytesIO(workspace_response.content), 'application/zip')},
        follow_redirects=False,
    )
    assert imported_response.status_code == 303
    assert 'import_export_notice=' in imported_response.headers['location']
    assert any(workspace.name == 'Default Workspace' for workspace in app_module.workspace_registry.list())

    full_response = client.get('/admin/import-export/export?export_target=full-environment')
    assert full_response.status_code == 200
    with zipfile.ZipFile(BytesIO(full_response.content)) as archive:
        full_manifest = json.loads(archive.read('manifest.json'))
        assert full_manifest['kind'] == 'full-environment'
        assert len(full_manifest['workspaces']) == 1
        assert 'config/workspace-registry.db' not in archive.namelist()
    full_import_response = client.post(
        '/admin/import-export/import',
        data={'confirmed_import': 'true'},
        files={'package': ('full-environment.zip', BytesIO(full_response.content), 'application/zip')},
        follow_redirects=False,
    )
    assert full_import_response.status_code == 303
    assert len(app_module.workspace_registry.list()) == 1


def test_config_import_replaces_global_users_and_preserves_user_ids(client) -> None:
    import src.DashboardAnalytic as app_module

    login_super(client)
    created = client.post(
        '/admin/users',
        data={'username': 'exported-user', 'password': 'exported123', 'role': 'user'},
        follow_redirects=False,
    )
    assert created.status_code == 303
    expected_users = [
        (int(row['id']), row['username'], row['role'], bool(row['active']))
        for row in app_module.repository.list_users()
    ]

    exported = client.get('/admin/import-export/export?export_target=config')
    assert exported.status_code == 200
    with zipfile.ZipFile(BytesIO(exported.content)) as archive:
        assert 'config/application.db' in archive.namelist()

    exported_user = next(row for row in app_module.repository.list_users() if row['username'] == 'exported-user')
    changed_password = client.post(
        f"/admin/users/{exported_user['id']}/update",
        data={'username': 'exported-user', 'password': 'local-change123', 'role': 'user', 'active': '1'},
        follow_redirects=False,
    )
    assert changed_password.status_code == 303

    local_only = client.post(
        '/admin/users',
        data={'username': 'local-only-user', 'password': 'local123', 'role': 'user'},
        follow_redirects=False,
    )
    assert local_only.status_code == 303
    assert app_module.repository.get_user('local-only-user') is not None

    # Simulate a different deployment whose application database contains
    # conflicting records.  Importing configuration must replace it exactly,
    # including user IDs and roles, rather than merge it with the source.
    with app_module.repository.global_connection() as conn:
        conn.execute("UPDATE users SET username = 'destination-user', role = 'admin' WHERE username = 'exported-user'")

    imported = client.post(
        '/admin/import-export/import',
        data={'confirmed_import': 'true'},
        files={'package': ('configuration.zip', BytesIO(exported.content), 'application/zip')},
        follow_redirects=False,
    )
    assert imported.status_code == 303
    assert app_module.repository.get_user('local-only-user') is None
    assert app_module.repository.get_user('destination-user') is None
    restored_user = app_module.repository.get_user('exported-user')
    assert restored_user is not None
    assert app_module.verify_password('exported123', restored_user.password_hash)
    assert [
        (int(row['id']), row['username'], row['role'], bool(row['active']))
        for row in app_module.repository.list_users()
    ] == expected_users


def test_admin_export_job_creates_a_disk_backed_download(client) -> None:
    login_super(client)
    started = client.post('/admin/import-export/export/jobs', data={'export_target': 'config'})
    assert started.status_code == 200
    status_url = started.json()['status_url']
    payload = {}
    for _ in range(100):
        status_response = client.get(status_url)
        assert status_response.status_code == 200
        payload = status_response.json()
        if payload['status'] in {'ready', 'failed'}:
            break
        time.sleep(0.01)
    assert payload['status'] == 'ready'
    assert payload['size'] > 0
    download = client.get(payload['download_url'])
    assert download.status_code == 200
    with zipfile.ZipFile(BytesIO(download.content)) as archive:
        assert json.loads(archive.read('manifest.json'))['kind'] == 'config'


def test_admin_import_export_is_limited_to_slides_templates(client) -> None:
    login(client)

    panel = client.get('/admin')
    assert panel.status_code == 200
    assert 'data-panel-state-key="admin:import-export"' in panel.text
    assert 'Slides Templates</option>' in panel.text
    assert 'value="config" disabled>Config</option>' in panel.text
    assert 'value="workspace:default" disabled>Workspace: Default Workspace</option>' in panel.text

    blocked_export = client.get('/admin/import-export/export?export_target=config')
    assert blocked_export.status_code == 403

    templates_export = client.get('/admin/import-export/export?export_target=slides-templates')
    assert templates_export.status_code == 200
    with zipfile.ZipFile(BytesIO(templates_export.content)) as archive:
        assert json.loads(archive.read('manifest.json'))['kind'] == 'slides-templates'

    config_package = BytesIO()
    with zipfile.ZipFile(config_package, 'w') as archive:
        archive.writestr('manifest.json', json.dumps({
            'format': 'dashboard-analytic-export', 'version': 1, 'kind': 'config',
        }))
    config_package.seek(0)
    blocked_import = client.post(
        '/admin/import-export/inspect',
        files={'package': ('config.zip', config_package, 'application/zip')},
    )
    assert blocked_import.status_code == 403


def test_admin_can_login_upload_and_see_automatic_dashboard(client) -> None:
    login(client)
    csv_content = b"market,period,score,gap\nES,2026-Q1,91,2.1\nES,2026-Q1,87,3.3\nDE,2026-Q2,76,5.2\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    dashboard_response = client.get(upload_response.headers["location"])
    assert dashboard_response.status_code == 200
    assert "sample.csv" in dashboard_response.text
    assert "Data Ingestion" in dashboard_response.text
    assert "Workspace" in dashboard_response.text
    assert "Workspace opened from cache" in dashboard_response.text

    analysis_redirect = client.post(
        "/dashboard/analyze",
        data={
            "dataset_id": 1,
            "metric": "score",
            "market": "ES",
            "period": "2026-Q1",
            "aggregation": "all",
        },
        follow_redirects=False,
    )
    assert analysis_redirect.status_code == 303
    filtered_dashboard = client.get(analysis_redirect.headers["location"])
    assert filtered_dashboard.status_code == 200
    assert "Processed Metrics" in filtered_dashboard.text
    assert "CDF Curve" in filtered_dashboard.text
    assert "89" in filtered_dashboard.text


def test_dashboard_disables_metrics_without_non_null_values(client) -> None:
    login(client)
    csv_content = (
        b"market,period,operator,region,latency_ms,score\n"
        b"ES,2026-Q1,Vodafone,North,,91\n"
        b"ES,2026-Q1,Orange,South,,87\n"
    )
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert response.status_code == 200
    assert 'value="score"' in response.text
    assert 'value="latency_ms" disabled' in response.text
    assert "data-table-wrap" in response.text
    assert "Global Aggregation" in response.text


def test_admin_can_retry_stuck_dataset(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )

    import src.DashboardAnalytic as app_module

    app_module.repository.update_dataset_profile(1, status="failed", progress=100, dataset_kind=None, row_count=None, column_count=None, default_metric=None)
    retry_response = client.post("/dashboard/retry/1", follow_redirects=False)
    assert retry_response.status_code == 303

    dashboard_response = client.get(retry_response.headers["location"])
    assert dashboard_response.status_code == 200
    assert "Workspace opened from cache" in dashboard_response.text


def test_admin_cannot_retry_queued_dataset(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )

    response = client.post("/dashboard/retry/1")
    assert response.status_code == 400
    assert "Only failed or stopped datasets can be retried" in response.text


def test_admin_can_delete_queued_dataset(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )

    import src.DashboardAnalytic as app_module

    dataset = app_module.repository.get_dataset(1)
    assert dataset is not None
    dataset_path = Path(dataset["stored_path"])
    assert dataset_path.exists()
    response = client.post("/dashboard/delete/1", follow_redirects=False)
    assert response.status_code == 303
    assert app_module.repository.get_dataset(1) is None
    assert not dataset_path.exists()


def test_admin_can_stop_processing_dataset(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )

    import src.DashboardAnalytic as app_module

    app_module.repository.update_dataset_profile(1, status="processing", progress=33)
    response = client.post("/dashboard/stop/1", follow_redirects=False)
    assert response.status_code == 303

    dataset = app_module.repository.get_dataset(1)
    assert dataset is not None
    assert dataset["status"] == "stopped"


def test_reupload_same_file_reuses_existing_dataset_entry(client) -> None:
    login(client)
    payload = b"market,period,score\nES,2026-Q1,91\n"
    first_upload = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(payload), "text/csv")},
        follow_redirects=False,
    )
    second_upload = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(payload), "text/csv")},
        follow_redirects=False,
    )
    assert first_upload.status_code == 303
    assert second_upload.status_code == 303

    import src.DashboardAnalytic as app_module

    datasets = app_module.repository.list_datasets()
    assert len(datasets) == 1


def test_reupload_preserves_original_upload_date_for_dataset_ordering(client) -> None:
    login(client)
    payload = b"market,period,score\nES,2026-Q1,91\n"
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(payload), "text/csv")},
        follow_redirects=False,
    )
    import src.DashboardAnalytic as app_module

    original_upload = '2025-01-02 03:04:05'
    with app_module.repository.connection() as conn:
        conn.execute('UPDATE datasets SET uploaded_at = ? WHERE id = 1', (original_upload,))
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(payload), "text/csv")},
        follow_redirects=False,
    )

    dataset = app_module.repository.get_dataset(1)
    assert dataset is not None
    assert dataset['uploaded_at'] == original_upload

    workspace = client.get('/workspace')
    assert workspace.text.index('<th>ID</th>') < workspace.text.index('<th>Dataset</th>')
    assert '<td data-queue-id>1</td>' in workspace.text
    assert '<th>Uploaded</th>' in workspace.text
    assert '<th>Updated</th>' in workspace.text
    assert workspace.text.index('<th>Uploaded</th>') < workspace.text.index('<th>Updated</th>')


def test_workspace_management_saves_user_access_for_the_selected_workspace(client) -> None:
    import src.DashboardAnalytic as app_module

    login_super(client)
    created = client.post('/workspace/create', data={'name': 'Germany'}, follow_redirects=False)
    assert created.status_code == 303
    germany = next(item for item in app_module.workspace_registry.list() if item.name == 'Germany')
    demo = next(row for row in app_module.repository.list_users() if row['username'] == 'demo')

    response = client.post(
        '/workspace/access',
        data={'workspace_id': germany.id, 'usernames': [str(demo['username'])]},
        headers={'X-Requested-With': 'XMLHttpRequest'},
    )

    assert response.status_code == 200
    assert response.json() == {'ok': True, 'notice': 'Workspace access updated.'}
    assert app_module.repository.list_user_workspace_ids(int(demo['id'])) == [germany.id]
    users_page = client.get('/admin')
    demo_row = users_page.text.split(f'aria-label="Filter workspaces for {demo["username"]}"', 1)[1].split('</details>', 1)[0]
    assert f'value="{germany.id}"' in demo_row
    assert 'checked' in demo_row


def test_workspace_management_isolates_dataset_databases_and_remembers_last_opened_workspace(client) -> None:
    import src.DashboardAnalytic as app_module

    login(client)
    admin = next(row for row in app_module.repository.list_users() if row['username'] == 'admin')
    app_module.repository.set_user_workspace_access(int(admin['id']), ['default'])
    assert app_module.workspace_registry.registry_path.name == 'workspace-registry.db'
    assert app_module.workspace_registry.registry_path == app_module.settings.input_dir.parent.parent / 'workspace-registry.db'
    payload = b"market,period,score\nES,2026-Q1,91\n"
    client.post(
        '/dashboard/upload', data={'dataset_kinds': 'data'},
        files={'dataset_files': ('default.csv', BytesIO(payload), 'text/csv')},
    )
    default_db = app_module.repository.db_path
    assert default_db.parent == app_module.settings.input_dir.parent
    assert default_db.name == 'Default Workspace.db'

    created = client.post('/workspace/create', data={'name': 'Campaign benchmark'}, follow_redirects=False)
    assert created.status_code == 303
    assert app_module.active_workspace is not None
    assert app_module.active_workspace.name == 'Campaign benchmark'
    assert app_module.repository.db_path != default_db
    assert app_module.repository.db_path.name == 'Campaign benchmark.db'
    assert app_module.repository.db_path.parent.name == 'Campaign benchmark'
    assert app_module.repository.list_datasets() == []

    page = client.get('/workspace')
    assert 'Campaign benchmark' in page.text
    assert 'Manage workspaces' in page.text
    assert 'data-workspace-open disabled>Open</button>' in page.text
    assert page.text.index('>Campaign benchmark</option>') < page.text.index('>Default Workspace</option>')
    workspace_id = app_module.active_workspace.id
    renamed = client.post('/workspace/rename', data={'workspace_id': workspace_id, 'name': 'Campaign benchmark Q3'})
    assert renamed.status_code == 200
    assert 'Campaign benchmark Q3' in renamed.text
    assert app_module.repository.db_path.name == 'Campaign benchmark Q3.db'
    assert app_module.repository.db_path.parent.name == 'Campaign benchmark Q3'
    assert app_module.repository.db_path.exists()

    selected = client.post('/workspace/select', data={'workspace_id': 'default'}, follow_redirects=False)
    assert selected.status_code == 303
    assert app_module.repository.db_path == default_db
    assert len(app_module.repository.list_datasets()) == 1

    duplicated = client.post('/workspace/duplicate', data={'workspace_id': 'default'}, follow_redirects=False)
    assert duplicated.status_code == 303
    copied_workspace = next(item for item in app_module.workspace_registry.list() if item.name == 'Default Workspace - Copy')
    assert copied_workspace.database_path.name == 'Default Workspace - Copy.db'
    with app_module.repository.connection() as conn:
        assert conn.execute('SELECT COUNT(*) FROM datasets').fetchone()[0] == 1

    cannot_remove_open = client.post('/workspace/delete', data={'workspace_id': 'default'}, follow_redirects=False)
    assert cannot_remove_open.headers['location'].startswith('/workspace?workspace_warning=')

    closed = client.post('/workspace/close', data={'workspace_id': 'default'}, follow_redirects=False)
    assert closed.status_code == 303
    assert app_module.active_workspace is None
    closed_workspace = client.get('/workspace')
    assert 'Data Ingestion' not in closed_workspace.text
    assert 'module-tab-disabled' in closed_workspace.text
    assert client.get('/dashboard', follow_redirects=False).status_code == 303
    login_page = client.get('/login')
    assert login_page.text.index('>Campaign benchmark Q3</option>') < login_page.text.index('>Default Workspace</option>')
    assert '<option value="default" selected>Default Workspace</option>' in login_page.text

    deleted = client.post('/workspace/delete', data={'workspace_id': workspace_id}, follow_redirects=False)
    assert deleted.status_code == 303
    assert app_module.workspace_registry.get(workspace_id) is None


def test_queued_import_continues_after_its_workspace_is_closed(client) -> None:
    import src.DashboardAnalytic as app_module

    login(client)
    source_path = app_module.settings.input_dir / 'continue-after-close.csv'
    source_path.write_bytes(b'market,period,score\nES,2026-Q1,91\n')
    dataset_id, _ = app_module.repository.add_dataset(source_path.name, str(source_path), 'admin')
    app_module.repository.update_dataset_profile(dataset_id, dataset_kind='data')
    workspace_database = app_module.repository.db_path

    tasks = BackgroundTasks()
    app_module.enqueue_dataset_processing(tasks, dataset_id, source_path, 'admin')
    closed = client.post('/workspace/close', data={'workspace_id': 'default'}, follow_redirects=False)
    assert closed.status_code == 303
    assert app_module.active_workspace is None

    queued_task = tasks.tasks[0]
    queued_task.func(*queued_task.args, **queued_task.kwargs)

    completed = app_module.Repository(workspace_database).get_dataset(dataset_id)
    assert completed is not None
    assert completed['status'] == 'ready'


def test_admin_panel_is_available_for_admin(client) -> None:
    login(client)
    response = client.get("/admin")
    assert response.status_code == 200
    assert "Admin panel" in response.text
    assert 'value="super-admin" aria-label="Role for super" readonly' in response.text
    assert 'name="username" value="super" form="user-update-1" autocomplete="off" readonly data-user-autofill-guard required disabled' in response.text
    assert 'name="password" value="" placeholder="••••••••" autocomplete="new-password" readonly form="user-update-1" data-user-password data-user-autofill-guard disabled' in response.text
    assert 'data-user-password-toggle' in response.text
    assert '<svg viewBox="0 0 24 24" aria-hidden="true"><path d="M2.2 12' in response.text
    assert 'Keep current' not in response.text
    assert 'restoreUsersFromServerMarkup' in response.text
    assert 'form="user-update-1" disabled title="Only super-admins can modify super-admin accounts">Save</button>' not in response.text


def test_login_and_admin_remain_available_after_closing_the_active_workspace(client) -> None:
    import src.DashboardAnalytic as app_module

    login(client)
    app_module.close_active_workspace()

    login_page = client.get('/login')
    assert login_page.status_code == 200

    admin = client.get('/admin')
    assert admin.status_code == 200
    assert 'Open a workspace from Workspace Management before managing its datasets.' in admin.text
    assert 'Open a workspace from Workspace Management before viewing or editing its database.' in admin.text


def test_admin_database_management_lists_and_updates_active_workspace_tables(client) -> None:
    import src.DashboardAnalytic as app_module

    login(client)
    client.post(
        "/admin/users",
        data={"username": "database-editor", "password": "start123", "role": "user"},
        follow_redirects=False,
    )
    app_module.repository.replace_dataset_rows(987, pd.DataFrame({"obsolete": ["row"]}))
    app_module.repository.replace_reporting_rows(987, 'data', pd.DataFrame({"Campaign": ["legacy"]}))
    assert app_module.repository.dataset_rows_table_exists(987)
    admin = client.get("/admin")
    assert admin.status_code == 200
    assert "Database Management" in admin.text
    assert 'data-database-table-select' in admin.text
    assert 'Workspace: Default Workspace' in admin.text
    assert 'Clean orphaned rows' in admin.text
    assert 'value="users"' in admin.text
    assert app_module.repository.dataset_rows_table_exists(987)
    assert app_module.repository.database_table_page('reporting_rows_data')['total_rows'] == 1
    cleanup = client.post('/admin/database/cleanup', follow_redirects=False)
    assert cleanup.status_code == 303
    assert not app_module.repository.dataset_rows_table_exists(987)
    assert app_module.repository.database_table_page('reporting_rows_data')['total_rows'] == 0

    users = client.get("/admin/database/table", params={"table": "users", "limit": 100})
    assert users.status_code == 200
    payload = users.json()
    assert any(column["name"] == "id" and column["primary_key"] for column in payload["columns"])
    editor = next(row for row in payload["rows"] if row["username"] == "database-editor")

    values = client.get(
        "/admin/database/table/values",
        params={"table": "users", "column": "username", "search": "database-editor"},
    )
    assert values.status_code == 200
    assert values.json()["values"] == ["database-editor"]

    filtered = client.post(
        "/admin/database/table/query",
        json={"table": "users", "offset": 0, "limit": 100, "filters": {"username": ["database-editor"]}},
    )
    assert filtered.status_code == 200
    assert [row["username"] for row in filtered.json()["rows"]] == ["database-editor"]
    assert filtered.json()["total_rows"] == 1
    assert filtered.json()["all_rows"] > filtered.json()["total_rows"]

    saved = client.post(
        "/admin/database/table",
        json={"table": "users", "rowid": editor["__database_rowid__"], "updates": {"active": "0"}},
    )
    assert saved.status_code == 200
    assert app_module.repository.get_user("database-editor").active is False

    protected = client.post(
        "/admin/database/table",
        json={"table": "users", "rowid": editor["__database_rowid__"], "updates": {"id": "999"}},
    )
    assert protected.status_code == 400
    assert "Primary-key values cannot be edited" in protected.json()["detail"]

    deleted = client.post(
        "/admin/database/table/delete",
        json={"table": "users", "rowid": editor["__database_rowid__"]},
    )
    assert deleted.status_code == 200
    assert app_module.repository.get_user("database-editor") is None


def test_admin_dataset_management_renames_dataset_file_and_materialised_source_labels(client) -> None:
    import src.DashboardAnalytic as app_module

    login(client)
    upload = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'data'},
        files={'dataset_files': ('original-cdr.csv', BytesIO(b'Campaign,LQ\nUK_Q2_SA_2026,3.8\n'), 'text/csv')},
        follow_redirects=False,
    )
    assert upload.status_code == 303
    before = app_module.repository.get_dataset(1)
    assert before is not None
    old_path = Path(before['stored_path'])
    app_module.repository.copy_dataset_rows_to_reporting(1, 'data', ['source_file'])

    renamed = client.post(
        '/admin/datasets/1/rename', data={'file_name': 'renamed-cdr.csv'}, follow_redirects=False,
    )
    assert renamed.status_code == 303
    after = app_module.repository.get_dataset(1)
    assert after is not None
    new_path = Path(after['stored_path'])
    assert after['file_name'] == 'renamed-cdr.csv'
    assert new_path.name == 'renamed-cdr.csv'
    assert not old_path.exists()
    assert new_path.exists()

    with app_module.repository.connection() as conn:
        dataset_source = conn.execute('SELECT DISTINCT source_file FROM "dataset_rows_1"').fetchall()
        reporting_source = conn.execute(
            'SELECT DISTINCT source_file FROM "reporting_rows_data" WHERE dataset_id = 1'
        ).fetchall()
    assert [row['source_file'] for row in dataset_source] == ['renamed-cdr.csv']
    assert [row['source_file'] for row in reporting_source] == ['renamed-cdr.csv']

    workspace = client.get('/workspace')
    assert 'data-dataset-name-editor' not in workspace.text
    admin = client.get('/admin')
    assert 'Datasets Management' in admin.text
    assert '<th>Uploaded</th>' in admin.text
    assert '<th>Updated</th>' in admin.text
    assert admin.text.index('<th>Uploaded</th>') < admin.text.index('<th>Updated</th>')
    assert 'dataset-rename-1' in admin.text
    assert 'data-admin-dataset-rename-save' in admin.text
    assert 'Save name' not in admin.text
    assert 'Show Dashboard' in admin.text
    assert 'Preview' in admin.text


def test_dashboard_upload_accepts_multiple_files(client) -> None:
    login(client)

    response = client.post(
        "/dashboard/upload",
        files=[
            ("dataset_files", ("sample-a.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")),
            ("dataset_files", ("sample-b.csv", BytesIO(b"market,period,score\nDE,2026-Q2,78\n"), "text/csv")),
        ],
        follow_redirects=False,
    )
    assert response.status_code == 303

    import src.DashboardAnalytic as app_module

    datasets = app_module.repository.list_datasets()
    assert len(datasets) == 2


def test_workspace_upload_persists_selected_dataset_kind(client) -> None:
    login(client)

    response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "mapping_vodafone"},
        files={"dataset_files": ("operator_cells.csv", BytesIO(b"Cell ID,OP/ Vendor\n123,Ericsson\n"), "text/csv")},
        follow_redirects=False,
    )
    assert response.status_code == 303

    import src.DashboardAnalytic as app_module

    dataset = app_module.repository.get_dataset(1)
    assert dataset is not None
    assert dataset["dataset_kind"] == "mapping_vodafone"


def test_workspace_preview_and_cdr_dashboard_action(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("cdr_data.csv", BytesIO(b"operator,score\nVodafone UK,91\n"), "text/csv")},
        follow_redirects=False,
    )

    workspace_response = client.get("/workspace")
    assert workspace_response.status_code == 200
    assert 'data-queue-type-filter' in workspace_response.text
    assert 'value="">All Types' in workspace_response.text
    assert 'href="/workspace/preview/1" target="_blank" rel="noopener" data-preview-open-link data-loading-label="Generating dataset preview">Preview</a>' in workspace_response.text
    assert 'Show Dashboard</a>' in workspace_response.text

    preview_response = client.get("/workspace/preview/1")
    assert preview_response.status_code == 200
    assert "Dataset preview" in preview_response.text
    assert "Vodafone UK" in preview_response.text
    assert "Show Dashboard" in preview_response.text
    assert 'name="row_limit" value="100"' in preview_response.text
    assert 'data-preview-column-filter' in preview_response.text
    assert 'data-preview-row-filter' in preview_response.text
    assert 'data-preview-filter-table' in preview_response.text
    preview_script = client.get('/static/js/app.js')
    assert preview_script.status_code == 200
    assert 'preview-column-filter-trigger' in preview_script.text
    assert 'data-preview-value-option' in preview_script.text

    limited_preview_response = client.get("/workspace/preview/1?row_limit=25")
    assert limited_preview_response.status_code == 200
    assert 'name="row_limit" value="25"' in limited_preview_response.text

    dashboard_response = client.get('/dashboard?dataset_id=1&input_kind=data')
    assert dashboard_response.status_code == 200
    assert 'href="/workspace/preview/1" target="_blank" rel="noopener" data-preview-open-link data-loading-label="Generating dataset preview">Preview Dataset</a>' in dashboard_response.text


def test_cdr_preview_highlights_vendor_and_filters_cdr_dimensions(client) -> None:
    login(client)
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'data'},
        files={'dataset_files': (
            'cdr_data.csv',
            BytesIO(
                b'operator,vendor,RAT_A,Session_Type,Call_Status,score\n'
                b'Vodafone UK,Ericsson,ENDC,VoLTE,Completed,91\n'
                b'3,Nokia,NR,WhatsApp,Dropped,90\n'
            ),
            'text/csv',
        )},
        follow_redirects=False,
    )

    default_preview = client.get('/workspace/preview/1')
    assert default_preview.status_code == 200
    assert '<option value="Vodafone UK" selected>' in default_preview.text
    assert '<option value="3" selected>' in default_preview.text
    assert '<option value="Ericsson" selected>' in default_preview.text
    assert '<option value="Nokia" selected>' in default_preview.text

    preview = client.get(
        '/workspace/preview/1?cdr_operator=3&cdr_vendor=Nokia&cdr_rat=NR'
        '&cdr_session_type=WhatsApp&cdr_call_status=Dropped',
    )
    assert preview.status_code == 200
    assert 'name="cdr_operator"' in preview.text
    assert 'name="cdr_vendor"' in preview.text
    assert 'name="cdr_rat"' in preview.text
    assert 'name="cdr_session_type"' in preview.text
    assert 'name="cdr_call_status"' in preview.text
    assert 'class="vendor-column">vendor<' in preview.text
    preview_rows = preview.text.split('<tbody>', 1)[1].split('</tbody>', 1)[0]
    assert '>Nokia<' in preview_rows
    assert '>Ericsson<' not in preview_rows

    multi_preview = client.get('/workspace/preview/1?cdr_operator=Vodafone%20UK&cdr_operator=3')
    multi_rows = multi_preview.text.split('<tbody>', 1)[1].split('</tbody>', 1)[0]
    assert '>Nokia<' in multi_rows
    assert '>Ericsson<' in multi_rows
    assert 'name="cdr_operator" multiple' in multi_preview.text


def test_workspace_uses_persisted_vendor_flags_without_reloading_cdr_files(client, monkeypatch) -> None:
    login(client)
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'data'},
        files={'dataset_files': ('cdr_data.csv', BytesIO(b'operator,score\nVodafone UK,91\n'), 'text/csv')},
        follow_redirects=False,
    )
    import src.DashboardAnalytic as app_module

    def source_reload_should_not_run(*_args, **_kwargs):
        raise AssertionError('Workspace should use the persisted Vendor flags, not reload CDR files.')

    monkeypatch.setattr(app_module, 'load_cached_dataset', source_reload_should_not_run)
    assert client.get('/workspace').status_code == 200


def test_workspace_maps_unassigned_cdr_vendors_from_available_multivendor_mapping(client) -> None:
    login(client)
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'data'},
        files={'dataset_files': ('cdr_data.csv', BytesIO(b'operator,Cell_ID_A,score\n3,200 -> 200,91\n'), 'text/csv')},
        follow_redirects=False,
    )
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_three'},
        files={'dataset_files': ('Multivendor_Mapping_3UK.csv', BytesIO(b'Cid__ECI,Vendor\n200,Nokia\n'), 'text/csv')},
        follow_redirects=False,
    )

    workspace = client.get('/workspace')
    assert 'Map Vendors</button>' in workspace.text
    assert 'data-queue-status="ready"' in workspace.text
    assert 'name="cdr_dataset_ids"' in workspace.text
    assert 'name="three_mapping_dataset_id"' in workspace.text
    assert 'value="2" selected' in workspace.text
    assert 'data-loading-label="Mapping Vendors to CDR samples"' not in workspace.text
    assert 'Vendor mapping rule' in workspace.text
    assert 'the same non-empty Vendor at both endpoints returns' in workspace.text
    live_status = client.get('/api/datasets/status').json()['datasets']
    assert next(dataset for dataset in live_status if dataset['id'] == 1)['can_map_vendors'] is True

    response = client.post(
        '/workspace/map-vendors',
        data={'cdr_dataset_id': 1, 'three_mapping_dataset_id': 2},
        follow_redirects=False,
    )
    assert response.status_code == 303

    preview = client.get('/workspace/preview/1')
    assert preview.status_code == 200
    assert '>3_Nokia<' in preview.text

    workspace_after_mapping = client.get('/workspace').text.split('<tbody>', 1)[1].split('</tbody>', 1)[0]
    assert 'data-dataset-id="1"' in workspace_after_mapping
    assert 'Map Vendors</button>' not in workspace_after_mapping
    assert 'Clear Vendors</button>' in workspace_after_mapping
    assert 'data-vendor-clear-open' in workspace_after_mapping
    assert 'action="/workspace/clear-vendors"' in workspace.text
    live_status_after_mapping = client.get('/api/datasets/status').json()['datasets']
    assert next(dataset for dataset in live_status_after_mapping if dataset['id'] == 1)['can_clear_vendors'] is True

    clear_response = client.post('/workspace/clear-vendors/1', follow_redirects=False)
    assert clear_response.status_code == 303
    workspace_after_clear = client.get('/workspace').text.split('<tbody>', 1)[1].split('</tbody>', 1)[0]
    assert 'Map Vendors</button>' in workspace_after_clear
    assert 'Clear Vendors</button>' not in workspace_after_clear


def test_workspace_queues_vendor_mapping_for_multiple_cdrs(client) -> None:
    login(client)
    uploads = [
        ('cdr_data_q1.csv', 'data', b'operator,Cell_ID_A,Campaign\n3,200 -> 200,2026 Q1\n'),
        ('cdr_data_q2.csv', 'data', b'operator,Cell_ID_A,Campaign\n3,200 -> 200,2026 Q2\n'),
        ('Multivendor_Mapping_3UK.csv', 'mapping_three', b'Cid__ECI,Vendor\n200,Nokia\n'),
    ]
    for file_name, kind, content in uploads:
        response = client.post(
            '/dashboard/upload',
            data={'dataset_kinds': kind},
            files={'dataset_files': (file_name, BytesIO(content), 'text/csv')},
            follow_redirects=False,
        )
        assert response.status_code == 303

    response = client.post(
        '/workspace/map-vendors',
        data={'cdr_dataset_ids': ['1', '2'], 'three_mapping_dataset_id': '3'},
        follow_redirects=False,
    )
    assert response.status_code == 303
    datasets = client.get('/api/datasets/status').json()['datasets']
    assert all(dataset['vendor_mapping_applied'] for dataset in datasets if dataset['id'] in {1, 2})

    response = client.post(
        '/workspace/clear-vendors',
        data={'cdr_dataset_ids': ['1', '2']},
        follow_redirects=False,
    )
    assert response.status_code == 303
    datasets = client.get('/api/datasets/status').json()['datasets']
    assert all(not dataset['vendor_mapping_applied'] for dataset in datasets if dataset['id'] in {1, 2})


def test_failed_vendor_mapping_keeps_the_cdr_available(client) -> None:
    login(client)
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'data'},
        files={'dataset_files': ('cdr_without_cell_id.csv', BytesIO(b'Operator,score\n3,91\n'), 'text/csv')},
        follow_redirects=False,
    )
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_three'},
        files={'dataset_files': ('Multivendor_Mapping_3UK.csv', BytesIO(b'Cid__ECI,Vendor\n200,Nokia\n'), 'text/csv')},
        follow_redirects=False,
    )

    response = client.post('/workspace/map-vendors', data={'cdr_dataset_id': 1, 'three_mapping_dataset_id': 2})
    assert response.status_code == 200
    dataset = next(item for item in client.get('/api/datasets/status').json()['datasets'] if item['id'] == 1)
    assert dataset['status'] == 'ready'
    assert client.get('/workspace/preview/1').status_code == 200
    assert 'Vendor mapping failed for CDR dataset 1' in client.get('/workspace?dataset_id=1').text


def test_workspace_recovers_legacy_vendor_mapping_failures(client) -> None:
    login(client)
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'data'},
        files={'dataset_files': ('legacy_cdr.csv', BytesIO(b'Operator,score\n3,91\n'), 'text/csv')},
        follow_redirects=False,
    )
    import src.DashboardAnalytic as app_module

    app_module.repository.update_dataset_profile(
        1,
        status='failed',
        progress=100,
        last_error='The selected CDR must contain Operator and Cell_ID_A to assign vendors.',
    )

    response = client.get('/workspace')
    assert response.status_code == 200
    dataset = next(item for item in client.get('/api/datasets/status').json()['datasets'] if item['id'] == 1)
    assert dataset['status'] == 'ready'
    assert client.get('/workspace/preview/1').status_code == 200


def test_workspace_upload_can_map_selected_cdr_vendor_during_processing(client) -> None:
    login(client)
    mapping_response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_three'},
        files={'dataset_files': ('Multivendor_Mapping_3UK.csv', BytesIO(b'Cid__ECI,Vendor\n200,Nokia\n'), 'text/csv')},
        follow_redirects=False,
    )
    assert mapping_response.status_code == 303

    workspace = client.get('/workspace')
    assert 'data-three-mapping-options=' in workspace.text
    assert 'Multivendor_Mapping_3UK.csv' in workspace.text
    assert 'No Map Vendor Column' in workspace.text
    assert 'name = \'three_mapping_dataset_ids\'' not in workspace.text
    assert "mappingSelect.name = fieldName" in workspace.text

    cdr_response = client.post(
        '/dashboard/upload',
        data={
            'dataset_kinds': 'data',
            'vodafone_mapping_dataset_ids': '',
            'three_mapping_dataset_ids': '1',
        },
        files={'dataset_files': ('cdr_data.csv', BytesIO(b'Operator,Cell_ID_A,score\n3,200 -> 200,91\n'), 'text/csv')},
        follow_redirects=False,
    )
    assert cdr_response.status_code == 303

    preview = client.get('/workspace/preview/2')
    assert preview.status_code == 200
    assert '>3_Nokia<' in preview.text

    import src.DashboardAnalytic as app_module
    dataset = app_module.serialize_dataset_row(app_module.repository.get_dataset(2))
    assert dataset['vendor_mapping_applied'] is True


def test_workspace_batch_upload_keeps_vendor_mapping_choices_aligned_per_file(client) -> None:
    login(client)
    client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_three'},
        files={'dataset_files': ('Multivendor_Mapping_3UK.csv', BytesIO(b'Cid__ECI,Vendor\n200,Nokia\n'), 'text/csv')},
        follow_redirects=False,
    )

    response = client.post(
        '/dashboard/upload',
        data={
            'dataset_kinds': ['data', 'generic'],
            'vodafone_mapping_dataset_ids': ['', ''],
            'three_mapping_dataset_ids': ['1', ''],
        },
        files=[
            ('dataset_files', ('cdr_data.csv', BytesIO(b'Operator,Cell_ID_A,score\n3,200 -> 200,91\n'), 'text/csv')),
            ('dataset_files', ('other.csv', BytesIO(b'name,value\nother,1\n'), 'text/csv')),
        ],
        follow_redirects=False,
    )
    assert response.status_code == 303
    assert '>3_Nokia<' in client.get('/workspace/preview/2').text


def test_vfuk_preview_limits_mapping_sheets_and_displays_materialised_gcid(client) -> None:
    login(client)
    workbook = BytesIO()
    with pd.ExcelWriter(workbook, engine='openpyxl') as writer:
        pd.DataFrame({
            'eNodeB ID': [13008],
            'Local Cell ID': [1],
            'OP/ Vendor': ['Samsung'],
        }).to_excel(writer, sheet_name='4G', index=False)
        pd.DataFrame({
            'gNodeB ID': [53986],
            'Local Cell ID': [302],
            'OP/ Vendor': ['Samsung'],
        }).to_excel(writer, sheet_name='5G', index=False)
        pd.DataFrame({'Cell ID': [1]}).to_excel(writer, sheet_name='2G', index=False)
    workbook.seek(0)

    response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_vodafone'},
        files={'dataset_files': ('Multivendor_Mapping_VFUK.xlsx', workbook, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')},
        follow_redirects=False,
    )
    assert response.status_code == 303

    # Simulate a mapping that was processed before the GCID materialisation
    # was introduced; opening its preview must upgrade the stored rows.
    import src.DashboardAnalytic as app_module
    stored_rows = app_module.repository.load_dataset_rows(
        1,
        app_module.repository.list_dataset_row_columns(1),
        {},
    ).drop(columns=['GCID'])
    app_module.repository.replace_dataset_rows(1, stored_rows)

    default_preview = client.get('/workspace/preview/1')
    assert default_preview.status_code == 200
    assert 'name="source_sheet"' in default_preview.text
    assert '<option value="4G" selected>4G</option>' in default_preview.text
    assert '<option value="5G">5G</option>' in default_preview.text
    assert '2G' not in default_preview.text
    assert '3330049' in default_preview.text
    assert '>source_sheet<' not in default_preview.text
    assert 'class="gcid-column"' in default_preview.text

    five_g_preview = client.get('/workspace/preview/1?source_sheet=5G')
    assert five_g_preview.status_code == 200
    assert '<option value="5G" selected>5G</option>' in five_g_preview.text
    assert '221126958' in five_g_preview.text


def test_three_mapping_preview_excludes_empty_normalized_columns(client) -> None:
    login(client)
    response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_three'},
        files={'dataset_files': ('Multivendor_Mapping_3UK.csv', BytesIO(b'MBNL_ID,Vendor,Site_Name,Cid__ECI\nAAB013,Ericsson,United Reformed Church,123\n'), 'text/csv')},
        follow_redirects=False,
    )
    assert response.status_code == 303

    preview = client.get('/workspace/preview/1')
    assert preview.status_code == 200
    assert 'MBNL_ID' in preview.text
    assert 'Cid__ECI' in preview.text
    assert 'Vendor' in preview.text
    assert '>operator<' not in preview.text
    assert '>vendor__2<' not in preview.text
    assert '>technology_primary<' not in preview.text
    assert '>GCID<' in preview.text
    assert '>123<' in preview.text
    assert preview.text.index('>GCID<') < preview.text.index('>MBNL_ID<')
    assert 'class="vendor-column">Vendor<' in preview.text


def test_mapping_preview_shows_every_source_column(client) -> None:
    login(client)
    source_columns = ['CId___ECI', 'Vendor', *(f'Extra_{number:02d}' for number in range(1, 27))]
    source_row = ['123', 'Ericsson', *(str(number) for number in range(1, 27))]
    content = (','.join(source_columns) + '\n' + ','.join(source_row) + '\n').encode()
    response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_three'},
        files={'dataset_files': ('Multivendor_Mapping_3UK.csv', BytesIO(content), 'text/csv')},
        follow_redirects=False,
    )
    assert response.status_code == 303

    preview = client.get('/workspace/preview/1')
    assert preview.status_code == 200
    assert '>Extra_26<' in preview.text


def test_mapping_preview_formats_integral_gcid_without_decimal_suffix(client) -> None:
    login(client)
    response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_vodafone'},
        files={'dataset_files': ('Multivendor_Mapping_VFUK.csv', BytesIO(b'source_sheet,gNodeB ID,Local Cell ID,OP/ Vendor\n5G,53986,302,Ericsson\n5G,,,Ericsson\n'), 'text/csv')},
        follow_redirects=False,
    )
    assert response.status_code == 303

    preview = client.get('/workspace/preview/1?source_sheet=5G')
    assert preview.status_code == 200
    assert '>221126958<' in preview.text
    assert '>221126958.0<' not in preview.text


def test_mapping_preview_hides_unnamed_columns_but_keeps_cell_name(client) -> None:
    login(client)
    response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_vodafone'},
        files={'dataset_files': ('Multivendor_Mapping_VFUK.csv', BytesIO(b'source_sheet,eNodeB ID,Local Cell ID,Cell Name,Unnamed_2,OP/ Vendor\n4G,13008,1,Cell A,,Samsung\n'), 'text/csv')},
        follow_redirects=False,
    )
    assert response.status_code == 303

    preview = client.get('/workspace/preview/1')
    assert preview.status_code == 200
    assert '>Cell Name<' in preview.text
    assert '>Cell A<' in preview.text
    assert '>Unnamed_2<' not in preview.text
    assert '>source_sheet<' not in preview.text


def test_vfuk_preview_uses_only_the_selected_source_sheet_columns(client) -> None:
    login(client)
    workbook = BytesIO()
    with pd.ExcelWriter(workbook, engine='openpyxl') as writer:
        pd.DataFrame({
            'Cell Name': ['4G Cell'],
            'eNodeB ID': [13008],
            'Local Cell ID': [1],
            'OP/ Vendor': ['Samsung'],
        }).to_excel(writer, sheet_name='4G', index=False)
        pd.DataFrame({
            'Cell Name': ['5G Cell'],
            'gNodeB ID': [53986],
            'Local Cell ID': [302],
            'Only 5G': ['present only in 5G'],
            'OP/ Vendor': ['Ericsson'],
        }).to_excel(writer, sheet_name='5G', index=False)
    workbook.seek(0)
    response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_vodafone'},
        files={'dataset_files': ('Multivendor_Mapping_VFUK.xlsx', workbook, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')},
        follow_redirects=False,
    )
    assert response.status_code == 303

    four_g_preview = client.get('/workspace/preview/1?source_sheet=4G')
    assert four_g_preview.status_code == 200
    assert '>Cell Name<' in four_g_preview.text
    assert '>4G Cell<' in four_g_preview.text
    assert '>Only 5G<' not in four_g_preview.text

    five_g_preview = client.get('/workspace/preview/1?source_sheet=5G')
    assert five_g_preview.status_code == 200
    assert '>Only 5G<' in five_g_preview.text
    assert '>present only in 5G<' in five_g_preview.text
    assert 'all 6 available columns' in five_g_preview.text
    assert '<span>Available Columns</span><strong>6</strong>' in five_g_preview.text


def test_mapping_preview_filters_by_vendor_and_gcid(client) -> None:
    login(client)
    response = client.post(
        '/dashboard/upload',
        data={'dataset_kinds': 'mapping_three'},
        files={'dataset_files': ('Multivendor_Mapping_3UK.csv', BytesIO(b'Vendor,CId___ECI,Site_Name\nEricsson,123,Site A\nNokia,456,Site B\n'), 'text/csv')},
        follow_redirects=False,
    )
    assert response.status_code == 303

    vendor_preview = client.get('/workspace/preview/1?mapping_vendor=Nokia')
    assert vendor_preview.status_code == 200
    vendor_rows = vendor_preview.text.split('<tbody>', 1)[1].split('</tbody>', 1)[0]
    assert '>Nokia<' in vendor_rows
    assert '>Ericsson<' not in vendor_rows
    assert 'name="mapping_vendor"' in vendor_preview.text

    gcid_preview = client.get('/workspace/preview/1?gcid=123')
    assert gcid_preview.status_code == 200
    gcid_rows = gcid_preview.text.split('<tbody>', 1)[1].split('</tbody>', 1)[0]
    assert '>123<' in gcid_rows
    assert '>456<' not in gcid_rows


def test_workspace_queue_type_filter_lists_all_supported_types_in_order(client) -> None:
    login(client)

    response = client.get("/workspace")
    assert response.status_code == 200
    filter_html = response.text.split('<select data-queue-type-filter>', 1)[1].split('</select>', 1)[0]
    labels = [
        'All Types',
        'CDR-Data',
        'CDR-Speech',
        'CDR-Voice',
        'Multivendor Mapping — Three UK (3UK)',
        'Multivendor Mapping — Vodafone UK (VFUK)',
        'Other supported dataset',
        'Smart Orchestrator Logs',
    ]
    positions = [filter_html.index(label) for label in labels]
    assert positions == sorted(positions)
    assert 'value="data" disabled' in filter_html
    assert '<summary class="collapsible-summary">\n      <div>\n        <p class="eyebrow">Data Processing</p>' in response.text


def test_dataset_selector_shows_all_datasets_when_no_input_kind_filter_is_set(client) -> None:
    login(client)

    client.post(
        "/dashboard/upload",
        files={"dataset_files": ("voice.csv", BytesIO(b"POLQA_LQ_Avg,market,period\n4.2,ES,2026-Q1\n"), "text/csv")},
        follow_redirects=False,
    )
    client.post(
        "/dashboard/upload",
        files={"dataset_files": ("data.csv", BytesIO(b"Mean_Data_Rate,market,period\n25.1,DE,2026-Q2\n"), "text/csv")},
        follow_redirects=False,
    )

    response = client.get("/dashboard?dataset_id=2")
    assert response.status_code == 200
    assert '<option value="1"' in response.text
    assert '<option value="2"' in response.text


def test_dataset_selector_only_lists_ready_datasets(client) -> None:
    login(client)

    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("ready.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("stopped.csv", BytesIO(b"market,period,score\nDE,2026-Q2,78\n"), "text/csv")},
        follow_redirects=False,
    )

    import src.DashboardAnalytic as app_module

    app_module.repository.update_dataset_profile(2, status="stopped", progress=50)

    response = client.get("/dashboard")
    selector_fragment = response.text.split('data-dataset-select', 1)[1].split('</select>', 1)[0]
    assert response.status_code == 200
    assert 'value="1"' in selector_fragment
    assert 'ready.csv' in selector_fragment
    assert 'value="2"' not in selector_fragment
    assert 'stopped.csv' not in selector_fragment


def test_dashboard_selector_excludes_mapping_and_other_dataset_types(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("cdr-data.csv", BytesIO(b"Mean_Data_Rate,Operator\n12.5,EE\n"), "text/csv")},
        follow_redirects=False,
    )
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "mapping_vodafone"},
        files={"dataset_files": ("VFUK.csv", BytesIO(b"eNodeB ID,Local Cell ID,OP/ Vendor\n1,1,Ericsson\n"), "text/csv")},
        follow_redirects=False,
    )

    response = client.get("/dashboard?dataset_id=2")
    selector_fragment = response.text.split('data-dataset-select', 1)[1].split('</select>', 1)[0]
    assert 'cdr-data.csv' in selector_fragment
    assert 'VFUK.csv' not in selector_fragment
    assert 'All CDR Types' in response.text


def test_dashboard_ignores_non_ready_dataset_id_in_selector_flow(client) -> None:
    login(client)

    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("ready.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("failed.csv", BytesIO(b"market,period,score\nDE,2026-Q2,78\n"), "text/csv")},
        follow_redirects=False,
    )

    import src.DashboardAnalytic as app_module

    app_module.repository.update_dataset_profile(2, status="failed", progress=100, last_error="broken")

    response = client.get("/dashboard?dataset_id=2")
    selector_fragment = response.text.split('data-dataset-select', 1)[1].split('</select>', 1)[0]
    assert response.status_code == 200
    assert 'option value="1"' in selector_fragment
    assert 'option value="2"' not in selector_fragment


def test_reporting_preselects_latest_ready_cdr_of_each_type(client) -> None:
    login(client)
    uploads = [
        ('old-data.csv', 'data', b'Mean_Data_Rate,RAT_A\n10,ENDC\n'),
        ('voice.csv', 'voice', b'Call_Setup_Time,RAT_A\n1.2,ENDC\n'),
        ('latest-data.csv', 'data', b'Mean_Data_Rate,RAT_A\n20,ENDC\n'),
        ('speech.csv', 'speech', b'LQ,RAT_A\n3.8,ENDC\n'),
    ]
    for filename, kind, content in uploads:
        response = client.post(
            '/dashboard/upload',
            data={'dataset_kinds': kind},
            files={'dataset_files': (filename, BytesIO(content), 'text/csv')},
            follow_redirects=False,
        )
        assert response.status_code == 303

    reporting = client.get('/reporting')
    data_select = reporting.text.split('name="data_dataset_id"', 1)[1].split('</select>', 1)[0]
    voice_select = reporting.text.split('name="voice_dataset_id"', 1)[1].split('</select>', 1)[0]
    speech_select = reporting.text.split('name="speech_dataset_id"', 1)[1].split('</select>', 1)[0]

    assert 'value="3" data-vendor-mapped="false" selected' in data_select
    assert 'value="2" data-vendor-mapped="false" selected' in voice_select
    assert 'value="4" data-vendor-mapped="false" selected' in speech_select
    assert 'if (active) catalogue.value = active.value;' in reporting.text


def test_dashboard_explicit_dataset_id_overrides_mismatched_input_kind_filter(client) -> None:
    login(client)

    client.post(
        "/dashboard/upload",
        files={"dataset_files": ("voice.csv", BytesIO(b"POLQA_LQ_Avg,market,period\n4.2,ES,2026-Q1\n"), "text/csv")},
        follow_redirects=False,
    )
    client.post(
        "/dashboard/upload",
        files={"dataset_files": ("data.csv", BytesIO(b"Mean_Data_Rate,market,period,test_name,vendor,region\n25.1,DE,2026-Q2,Speed,Nokia,North\n"), "text/csv")},
        follow_redirects=False,
    )

    response = client.get("/dashboard?dataset_id=2&input_kind=voice")
    assert response.status_code == 200
    assert "<h2>data.csv</h2>" in response.text
    assert 'option value="2" data-dataset-kind="data" selected' in response.text


def test_dashboard_data_filters_show_test_name_between_vendor_and_region(client) -> None:
    login(client)

    client.post(
        "/dashboard/upload",
        files={"dataset_files": ("data.csv", BytesIO(b"Mean_Data_Rate,market,period,test_name,vendor,region\n25.1,DE,2026-Q2,Speed,Nokia,North\n"), "text/csv")},
        follow_redirects=False,
    )

    response = client.get("/dashboard?dataset_id=1")
    assert response.status_code == 200
    vendor_pos = response.text.index("Vendor")
    test_name_pos = response.text.index("Test Name")
    region_pos = response.text.index("Region")
    assert vendor_pos < test_name_pos < region_pos


def test_admin_can_update_user_identity_fields(client) -> None:
    login(client)

    create_response = client.post(
        "/admin/users",
        data={"username": "analyst", "password": "start123", "role": "user"},
        follow_redirects=False,
    )
    assert create_response.status_code == 303

    import src.DashboardAnalytic as app_module

    users = app_module.repository.list_users()
    analyst = next(row for row in users if row["username"] == "analyst")

    update_response = client.post(
        f"/admin/users/{analyst['id']}/update",
        data={"username": "analyst-updated", "password": "newpass456", "role": "admin"},
        headers={'X-Requested-With': 'XMLHttpRequest', 'Accept': 'application/json'},
    )
    assert update_response.status_code == 200
    assert update_response.json()['user'] == {
        'id': analyst['id'], 'username': 'analyst-updated', 'role': 'admin', 'active': False, 'workspace_ids': [],
    }

    updated = app_module.repository.get_user("analyst-updated")
    assert updated is not None
    assert next(row for row in app_module.repository.list_users() if row['username'] == 'analyst-updated')['id'] == analyst['id']
    assert updated.role == "admin"
    assert updated.active is False
    assert app_module.verify_password("newpass456", updated.password_hash)


def test_automatic_user_field_update_preserves_the_canonical_username(client) -> None:
    import src.DashboardAnalytic as app_module

    login_super(client)
    created = client.post(
        '/admin/users',
        data={'username': 'analyst', 'password': 'start123', 'role': 'user'},
        follow_redirects=False,
    )
    assert created.status_code == 303
    analyst = next(row for row in app_module.repository.list_users() if row['username'] == 'analyst')

    response = client.post(
        f"/admin/users/{analyst['id']}/update",
        data={'username': 'admin', 'password': '', 'role': 'admin', 'active': '1', 'edited_field': 'role'},
        headers={'X-Requested-With': 'XMLHttpRequest', 'Accept': 'application/json'},
    )
    assert response.status_code == 200
    assert response.json()['user']['username'] == 'analyst'
    assert response.json()['user']['role'] == 'admin'


def test_admin_can_delete_user(client) -> None:
    login(client)

    create_response = client.post(
        "/admin/users",
        data={"username": "temporary", "password": "temp123", "role": "user"},
        follow_redirects=False,
    )
    assert create_response.status_code == 303

    import src.DashboardAnalytic as app_module

    user_row = next(row for row in app_module.repository.list_users() if row["username"] == "temporary")
    delete_response = client.post(f"/admin/users/{user_row['id']}/delete", follow_redirects=False)
    assert delete_response.status_code == 303
    assert app_module.repository.get_user("temporary") is None


def test_password_reset_uses_the_expected_account_defaults(client) -> None:
    import src.DashboardAnalytic as app_module

    login_super(client)
    custom = client.post(
        '/admin/users',
        data={'username': 'analyst', 'password': 'initial-password', 'role': 'user'},
        follow_redirects=False,
    )
    assert custom.status_code == 303

    expected_passwords = {
        'super': 'super123',
        'admin': 'admin123',
        'demo': 'demo123',
        'analyst': 'Ericsson123',
    }
    for username, expected_password in expected_passwords.items():
        record = next(row for row in app_module.repository.list_users() if row['username'] == username)
        app_module.repository.update_password(username, 'different-password')
        response = client.post(f"/admin/users/{record['id']}/reset-password", follow_redirects=False)
        assert response.status_code == 303
        assert app_module.verify_password(expected_password, app_module.repository.get_user(username).password_hash)


def test_admin_cannot_delete_current_signed_in_user(client) -> None:
    login(client)

    import src.DashboardAnalytic as app_module

    admin_row = next(row for row in app_module.repository.list_users() if row["username"] == "admin")
    response = client.post(f"/admin/users/{admin_row['id']}/delete")
    assert response.status_code == 400
    assert "You cannot delete the current signed-in admin user" in response.text


def test_super_admin_cannot_demote_or_deactivate_last_active_super_admin(client) -> None:
    login_super(client)

    import src.DashboardAnalytic as app_module

    admin_row = next(row for row in app_module.repository.list_users() if row["username"] == "super")
    response = client.post(
        f"/admin/users/{admin_row['id']}/update",
        data={"username": "super", "password": "", "role": "user"},
    )
    assert response.status_code == 400
    assert "At least one active super-admin must remain" in response.text

    response = client.post(
        f"/admin/users/{admin_row['id']}/update",
        data={"username": "super", "password": "", "role": "admin"},
    )
    assert response.status_code == 400
    assert "At least one active super-admin must remain" in response.text


def test_admin_cannot_delete_super_admin_even_if_not_current_user(client) -> None:
    login_super(client)

    import src.DashboardAnalytic as app_module

    create_response = client.post(
        "/admin/users",
        data={"username": "backup-admin", "password": "backup123", "role": "admin"},
        follow_redirects=False,
    )
    assert create_response.status_code == 303

    users = app_module.repository.list_users()
    backup_admin = next(row for row in users if row["username"] == "backup-admin")
    admin_row = next(row for row in users if row["username"] == "super")

    switch_session = client.post(
        "/login",
        data={"username": "backup-admin", "password": "backup123"},
        follow_redirects=False,
    )
    assert switch_session.status_code == 303

    disable_backup = client.post(
        f"/admin/users/{backup_admin['id']}/update",
        data={"username": "backup-admin", "password": "", "role": "admin"},
        follow_redirects=False,
    )
    assert disable_backup.status_code == 303

    response = client.post(f"/admin/users/{admin_row['id']}/delete")
    assert response.status_code == 403
    assert "Only super-admins can assign or modify super-admin accounts" in response.text


def test_admin_cannot_assign_or_modify_super_admin_roles(client) -> None:
    login(client)

    create_super = client.post(
        "/admin/users",
        data={"username": "forbidden-super", "password": "password123", "role": "super-admin"},
    )
    assert create_super.status_code == 400
    assert "Only super-admins can create super-admin users" in create_super.text

    created = client.post(
        "/admin/users",
        data={"username": "managed-user", "password": "password123", "role": "user"},
        follow_redirects=False,
    )
    assert created.status_code == 303

    import src.DashboardAnalytic as app_module

    managed_user = next(row for row in app_module.repository.list_users() if row["username"] == "managed-user")
    promoted = client.post(
        f"/admin/users/{managed_user['id']}/update",
        data={"username": "managed-user", "password": "", "role": "super-admin", "active": "1"},
    )
    assert promoted.status_code == 403
    assert "Only super-admins can assign or modify super-admin accounts" in promoted.text

    super_row = next(row for row in app_module.repository.list_users() if row["username"] == "super")
    modified_super = client.post(
        f"/admin/users/{super_row['id']}/update",
        data={"username": "renamed-super", "password": "forbidden-password", "role": "admin", "active": "1"},
    )
    assert modified_super.status_code == 403
    assert "Only super-admins can assign or modify super-admin accounts" in modified_super.text
    assert app_module.repository.get_user('renamed-super') is None
    assert app_module.repository.get_user('super') is not None
    assert app_module.verify_password('super123', app_module.repository.get_user('super').password_hash)

    deleted_super = client.post(f"/admin/users/{super_row['id']}/delete")
    assert deleted_super.status_code == 403
    assert "Only super-admins can assign or modify super-admin accounts" in deleted_super.text


def test_top_navigation_shows_document_links(client) -> None:
    login(client)
    response = client.get("/workspace")
    assert response.status_code == 200
    assert "<h1>Dashboard Analytic</h1>" in response.text
    assert f"v0.2.0 · {__release_date__}" in response.text
    assert 'href="/documents/view/readme"' in response.text
    assert 'href="/documents/view/changelog"' in response.text
    assert 'href="/documents/view/help"' in response.text
    assert 'target="_blank"' not in response.text
    assert 'href="/dashboard"' in response.text
    assert 'class="module-tabs"' in response.text
    assert 'class="module-tabs-secondary"' in response.text
    assert 'module-hero-dashboard' not in response.text
    assert 'class="module-tab module-tab-workspace active" href="/workspace"' in response.text
    assert '<span class="module-tab-label-desktop">E2E PowerPoint Reporting</span>' in response.text
    assert '<span class="module-tab-label-mobile">Reporting</span>' in response.text
    assert '<span class="module-tab-label-mobile">Dashboard</span>' in response.text
    assert 'href="/logout"' in response.text
    assert '<span class="topnav-link topnav-user-badge topnav-user-badge-admin">User: admin</span>' in response.text

    dashboard = client.get("/dashboard")
    assert dashboard.status_code == 200
    assert 'module-hero-dashboard' in dashboard.text
    assert 'linear-gradient(135deg, #0c4c8c, #68b8ff)' in dashboard.text

    reporting = client.get("/reporting")
    assert reporting.status_code == 200
    assert 'module-hero-reporting' in reporting.text
    assert 'linear-gradient(135deg, #4b208a, #bd90ff)' in reporting.text

    admin = client.get("/admin")
    assert admin.status_code == 200
    assert 'module-hero-admin' in admin.text
    assert 'linear-gradient(135deg, #ff8070, #861919)' in admin.text


def test_non_admin_navigation_hides_admin_tab(client) -> None:
    response = client.post("/login", data={"username": "demo", "password": "demo123"}, follow_redirects=False)
    assert response.status_code == 303

    workspace = client.get("/workspace")
    assert workspace.status_code == 200
    assert 'class="module-tabs-secondary"' in workspace.text
    assert 'href="/documents/view/readme"' in workspace.text
    assert 'href="/documents/view/changelog"' in workspace.text
    assert 'href="/documents/view/help"' in workspace.text
    assert 'href="/admin"' not in workspace.text
    assert 'User: demo' in workspace.text


def test_admin_imports_report_catalogue_and_synchronizes_help(client, tmp_path, monkeypatch) -> None:
    from src.modules.cdr_reporting import CATALOG_HEADERS
    import src.DashboardAnalytic as app_module

    login(client)
    help_document = tmp_path / 'powerpoint-reporting.md'
    help_document.write_text(
        '# Reporting\n\n<!-- SLIDES_TEMPLATES:START -->\nold\n<!-- SLIDES_TEMPLATES:END -->\n',
        encoding='utf-8',
    )
    monkeypatch.setattr(app_module, 'REPORT_CATALOGUE_DOCUMENT', help_document)
    content = (
        ','.join(CATALOG_HEADERS)
        + '\n8,Completed Call Ratio,Voice quality,Title and 1 column + Comments,Completed call ratio,CDR-Voice,Call_Status,100% Stacked Vertical Bars,Call Family = VoLTE,Operator,Campaign,Completed/Dropped/Failed,\n'
    ).encode('utf-8')

    response = client.post(
        '/admin/report-catalogues/nsa',
        data={'catalogue_name': 'Test baseline'},
        files={'catalogue_file': ('nsa-slides-template.csv', BytesIO(content), 'text/csv')},
        follow_redirects=False,
    )

    assert response.status_code == 303
    assert app_module.reporting_catalog_path('nsa').read_bytes() == content
    assert 'Completed Call Ratio' in help_document.read_text(encoding='utf-8')

    exported = client.get('/admin/report-catalogues/nsa/export')
    assert exported.status_code == 200
    assert exported.content == content.rstrip(b'\n').rstrip(b',') + b',Top\n'

    confirmation = client.get('/admin?catalogue_notice=Imported%20Test%20baseline%20%28NSA%29.')
    assert 'data-catalogue-import-notice' in confirmation.text
    assert 'catalogue-management-notice' not in confirmation.text
    assert 'id="info-overlay"' in confirmation.text


def test_admin_import_converts_a_legacy_catalogue_when_requested(client, tmp_path, monkeypatch) -> None:
    import src.DashboardAnalytic as app_module

    login(client)
    help_document = tmp_path / 'powerpoint-reporting.md'
    help_document.write_text('# Reporting\n\n<!-- SLIDES_TEMPLATES:START -->\nold\n<!-- SLIDES_TEMPLATES:END -->\n', encoding='utf-8')
    monkeypatch.setattr(app_module, 'REPORT_CATALOGUE_DOCUMENT', help_document)
    legacy = (
        'Slide,Slide title,Slide subtitle,Layout,CDR Source,KPI,Chart Type,Filters,Grouping\n'
        '8,Legacy quality,,Title and 1 column + Comments,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator × Campaign\n'
    ).encode('utf-8')

    response = client.post(
        '/admin/report-catalogues/nsa',
        data={'catalogue_name': 'Legacy baseline', 'convert_catalogue': '1'},
        files={'catalogue_file': ('legacy.csv', BytesIO(legacy), 'text/csv')},
        follow_redirects=False,
    )

    assert response.status_code == 303
    stored = app_module.reporting_catalog_path('nsa').read_text(encoding='utf-8')
    assert 'Chart Tittle' in stored
    assert 'Legacy quality' in stored


def test_admin_stores_multiple_named_report_catalogues_and_can_activate_one(client, tmp_path, monkeypatch) -> None:
    from src.modules.cdr_reporting import CATALOG_HEADERS
    import src.DashboardAnalytic as app_module

    login(client)
    help_document = tmp_path / 'powerpoint-reporting.md'
    help_document.write_text('# Reporting\n\n<!-- SLIDES_TEMPLATES:START -->\nold\n<!-- SLIDES_TEMPLATES:END -->\n', encoding='utf-8')
    monkeypatch.setattr(app_module, 'REPORT_CATALOGUE_DOCUMENT', help_document)
    first = (','.join(CATALOG_HEADERS) + '\n8,First,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,\n').encode('utf-8')
    second = (','.join(CATALOG_HEADERS) + '\n8,Second,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,\n').encode('utf-8')

    for name, content in [('Baseline Q4', first), ('Updated Q4', second)]:
        response = client.post(
            '/admin/report-catalogues/nsa',
            data={'catalogue_name': name},
            files={'catalogue_file': ('nsa.csv', BytesIO(content), 'text/csv')},
            follow_redirects=False,
        )
        assert response.status_code == 303

    admin = client.get('/admin')
    assert 'Baseline Q4' in admin.text
    assert 'Updated Q4' in admin.text
    assert 'name="catalogue_selection"' in admin.text
    assert 'Choose a template to edit' in admin.text
    assert 'data-catalogue-editor-table' not in admin.text
    assert '<th>Default</th>' in admin.text
    assert 'catalogue-default-mark is-default' in admin.text
    assert 'value="nsa:Baseline Q4"' in admin.text
    assert '/admin/report-catalogues/export-selected' in admin.text
    assert app_module.reporting_catalog_entries('nsa')[0].slide_title == 'Second'
    assert app_module.reporting_catalog_path('nsa').name == 'Updated Q4.csv'
    assert app_module.named_catalogue_path('nsa', 'Baseline Q4').exists()

    reporting = client.get('/reporting')
    assert 'value="nsa:Updated Q4" data-catalogue-technology="nsa" data-catalogue-active="true" selected' in reporting.text

    editor = client.get('/admin?catalogue_technology=nsa&catalogue_id=Baseline%20Q4')
    assert editor.status_code == 200
    assert 'Slides Templates Editor' in editor.text
    assert 'data-catalogue-editor-table' in editor.text
    assert 'data-catalogue-field="Layout"' in editor.text
    assert 'data-catalogue-editor-options' in editor.text
    assert 'data-catalogue-row-action="insert"' in editor.text
    assert 'data-catalogue-row-action="up"' in editor.text
    assert 'data-catalogue-row-action="down"' in editor.text
    assert 'data-catalogue-row-action="delete"' in editor.text
    assert 'data-catalogue-reenumerate' in editor.text
    assert editor.text.index('catalogue-row-actions-heading') < editor.text.index('>Slide</th>')
    assert 'value="nsa:Baseline Q4" selected' in editor.text
    assert 'Title and 1 column + Comments' in editor.text
    assert '<optgroup label="Layouts">' in editor.text
    assert '<optgroup label="Chart types">' in editor.text
    assert '<optgroup label="CDR fields">' in editor.text

    edited = (
        ','.join(CATALOG_HEADERS)
        + '\n9,Late,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,Right\n'
        + '\n8,Edited,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,Right\n'
    )
    saved = client.post('/admin/report-catalogues/nsa/Baseline%20Q4/save', data={'catalogue_content': edited}, follow_redirects=False)
    assert saved.status_code == 303
    saved_editor = client.get('/admin?catalogue_technology=nsa&catalogue_id=Baseline%20Q4')
    assert '>Edited</td>' in saved_editor.text
    assert saved_editor.text.index('>Edited</td>') < saved_editor.text.index('>Late</td>')

    activated = client.post('/admin/report-catalogues/nsa/Baseline%20Q4/activate', follow_redirects=False)
    assert activated.status_code == 303
    assert app_module.reporting_catalog_entries('nsa')[0].slide_title == 'Edited'
    assert app_module.reporting_catalog_path('nsa').name == 'Baseline Q4.csv'
    assert app_module.named_catalogue_path('nsa', 'Updated Q4').exists()

    protected_delete = client.post('/admin/report-catalogues/nsa/Baseline%20Q4/delete')
    assert protected_delete.status_code == 400
    assert 'The default template cannot be deleted.' in protected_delete.text

    reporting_after_activation = client.get('/reporting')
    assert 'value="nsa:Baseline Q4" data-catalogue-technology="nsa" data-catalogue-active="true" selected' in reporting_after_activation.text

    exported = client.get('/admin/report-catalogues/nsa/Updated%20Q4/export')
    assert exported.status_code == 200
    assert b'Second' in exported.content

def test_admin_catalogue_rename_supports_background_json_save(client) -> None:
    import src.DashboardAnalytic as app_module

    login(client)
    default_name = next(item['identifier'] for item in app_module.report_catalogue_options('nsa') if item['active'])
    response = client.post(
        f'/admin/report-catalogues/nsa/{quote(default_name)}/rename',
        data={'catalogue_name': 'Renamed default'},
        headers={'accept': 'application/json'},
    )

    assert response.status_code == 200
    assert response.json() == {'name': 'Renamed default', 'identifier': 'Renamed default'}
    assert next(item for item in app_module.report_catalogue_options('nsa') if item['identifier'] == 'Renamed default')['name'] == 'Renamed default'


def test_admin_renaming_named_catalogue_renames_its_csv_file(client, tmp_path, monkeypatch) -> None:
    from src.modules.cdr_reporting import CATALOG_HEADERS
    import src.DashboardAnalytic as app_module

    login(client)
    help_document = tmp_path / 'powerpoint-reporting.md'
    help_document.write_text('# Reporting\n\n<!-- SLIDES_TEMPLATES:START -->\nold\n<!-- SLIDES_TEMPLATES:END -->\n', encoding='utf-8')
    monkeypatch.setattr(app_module, 'REPORT_CATALOGUE_DOCUMENT', help_document)
    content = (
        ','.join(CATALOG_HEADERS)
        + '\n8,First,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,\n'
    ).encode('utf-8')
    imported = client.post(
        '/admin/report-catalogues/nsa',
        data={'catalogue_name': 'Original catalogue'},
        files={'catalogue_file': ('nsa.csv', BytesIO(content), 'text/csv')},
        follow_redirects=False,
    )
    assert imported.status_code == 303

    replacement = client.post(
        '/admin/report-catalogues/nsa',
        data={'catalogue_name': 'Replacement template'},
        files={'catalogue_file': ('replacement.csv', BytesIO(content), 'text/csv')},
        follow_redirects=False,
    )
    assert replacement.status_code == 303

    original_path = app_module.named_catalogue_path('nsa', 'Original catalogue')
    response = client.post(
        '/admin/report-catalogues/nsa/Original%20catalogue/rename',
        data={'catalogue_name': 'Renamed catalogue'},
        headers={'accept': 'application/json'},
    )

    renamed_path = app_module.named_catalogue_path('nsa', 'Renamed catalogue')
    assert response.status_code == 200
    assert response.json() == {'name': 'Renamed catalogue', 'identifier': 'Renamed catalogue'}
    assert not original_path.exists()
    assert renamed_path.read_bytes() == content
    assert not next(item for item in app_module.report_catalogue_options('nsa') if item['identifier'] == 'Renamed catalogue')['active']


def test_admin_duplicates_template_using_the_source_template_name(client) -> None:
    from src.modules.cdr_reporting import CATALOG_HEADERS
    import src.DashboardAnalytic as app_module

    login(client)
    content = (
        ','.join(CATALOG_HEADERS)
        + '\n8,First,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,\n'
    ).encode('utf-8')
    imported = client.post(
        '/admin/report-catalogues/nsa',
        data={'catalogue_name': 'Regional NSA Template'},
        files={'catalogue_file': ('nsa.csv', BytesIO(content), 'text/csv')},
        follow_redirects=False,
    )
    assert imported.status_code == 303

    duplicated = client.post('/admin/report-catalogues/nsa/Regional%20NSA%20Template/duplicate', follow_redirects=False)

    assert duplicated.status_code == 303
    copied = next(item for item in app_module.report_catalogue_options('nsa') if item['identifier'] == 'Regional NSA Template - Copy')
    assert copied['name'] == 'Regional NSA Template - Copy'
    assert copied['path'].name == 'Regional NSA Template - Copy.csv'


def test_template_registry_reconciles_an_unambiguous_manual_csv_rename(client) -> None:
    from src.modules.cdr_reporting import CATALOG_HEADERS
    import src.DashboardAnalytic as app_module

    # Keep this reconciliation scenario independent from whichever templates
    # are currently bundled with the project.
    default_dir = app_module.settings.slides_templates_dir / 'default' / 'nsa'
    current_default = next(default_dir.glob('*.csv'))
    base_default = default_dir / 'Base template.csv'
    current_default.rename(base_default)
    library_dir = app_module.settings.slides_templates_dir / 'library' / 'nsa'
    library_dir.mkdir(parents=True, exist_ok=True)
    for path in library_dir.glob('*.csv'):
        path.unlink()
    app_module.repository.add_report_template('nsa', 'Base template')
    app_module.repository.set_default_report_template('nsa', 'Base template')
    login(client)
    content = (
        ','.join(CATALOG_HEADERS)
        + '\n8,First,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,\n'
    ).encode('utf-8')
    for name in ('First template', 'Second template'):
        response = client.post(
            '/admin/report-catalogues/nsa',
            data={'catalogue_name': name},
            files={'catalogue_file': ('nsa.csv', BytesIO(content), 'text/csv')},
            follow_redirects=False,
        )
        assert response.status_code == 303

    original = app_module.named_catalogue_path('nsa', 'First template', 'First template')
    renamed = original.with_name('Historic baseline.csv')
    original.rename(renamed)

    options = app_module.report_catalogue_options('nsa')
    reconciled = next(item for item in options if item['identifier'] == 'Historic baseline')
    assert reconciled['name'] == 'Historic baseline'
    assert reconciled['path'] == renamed


def test_admin_importer_selects_template_type_and_moves_a_named_template(client) -> None:
    from src.modules.cdr_reporting import CATALOG_HEADERS
    import src.DashboardAnalytic as app_module

    login(client)
    content = (
        ','.join(CATALOG_HEADERS)
        + '\n8,First,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,\n'
    ).encode('utf-8')
    imported_sa = client.post(
        '/admin/slides-templates/import',
        data={'template_type': 'sa', 'catalogue_name': 'SA imported template'},
        files={'catalogue_file': ('sa.csv', BytesIO(content), 'text/csv')},
        follow_redirects=False,
    )
    assert imported_sa.status_code == 303
    assert app_module.reporting_catalog_path('sa').name == 'SA imported template.csv'

    # Add two NSA templates: the second becomes default, so the first remains
    # a movable library item.
    for name in ('Move me', 'NSA default replacement'):
        response = client.post(
            '/admin/report-catalogues/nsa',
            data={'catalogue_name': name},
            files={'catalogue_file': ('nsa.csv', BytesIO(content), 'text/csv')},
            follow_redirects=False,
        )
        assert response.status_code == 303
    moved = client.post(
        '/admin/report-catalogues/nsa/Move%20me/type',
        data={'template_type': 'sa'},
        follow_redirects=False,
    )
    assert moved.status_code == 303
    assert app_module.named_catalogue_path('sa', 'Move me').exists()
    assert not app_module.named_catalogue_path('nsa', 'Move me').exists()


def test_docs_routes_expose_readme_changelog_and_help(client) -> None:
    login(client)

    readme_view = client.get("/documents/view/readme")
    assert readme_view.status_code == 200
    assert "Loading document..." in readme_view.text
    assert "/api/documents/readme" in readme_view.text

    changelog_api = client.get("/api/documents/changelog")
    assert changelog_api.status_code == 200
    payload = changelog_api.json()
    assert payload["name"] == "CHANGELOG.md"
    assert "0.1.0" in payload["content"]

    help_view = client.get("/documents/view/help")
    assert help_view.status_code == 200
    assert "/api/documents/help" in help_view.text

    help_api = client.get("/api/documents/help")
    assert help_api.status_code == 200
    assert help_api.json()["name"] == "00-help.md"

    help_index = client.get("/api/documents/help-index")
    assert help_index.status_code == 200
    help_documents = help_index.json()["documents"]
    assert help_documents[0]["relative_path"] == "00-help.md"
    assert help_documents[0]["number"] == "00"
    assert help_documents[1]["label"] == "Configuration File"
    assert any(item["relative_path"] == "02-web-interface.md" for item in help_documents)
    assert any(
        item["relative_path"] == "04-e2e-dashboard-analysis.md"
        and item["label"] == "E2E Dashboard"
        for item in help_documents
    )
    assert any(
        item["relative_path"] == "05-e2e-ppt-reporting.md"
        and item["label"] == "E2E PowerPoint Reporting"
        for item in help_documents
    )
    excluded_help_documents = {
        "02-arguments-description.md",
        "02-arguments-description-short.md",
        "05-word-reporting.md",
        "09-github-actions.md",
        "10-testing.md",
    }
    assert not excluded_help_documents.intersection(item["relative_path"] for item in help_documents)

    help_article = client.get("/documents/view/help/02-web-interface.md")
    assert help_article.status_code == 200
    assert 'id="help-nav-list"' in help_article.text
    assert "/api/documents/help/02-web-interface.md" in help_article.text
    assert "`${number}. ${label}`" in help_article.text


def test_dashboard_analysis_reuses_cached_result_on_reload(client, monkeypatch) -> None:
    login(client)
    csv_content = b"market,period,score,gap\nES,2026-Q1,91,2.1\nES,2026-Q1,87,3.3\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    import src.DashboardAnalytic as app_module

    calls = {"count": 0}
    original_load_dataset = app_module.load_dataset
    app_module.ANALYSIS_CACHE.clear()
    app_module.DATAFRAME_CACHE.clear()
    assert app_module.repository.dataset_rows_table_exists(1)

    def counting_load_dataset(path):
        calls["count"] += 1
        return original_load_dataset(path)

    monkeypatch.setattr(app_module, "load_dataset", counting_load_dataset)

    first_response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert first_response.status_code == 200
    assert calls["count"] == 0

    second_response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert second_response.status_code == 200
    assert calls["count"] == 0


def test_dashboard_analysis_reuses_cached_dataset_frame_across_metric_changes(client, monkeypatch) -> None:
    login(client)
    csv_content = b"market,period,score,gap\nES,2026-Q1,91,2.1\nES,2026-Q1,87,3.3\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    import src.DashboardAnalytic as app_module

    calls = {"count": 0}
    original_load_dataset = app_module.load_dataset
    app_module.ANALYSIS_CACHE.clear()
    app_module.DATAFRAME_CACHE.clear()
    assert app_module.repository.dataset_rows_table_exists(1)

    def counting_load_dataset(path):
        calls["count"] += 1
        return original_load_dataset(path)

    monkeypatch.setattr(app_module, "load_dataset", counting_load_dataset)

    first_response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert first_response.status_code == 200
    assert calls["count"] == 0

    second_response = client.get("/dashboard?dataset_id=1&metric=gap&aggregation=all&load=1")
    assert second_response.status_code == 200
    assert calls["count"] == 0


def test_dashboard_renders_multiple_selected_metrics(client) -> None:
    login(client)
    csv_content = b"market,period,score,gap\nES,2026-Q1,91,2.1\nES,2026-Q1,87,3.3\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&metric=gap&aggregation=all&load=1")
    assert response.status_code == 200
    assert "Use the dropdown to select one, several, or all KPIs." in response.text
    assert response.text.count("Metric View") >= 2
    assert response.text.count("Selected Metric") >= 2
    assert "mean metric" in response.text.lower()
    assert "score" in response.text
    assert "gap" in response.text


def test_dashboard_shows_date_range_filters_and_applies_them(client) -> None:
    login(client)
    csv_content = (
        b"market,period,score,Call Start Time\n"
        b"ES,2026-Q1,91,2026-07-10 10:00:00\n"
        b"ES,2026-Q1,87,2026-07-11 12:00:00\n"
    )
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&date_from=2026-07-11&load=1")
    assert response.status_code == 200
    assert 'name="date_from"' in response.text
    assert 'name="date_to"' in response.text
    assert 'value="2026-07-11"' in response.text
    assert "2026-07-10" not in response.text


def test_dashboard_adaptive_filters_include_city_and_multi_select_fields(client) -> None:
    login(client)
    csv_content = (
        b"market,period,score,City,Region\n"
        b"ES,2026-Q1,91,Madrid,Central\n"
        b"ES,2026-Q1,87,Barcelona,East\n"
    )
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert response.status_code == 200
    assert 'select name="city" multiple' in response.text
    assert 'select name="region" multiple' in response.text
    assert ">Madrid<" in response.text
    assert ">Barcelona<" in response.text
    assert "All values are selected by default. Clearing all values applies an empty filter." in response.text


def test_dashboard_adaptive_filters_populate_netcheck_a_columns_for_existing_cdrs(client) -> None:
    login(client)
    csv_content = (
        b"RAT_A,Operator_A,Session_Type_A,Call_Status_A,Call_Duration\n"
        b"ENDC,Vodafone UK,VoLTE,Completed,61\n"
        b"NR,Three UK,VoNR,Dropped,42\n"
    )
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "voice"},
        files={"dataset_files": ("netcheck_voice.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    import src.DashboardAnalytic as app_module

    table_name = app_module.repository.dataset_rows_table_name(1)
    with app_module.repository.connection() as conn:
        conn.execute(
            f'''UPDATE "{table_name}" SET "operator" = NULL, "session_type" = NULL, "status" = NULL'''
        )
        conn.execute(
            """
            UPDATE dataset_profiles
            SET normalization_version = 4, filter_options_json = '{}'
            WHERE dataset_id = 1
            """
        )

    response = client.get("/dashboard?dataset_id=1&metric=Call_Duration&aggregation=all&load=1")
    assert response.status_code == 200
    assert 'select name="operator" multiple' in response.text
    assert 'value="Vodafone UK"' in response.text
    assert 'value="Three UK"' in response.text
    assert 'select name="session_type" multiple' in response.text
    assert 'value="VoLTE"' in response.text
    assert 'value="VoNR"' in response.text


def test_dashboard_adaptive_filters_label_technology_without_primary(client) -> None:
    login(client)
    csv_content = b"market,period,score,RAT\nES,2026-Q1,91,5G\nES,2026-Q1,87,LTE\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert response.status_code == 200
    assert ">Technology<" in response.text
    assert "Technology Primary" not in response.text


def test_dashboard_comparison_chart_exposes_per_metric_aggregation_override_control(client) -> None:
    login(client)
    csv_content = b"market,period,score,gap,operator,region\nES,2026-Q1,91,2.1,Vodafone,North\nES,2026-Q1,87,3.3,o2,South\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&metric=gap&aggregation=region&aggregation_overrides=score=operator&load=1")
    assert response.status_code == 200
    assert 'data-chart-aggregation-select' in response.text
    assert 'data-metric="score"' in response.text
    assert 'data-current-overrides="score=operator"' in response.text


def test_dashboard_exposes_global_and_per_metric_cdf_comparison_controls(client) -> None:
    login(client)
    csv_content = (
        b"market,period,score,vendor,region,operator,city\n"
        b"ES,2026-Q1,91,Nokia,North,Vodafone,Madrid\n"
        b"ES,2026-Q1,87,Huawei,South,Vodafone,Barcelona\n"
    )
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&load=1&cdf_grouping=vendor")
    assert response.status_code == 200
    assert "Global CDF Comparison" in response.text
    assert 'data-global-cdf-grouping-select' in response.text
    assert 'data-chart-cdf-grouping-select' in response.text
    assert 'data-cdf-range-control' in response.text
    assert 'Compare CDF by' in response.text
    assert 'input type="hidden" name="cdf_grouping" value="vendor"' in response.text


def test_dashboard_powerpoint_export_includes_visual_analytics_payload(client) -> None:
    login(client)
    csv_content = (
        b"market,period,score,gap,vendor,operator,region,city,Call Start Time\n"
        b"ES,2026-Q1,91,2.1,Nokia,Vodafone,North,Madrid,2026-07-10 10:00:00\n"
        b"ES,2026-Q1,87,3.3,Huawei,Orange,South,Barcelona,2026-07-11 11:00:00\n"
    )
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.post(
        "/dashboard/export/powerpoint",
        data={
            "dataset_id": "1",
            "metric": ["score", "gap"],
            "market": ["ES"],
            "aggregation": "operator",
            "cdf_grouping": "vendor",
            "date_from": "2026-07-10",
            "date_to": "2026-07-11",
            "extra_filters": "vendor=Nokia,Huawei; region=North,South",
        },
    )
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    from pptx import Presentation

    presentation = Presentation(BytesIO(response.content))
    assert len(presentation.slides) >= 4
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "sample.csv" in slide_text
    assert "Visual Analytics · score" in slide_text
    assert "Visual Analytics · gap" in slide_text
    assert "Date From: 2026-07-10" in slide_text
    assert "Mean" in slide_text
    assert "Avg" in slide_text
    assert "P10" in slide_text
    assert "P90" in slide_text
    assert "Min" in slide_text
    assert "Max" in slide_text


def test_workspace_logs_capture_analysis_warnings(client, monkeypatch) -> None:
    login(client)
    csv_content = b"market,period,score,gap\nES,2026-Q1,91,2.1\nES,2026-Q1,87,3.3\n"
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )

    import src.DashboardAnalytic as app_module

    original_build_analysis = app_module.build_analysis

    def warned_build_analysis(*args, **kwargs):
        warnings.warn("Synthetic analysis warning for workspace logs", UserWarning)
        return original_build_analysis(*args, **kwargs)

    monkeypatch.setattr(app_module, "build_analysis", warned_build_analysis)

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert response.status_code == 200

    logs = app_module.repository.list_workspace_logs(1)
    warning_logs = [log for log in logs if log["action"] == "analyze_dataset_warning"]
    assert warning_logs
    assert "Synthetic analysis warning for workspace logs" in warning_logs[0]["details_text"]


def test_dashboard_handles_empty_table_rows_without_template_failure(client) -> None:
    login(client)
    csv_content = b"market,period,operator,score\nES,2026-Q1,VDF,91\nES,2026-Q1,VDF,87\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=operator&market=DE&load=1")
    assert response.status_code == 200
    assert "No rows match the selected filters" in response.text or "No tabular rows match the selected filters" in response.text


def test_dashboard_materializes_legacy_ready_dataset_on_first_analysis(client) -> None:
    login(client)
    csv_content = b"market,period,operator,score\nES,2026-Q1,VDF,91\nES,2026-Q1,VDF,87\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    import src.DashboardAnalytic as app_module

    app_module.repository.drop_dataset_rows(1)
    assert not app_module.repository.dataset_rows_table_exists(1)

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=operator&load=1")
    assert response.status_code == 200
    assert "Charts and Scorecards" in response.text
    assert app_module.repository.dataset_rows_table_exists(1)


def test_dashboard_reuses_materialized_table_when_legacy_columns_only_differ_by_case(client, monkeypatch) -> None:
    login(client)
    csv_content = b"market,period,operator,score\nES,2026-Q1,VDF,91\nES,2026-Q1,ORG,87\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    import src.DashboardAnalytic as app_module

    legacy_frame = pd.DataFrame({
        "Market": ["ES", "ES"],
        "Period": ["2026-Q1", "2026-Q1"],
        "Operator": ["VDF", "ORG"],
        "score": [91, 87],
        "dataset_kind": ["generic", "generic"],
        "source_file": ["sample.csv", "sample.csv"],
    })
    app_module.repository.replace_dataset_rows(1, legacy_frame)
    app_module.DATAFRAME_CACHE.clear()
    app_module.ANALYSIS_CACHE.clear()

    calls = {"count": 0}
    original_load_dataset = app_module.load_dataset

    def counting_load_dataset(path):
        calls["count"] += 1
        return original_load_dataset(path)

    monkeypatch.setattr(app_module, "load_dataset", counting_load_dataset)

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=operator&load=1")
    assert response.status_code == 200
    assert "Charts and Scorecards" in response.text
    assert calls["count"] == 0


def test_dashboard_refreshes_stale_dataset_normalization_before_render(client) -> None:
    login(client)
    csv_content = b"market,period,score,RAT,PCell_RAT_Timeline\nES,2026-Q1,91,5G NSA,NR->LTE\nES,2026-Q1,87,LTE,LTE->NR\n"
    upload_response = client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    import src.DashboardAnalytic as app_module

    with app_module.repository.connection() as conn:
        conn.execute(
            """
            UPDATE dataset_profiles
            SET normalization_version = 1,
                filter_options_json = ?,
                available_aggregations_json = ?
            WHERE dataset_id = 1
            """,
            (
                json.dumps({"technology_primary": ["NR->LTE", "LTE->NR"]}),
                json.dumps(["technology_primary"]),
            ),
        )

    response = client.get("/dashboard?dataset_id=1&metric=score&aggregation=all&load=1")
    assert response.status_code == 200
    assert ">5G NSA<" in response.text
    assert ">LTE<" in response.text
    assert "NR-&gt;LTE" not in response.text

    refreshed = app_module.repository.get_dataset(1)
    assert refreshed is not None
    assert int(refreshed["normalization_version"]) == app_module.DATASET_NORMALIZATION_VERSION


def test_dataset_status_endpoint_returns_queue_payload(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )

    response = client.get("/api/datasets/status")
    assert response.status_code == 200
    payload = response.json()
    assert "datasets" in payload
    assert payload["datasets"][0]["file_name"] == "sample.csv"
    assert payload["datasets"][0]["size_mb_label"].endswith("MB")


def test_dashboard_handles_missing_source_file_without_500(client) -> None:
    login(client)

    import src.DashboardAnalytic as app_module

    with app_module.repository.connection() as conn:
        conn.execute(
            "INSERT INTO datasets (id, file_name, stored_path, uploaded_by) VALUES (?, ?, ?, ?)",
            (99, "missing.xlsx", "/tmp/does-not-exist.xlsx", "admin"),
        )
        conn.execute(
            """
            INSERT INTO dataset_profiles (
                dataset_id, status, progress, dataset_kind, default_metric, default_aggregation,
                available_metrics_json, available_aggregations_json, filter_options_json, summary_json, kpis_json
            ) VALUES (?, 'ready', 100, 'data', 'throughput_mbps', 'operator', '["throughput_mbps"]', '["operator"]', '{}', '{}', '{}')
            """,
            (99,),
        )

    response = client.get("/dashboard?dataset_id=99&metric=throughput_mbps&aggregation=operator&load=1")
    assert response.status_code == 200
    assert "source file is missing" in response.text


def test_materialized_dataset_handles_case_insensitive_duplicate_columns(client) -> None:
    login(client)
    csv_content = b"Campaign,campaign,score\nES_Q1_2026,manual-campaign,91\n"
    upload_response = client.post(
        "/dashboard/upload",
        files={"dataset_files": ("duplicate-columns.csv", BytesIO(csv_content), "text/csv")},
        follow_redirects=False,
    )
    assert upload_response.status_code == 303

    import src.DashboardAnalytic as app_module

    dataset = app_module.repository.get_dataset(1)
    assert dataset is not None
    assert dataset["status"] == "ready"
    assert dataset["last_error"] in (None, "")


def test_failed_dataset_shows_last_error_in_queue(client) -> None:
    login(client)

    import src.DashboardAnalytic as app_module

    with app_module.repository.connection() as conn:
        conn.execute(
            "INSERT INTO datasets (id, file_name, stored_path, uploaded_by) VALUES (?, ?, ?, ?)",
            (50, "broken.csv", "/tmp/broken.csv", "admin"),
        )
        conn.execute(
            """
            INSERT INTO dataset_profiles (
                dataset_id, status, progress, dataset_kind, last_error, available_metrics_json,
                available_aggregations_json, filter_options_json, summary_json, kpis_json
            ) VALUES (?, 'failed', 100, 'generic', 'duplicate column name: campaign', '[]', '[]', '{}', '{}', '{}')
            """,
            (50,),
        )

    response = client.get("/workspace")
    assert response.status_code == 200
    assert "duplicate column name: campaign" in response.text


def test_workspace_queue_shows_dataset_size_column(client) -> None:
    login(client)
    client.post(
        "/dashboard/upload",
        data={"dataset_kinds": "data"},
        files={"dataset_files": ("sample.csv", BytesIO(b"market,period,score\nES,2026-Q1,91\n"), "text/csv")},
        follow_redirects=False,
    )

    response = client.get("/workspace")
    assert response.status_code == 200
    assert "<th>Size</th>" in response.text
    assert "MB" in response.text


def test_workspace_shows_operational_logs_panel(client) -> None:
    login(client)

    import src.DashboardAnalytic as app_module

    app_module.repository.add_log(
        "admin",
        "process_dataset_failed",
        '{"dataset_id": 1, "file": "sample.csv", "error": "Synthetic processing failure"}',
    )

    response = client.get("/workspace")
    assert response.status_code == 200
    assert "Workspace Logs" in response.text
    assert "Execution and Error Trail" in response.text
    assert "All events" in response.text
    assert "Info only" in response.text
    assert "Error only" in response.text
    assert "Type" in response.text
    assert "Error" in response.text
    assert "Synthetic processing failure" in response.text
