from __future__ import annotations

import json
import re
import time
from datetime import datetime
import pandas as pd
import pytest
from io import BytesIO
from pathlib import Path
from unittest.mock import patch
from urllib.parse import urlencode
from pptx import Presentation
from pptx.dml.color import RGBColor

from src.modules.cdr_reporting import CATALOG_HEADERS, CatalogEntry, _apply_catalog_filters, _apply_catalog_grouping, _cdf_terminal_x_maximum, _hierarchical_unique_keys, _hierarchy_group_colours, _layout_chart_frames, _named_slide_layout, _render_cdf_line, _render_failure_count, _render_failure_count_hierarchy, _render_status_100, _series_colours, assign_cdr_vendors, classify_sessions, convert_catalog_csv, ensure_report_vendor_group, enrich_multivendor, load_catalog_csv, normalise_report_operator_aliases, parse_catalog_csv, parse_catalog_filters, parse_catalog_grouping, parse_legend_position, prepare_multivendor_catalog_entry, render_cdr_report, vendor_from_cells
from src.modules.repository import Repository


def wait_for_report_job(client, job_id: int) -> dict:
    for _ in range(100):
        response = client.get('/api/reporting/jobs')
        assert response.status_code == 200
        job = next(item for item in response.json()['jobs'] if item['id'] == job_id)
        if job['status'] not in {'queued', 'processing'}:
            return job
        time.sleep(0.02)
    raise AssertionError(f'Report job {job_id} did not finish in time.')


def wait_for_report_chart_job(client, job_id: int) -> dict:
    for _ in range(100):
        response = client.get('/api/reporting/chart-jobs')
        assert response.status_code == 200
        job = next(item for item in response.json()['jobs'] if item['id'] == job_id)
        if job['status'] not in {'queued', 'processing'}:
            return job
        time.sleep(0.02)
    raise AssertionError(f'Report Charts job {job_id} did not finish in time.')


def test_vendor_formula_keeps_vodafone_ericsson_null_exception_as_mixed() -> None:
    lookup = {'first': 'Ericsson'}

    assert vendor_from_cells('Vodafone UK', 'first -> unknown', lookup) == 'Vodafone_Mixed Vendor'
    assert vendor_from_cells('Vodafone UK', 'first -> first', lookup) == 'Vodafone_Ericsson'
    assert vendor_from_cells('3', 'first -> unknown', lookup) == '3_Mixed Vendor'
    assert vendor_from_cells('O2', 'first -> unknown', lookup) == 'O2'


def test_report_operator_aliases_share_filters_and_grouping_across_campaigns() -> None:
    frame = pd.DataFrame({
        "Operator": ["Vodafone", "Vodafone UK", "o2 - de", "O2(UK)", "Telefónica", "Three", "Three UK", "3 UK", "EE UK", "Everything Everywhere"],
        "Campaign": ["UK_Q4_2025", "UK_Q2_2026", "UK_Q4_2025", "UK_Q2_2026", "UK_Q2_2026", "UK_Q4_2025", "UK_Q2_2026", "UK_Q4_2025", "UK_Q2_2026", "UK_Q4_2025"],
        "Call_Status": ["Completed"] * 10,
    })
    entry = CatalogEntry(
        1, "", "", "", "", "CDR-Voice", "Call_Status", "100% Stacked Vertical Bars", "",
        "Operator IN (Vodafone, O2, 3, EE)", "Operator", "Campaign",
    )

    normalised = normalise_report_operator_aliases(frame)
    filtered = _apply_catalog_filters(normalised, entry, False, "Call_Status")
    grouped, primary, series = _apply_catalog_grouping(filtered, entry, False, "Call_Status")

    assert filtered["Operator"].tolist() == ["VF", "VF", "O2", "O2", "O2", "3", "3", "3", "EE", "EE"]
    assert set(grouped[primary]) == {"VF", "O2", "3", "EE"}
    assert set(grouped[series]) == {"2025-Q4", "2026-Q2"}


def test_h3g_is_not_normalised_as_operator_three() -> None:
    frame = pd.DataFrame({"Operator": ["H3G", "H3G UK", "Three UK"]})

    assert normalise_report_operator_aliases(frame)["Operator"].tolist() == ["H3G", "H3G UK", "3"]


def test_reporting_cache_resolves_separator_variants_and_refreshes_derived_dimensions(tmp_path) -> None:
    repository = Repository(tmp_path / 'workspace.db')
    repository.replace_dataset_rows(1, pd.DataFrame({
        'RAT_A': ['EN-DC'], 'G_Level_4': ['London'], 'Call Family': [None],
    }))
    repository.copy_dataset_rows_to_reporting(1, 'voice', ['RAT_A', 'G Level 4', 'Call Family'])

    # Simulate a CDR that received its derived field after its first reporting
    # copy. The cache must not retain the original null value indefinitely.
    repository.replace_dataset_rows(1, pd.DataFrame({
        'RAT_A': ['EN-DC'], 'G_Level_4': ['London'], 'Call Family': ['VoLTE'],
    }))
    repository.copy_dataset_rows_to_reporting(1, 'voice', ['RAT_A', 'G Level 4', 'Call Family'])

    loaded = repository.load_reporting_rows('voice', [1], ['G Level 4', 'Call Family'])
    assert loaded.to_dict(orient='records') == [{'G Level 4': 'London', 'Call Family': 'VoLTE'}]


def test_reporting_cache_repairs_an_empty_requested_source_column(tmp_path) -> None:
    repository = Repository(tmp_path / 'workspace.db')
    repository.replace_dataset_rows(1, pd.DataFrame({'RAT': ['EN-DC'], 'G_Level_4': ['London']}))
    repository.copy_dataset_rows_to_reporting(1, 'data', ['RAT', 'G Level 4'])
    with repository.connection() as connection:
        connection.execute('UPDATE reporting_rows_data SET G_Level_4 = NULL WHERE dataset_id = 1')

    repository.copy_dataset_rows_to_reporting(1, 'data', ['RAT', 'G Level 4'])

    loaded = repository.load_reporting_rows('data', [1], ['G Level 4'])
    assert loaded.to_dict(orient='records') == [{'G Level 4': 'London'}]


def test_session_classification_and_multivendor_enrichment() -> None:
    cdr = pd.DataFrame({
        'Operator': ['Vodafone UK', '3'],
        'RAT_A': ['LTE EN-DC', 'NR'],
        'Cell_ID_A': ['100 -> 100', '200 -> 201'],
    })
    vodafone_mapping = pd.DataFrame({
        'source_sheet': ['4G'],
        'eNodeB ID': [0],
        'Local Cell ID': [100],
        'OP/ Vendor': ['Ericsson'],
    })
    three_mapping = pd.DataFrame({
        'Cid__ECI': [200, 201],
        'Vendor': ['Nokia', 'Ericsson'],
    })

    nsa = classify_sessions(cdr, 'nsa')
    sa = classify_sessions(cdr, 'sa')

    assert len(nsa) == 1
    assert len(sa) == 1
    assert enrich_multivendor(nsa, vodafone_mapping, three_mapping)['report_vendor'].tolist() == ['Vodafone_Ericsson']
    assert enrich_multivendor(sa, vodafone_mapping, three_mapping)['report_vendor'].tolist() == ['3_Mixed Vendor']


def test_speech_session_classification_uses_call_mode_when_sample_rat_is_blank() -> None:
    speech = pd.DataFrame({
        'sample': ['native-volte', 'native-vonr', 'whatsapp-endc', 'whatsapp-nr'],
        'Sample_RAT_A': [None, None, 'EN-DC', 'NR SA'],
        'L1_Call_Mode_A': ['VoLTE', 'VoNR', 'VoIP', 'VoIP'],
    })

    assert classify_sessions(speech, 'nsa')['sample'].tolist() == ['native-volte', 'whatsapp-endc']
    assert classify_sessions(speech, 'sa')['sample'].tolist() == ['native-vonr', 'whatsapp-nr']


def test_workspace_vendor_assignment_writes_the_normalized_vendor_field() -> None:
    cdr = pd.DataFrame({
        'Operator': ['Vodafone UK', '3', 'O2 (UK)'],
        'Cell_ID_A': ['100 -> 100', '200 -> 201', '300'],
    })
    vodafone_mapping = pd.DataFrame({
        'source_sheet': ['4G'], 'eNodeB ID': [0], 'Local Cell ID': [100], 'OP/ Vendor': ['Ericsson'],
    })
    three_mapping = pd.DataFrame({'Cid__ECI': [200, 201], 'Vendor': ['Nokia', 'Ericsson']})

    mapped = assign_cdr_vendors(cdr, vodafone_mapping, three_mapping)

    assert mapped['vendor'].tolist() == ['Vodafone_Ericsson', '3_Mixed Vendor', pd.NA]
    assert mapped['report_vendor'].tolist() == ['Vodafone_Ericsson', '3_Mixed Vendor', 'O2 (UK)']
    assert mapped.columns[:2].tolist() == ['vendor', 'Operator']


def test_workspace_vendor_assignment_replaces_source_vendor_collisions() -> None:
    cdr = pd.DataFrame({
        'source_sheet': ['Vodafone'],
        'Vendor': ['legacy source value'],
        'vendor__2': ['normalised source value'],
        'Operator': ['3'],
        'Cell_ID_A': ['200 -> 200'],
    })
    three_mapping = pd.DataFrame({'Cid__ECI': [200], 'Vendor': ['Nokia']})

    mapped = assign_cdr_vendors(cdr, None, three_mapping)

    assert mapped.columns[:2].tolist() == ['source_sheet', 'vendor']
    assert 'Vendor' not in mapped.columns
    assert 'vendor__2' not in mapped.columns
    assert mapped.loc[0, 'vendor'] == '3_Nokia'
    assert mapped.columns[-1] == 'report_vendor'


def test_workspace_vendor_assignment_supports_a_single_selected_mapping() -> None:
    cdr = pd.DataFrame({
        'Operator': ['Vodafone UK', '3'],
        'Cell_ID_A': ['100 -> 100', '200 -> 200'],
    })
    vodafone_mapping = pd.DataFrame({
        'source_sheet': ['4G'], 'eNodeB ID': [0], 'Local Cell ID': [100], 'OP/ Vendor': ['Ericsson'],
    })

    mapped = assign_cdr_vendors(cdr, vodafone_mapping, None)

    assert mapped.loc[0, 'vendor'] == 'Vodafone_Ericsson'
    assert pd.isna(mapped.loc[1, 'vendor'])
    assert mapped['report_vendor'].tolist() == ['Vodafone_Ericsson', '3']


def test_workspace_vendor_assignment_accepts_equivalent_global_cell_id_columns() -> None:
    cdr = pd.DataFrame({
        'Operator': ['3'],
        'Global CI': ['200 -> 200'],
    })
    three_mapping = pd.DataFrame({'Cid__ECI': [200], 'Vendor': ['Nokia']})

    mapped = assign_cdr_vendors(cdr, None, three_mapping)

    assert mapped['vendor'].tolist() == ['3_Nokia']


def test_catalogue_converter_migrates_legacy_headers_and_grouping() -> None:
    legacy = (
        'Slide,Slide title,Slide subtitle,Layout,CDR Source,KPI,Chart Type,Filters,Grouping\n'
        '8,Quality,Voice,Title and 1 column + Comments,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator × Campaign\n'
    )

    converted = convert_catalog_csv(legacy, 'nsa')
    entries = parse_catalog_csv(converted, 'nsa')

    assert converted.decode('utf-8').splitlines()[0] == ','.join(CATALOG_HEADERS)
    assert entries[0].slide_title == 'Quality'
    assert entries[0].chart_title == ''
    assert entries[0].legend == ''
    assert entries[0].grouping_rows == 'Operator'
    assert entries[0].grouping_columns == 'Campaign'


def test_catalogue_converter_assigns_layouts_for_missing_legacy_layouts() -> None:
    legacy = (
        'Slide,Slide title,Slide subtitle,Layout,CDR Source,KPI,Chart Type,Filters,Grouping\n'
        '9,Quality,,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator × Campaign\n'
        '9,Quality,,,CDR-Voice,Call_Setup_Time,Average Vertical Bars,,Operator × Campaign\n'
    )

    entries = parse_catalog_csv(convert_catalog_csv(legacy, 'nsa'), 'nsa')

    assert {entry.layout for entry in entries} == {'Title and 2 columns + Comments'}


def test_legacy_workspace_mapping_gets_the_report_group_without_writing_operator_as_vendor() -> None:
    frame = pd.DataFrame({'Operator': ['Vodafone UK', 'O2 (UK)'], 'vendor': ['Vodafone_Ericsson', pd.NA]})

    grouped = ensure_report_vendor_group(frame)

    assert grouped['report_vendor'].tolist() == ['Vodafone_Ericsson', 'O2 (UK)']


def test_vodafone_mapping_derives_gcid_from_4g_enodeb_and_local_cell() -> None:
    cdr = pd.DataFrame({'Operator': ['Vodafone UK'], 'Cell_IDs_A': ['3330049']})
    vodafone_mapping = pd.DataFrame({
        'source_sheet': ['4G'],
        'eNodeB ID': [13008],
        'Local Cell ID': [1],
        'OP/ Vendor': ['Samsung'],
    })
    three_mapping = pd.DataFrame({'Cid__ECI': [1], 'Vendor': ['Nokia']})

    assert enrich_multivendor(cdr, vodafone_mapping, three_mapping)['report_vendor'].tolist() == ['Vodafone_Samsung']


def test_catalogue_csv_requires_the_report_chart_contract_columns() -> None:
    catalogue = (
        ','.join(CATALOG_HEADERS)
        + '\n8,Completed Call Ratio,Voice quality,Title and 1 column + Comments,Completed call ratio,CDR-Voice,Call_Status,100% Stacked Vertical Bars,Call Family = VoLTE,Operator,Campaign,Completed/Dropped/Failed,\n'
    ).encode('utf-8')

    entries = parse_catalog_csv(catalogue, 'nsa')

    assert entries[0].slide == 8
    assert entries[0].slide_title == 'Completed Call Ratio'
    assert entries[0].slide_subtitle == 'Voice quality'
    assert entries[0].chart_title == 'Completed call ratio'
    assert entries[0].legend == 'Completed/Dropped/Failed'
    assert entries[0].source_kind == 'voice'
    assert entries[0].chart_type == '100% Stacked Vertical Bars'


def test_catalogue_parses_legend_position_and_accepts_prior_schema() -> None:
    current = (
        ','.join(CATALOG_HEADERS)
        + '\n8,Quality,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,Completed/Dropped/Failed,Left\n'
    )
    entry = parse_catalog_csv(current, 'nsa')[0]

    assert entry.legend_position == 'left'
    assert parse_legend_position('Bottom') == 'bottom'
    with pytest.raises(ValueError, match='Legend Position'):
        parse_legend_position('Centre')

    previous = 'Slide,Slide tittle,Slide Subtittle,Layout,Chart Tittle,CDR source,KPI,Chart type,Legend,Filters,Grouping_Rows,Grouping_Columns' + '\n8,Quality,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,,Operator,Campaign\n'
    assert parse_catalog_csv(previous, 'nsa')[0].legend_position == 'top'
    two_columns = ','.join(CATALOG_HEADERS) + '\n8,Quality,,Title and 2 columns + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,Operator,Campaign,,\n'
    assert parse_catalog_csv(two_columns, 'nsa')[0].legend_position == 'top'


def test_status_chart_draws_legend_at_the_catalogue_position() -> None:
    frame = pd.DataFrame({
        'Operator': ['Vodafone'],
        'Campaign': ['2026 Q1'],
        'Call_Status': ['Completed'],
    })

    with patch('src.modules.cdr_reporting._draw_chart_legend') as draw_legend:
        _render_status_100('Status', frame, 'Operator', 'Campaign', legend_position='left')

    assert draw_legend.call_args.args[2] == 'left'


def test_catalogue_accepts_structural_slides_and_rejects_chart_configuration_on_them() -> None:
    title = parse_catalog_csv(
        ','.join(CATALOG_HEADERS) + '\n1,Quarterly report,NSA analysis,Title Page,,,,Title Slide,,,,,\n',
        'nsa',
    )[0]
    transition = parse_catalog_csv(
        ','.join(CATALOG_HEADERS) + '\n2,Voice analysis,Seven cities,Title Only,,,,Transition Slide,,,,,\n',
        'nsa',
    )[0]

    assert title.structural_type == 'title slide'
    assert transition.structural_type == 'transition slide'
    with pytest.raises(ValueError, match='cannot define chart, CDR, KPI'):
        parse_catalog_csv(
            ','.join(CATALOG_HEADERS) + '\n1,Quarterly report,,Title Page,,CDR-Data,LQ,Title Slide,,,,,\n',
            'nsa',
        )


def test_catalogue_filter_and_grouping_contract_is_parsed_and_applied() -> None:
    conditions = parse_catalog_filters('Session_Type IN (VoLTE, MultiRAB); LQ >= 1.6')
    grouping = parse_catalog_grouping('City × Operator × Campaign')
    assert [(item.column, item.operator, item.values) for item in conditions] == [
        ('Session_Type', 'IN', ('VoLTE', 'MultiRAB')), ('LQ', '>=', ('1.6',)),
    ]
    assert grouping.dimensions == ('City', 'Operator', 'Campaign')
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS) + '\n8,Quality,,Title and 1 column + Comments,Quality by city,CDR-Speech,LQ,Average Vertical Bars,Session_Type IN (VoLTE); LQ >= 1.6,City,Operator × Campaign,,\n',
        'nsa',
    )[0]
    frame = pd.DataFrame({
        'Session_Type': ['VoLTE', 'WhatsApp'], 'LQ': [3.2, 4.0], 'City': ['London', 'Leeds'],
        'Operator': ['EE', 'O2'], 'Campaign': ['Q1', 'Q1'],
    })
    filtered = _apply_catalog_filters(frame, entry, False, 'LQ')
    grouped, primary, series = _apply_catalog_grouping(filtered, entry, False, 'LQ')
    assert len(grouped) == 1
    assert grouped[primary].tolist() == ['London']
    assert grouped[series].tolist() == ['EE · Q1']
    assert '__catalog_stack' not in grouped.columns


def test_multivendor_rendering_rewrites_operator_display_and_grouping_but_not_filters() -> None:
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS)
        + '\n8,Operator comparison,Operator subtitle,Title and 1 column + Comments,Operator chart,CDR-Speech,LQ,Average Vertical Bars,Operator = Vodafone UK,Operator,Operator × Campaign,Operator,\n',
        'nsa',
    )[0]
    rendered = prepare_multivendor_catalog_entry(entry)

    assert rendered.slide_title == 'Vendor comparison'
    assert rendered.slide_subtitle == 'Vendor subtitle'
    assert rendered.chart_title == 'Vendor chart'
    assert rendered.legend == 'Campaign'
    assert rendered.grouping_rows == 'Vendor'
    assert rendered.grouping_columns == 'Vendor × Campaign'
    assert rendered.filters == 'Operator = Vodafone UK'

    frame = pd.DataFrame({
        'Operator': ['Vodafone UK', '3'],
        'report_vendor': ['Vodafone_Ericsson', '3_Nokia'],
        'Campaign': ['UK_Q2_SA_2026', 'UK_Q2_SA_2026'],
        'LQ': [3.8, 3.5],
    })
    filtered = _apply_catalog_filters(frame, rendered, True, 'LQ')
    grouped, primary, series = _apply_catalog_grouping(filtered, rendered, True, 'LQ')
    assert grouped[primary].tolist() == ['Vodafone_Ericsson']
    assert grouped[series].tolist() == ['Vodafone_Ericsson · 2026-Q2']


def test_rows_only_grouping_uses_one_all_series_without_repeating_the_category() -> None:
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS) + '\n8,Quality,,Title and 1 column + Comments,Quality,CDR-Speech,LQ,Average Vertical Bars,,Operator,,,\n',
        'nsa',
    )[0]
    frame = pd.DataFrame({'Operator': ['EE', 'O2'], 'LQ': [3.2, 3.8]})

    grouped, primary, series = _apply_catalog_grouping(frame, entry, False, 'LQ')

    assert grouped[primary].tolist() == ['EE', 'O2']
    assert grouped[series].tolist() == ['(all)', '(all)']


def test_campaign_grouping_displays_only_year_and_quarter() -> None:
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS) + '\n8,Quality,,Title and 1 column + Comments,Quality,CDR-Speech,LQ,Average Vertical Bars,,Operator,Campaign,,\n',
        'nsa',
    )[0]
    frame = pd.DataFrame({
        'Operator': ['EE', '3', 'Vodafone UK'],
        'Campaign': ['UK_Q2_SA_2026', 'UK_Q4_2025', '2024 Q3 NSA'],
        'LQ': [3.2, 3.8, 4.0],
    })

    grouped, _primary, series = _apply_catalog_grouping(frame, entry, False, 'LQ')

    assert grouped['Campaign'].tolist() == ['UK_Q2_SA_2026', 'UK_Q4_2025', '2024 Q3 NSA']
    assert grouped['__catalog_column_0'].tolist() == ['2026-Q2', '2025-Q4', '2024-Q3']
    assert grouped[series].tolist() == ['2026-Q2', '2025-Q4', '2024-Q3']


def test_cdf_renders_a_curve_for_each_complete_rows_and_columns_combination() -> None:
    frame = pd.DataFrame({
        '__catalog_row_0': ['Vodafone', 'Vodafone', 'Vodafone', 'Vodafone', 'O2', 'O2', 'O2', 'O2'],
        '__catalog_column_0': ['2025', '2025', '2026', '2026', '2025', '2025', '2026', '2026'],
        '__catalog_primary': ['unused'] * 8,
        '__catalog_series': ['unused'] * 8,
        'Campaign': ['2025', '2025', '2026', '2026', '2025', '2025', '2026', '2026'],
        'Metric': [1.0, 2.0, 1.5, 2.5, 1.2, 2.2, 1.7, 2.7],
    })

    with patch('src.modules.cdr_reporting._draw_chart_legend') as draw_legend:
        _render_cdf_line('CDF', frame, '__catalog_primary', '__catalog_series', 'Metric')

    legend_items = draw_legend.call_args.args[1]
    assert [item[0] for item in legend_items] == [
        'Vodafone · 2025', 'Vodafone · 2026', 'O2 · 2025', 'O2 · 2026',
    ]


def test_cdf_uses_emphasised_lines_when_only_one_campaign_is_rendered() -> None:
    frame = pd.DataFrame({
        '__catalog_row_0': ['Vodafone', 'Vodafone', 'O2', 'O2'],
        '__catalog_primary': ['unused'] * 4, '__catalog_series': ['unused'] * 4,
        'Campaign': ['2026 Q2'] * 4, 'Metric': [1.0, 2.0, 1.2, 2.2],
    })

    with patch('src.modules.cdr_reporting._draw_chart_legend') as draw_legend:
        _render_cdf_line('CDF', frame, '__catalog_primary', '__catalog_series', 'Metric')

    assert {item[2] for item in draw_legend.call_args.args[1]} == {4}


def test_cdf_trims_only_a_converged_tail_after_three_curves_reach_98_percent() -> None:
    assert _cdf_terminal_x_maximum(
        [[1.0] * 99 + [10.0], [2.0] * 99 + [10.0], [3.0] * 99 + [10.0], [4.0] * 70 + [30.0] * 30],
        1.0, 30.0, minimum_separation=0.08,
    ) == 3.0
    # Fewer than three completed curves never make a tail eligible.
    assert _cdf_terminal_x_maximum([[1.0, 2.0, 3.0, 10.0], [1.0, 2.0, 3.0, 10.0]], 1.0, 10.0) == 10.0
    # Exactly 98% does not satisfy the strictly-above-98% threshold.
    assert _cdf_terminal_x_maximum(
        [[1.0] * 99 + [10.0], [2.0] * 99 + [10.0], [3.0] * 98 + [10.0, 11.0]],
        1.0, 10.0,
    ) == 10.0


def test_operator_vendor_column_groups_keep_campaign_bars_in_their_operator_palette() -> None:
    colours = _hierarchy_group_colours([
        ('Vodafone_Ericsson', '2026 Q1'), ('Vodafone_Ericsson', '2026 Q2'),
        ('Vodafone_Huawei', '2026 Q1'), ('Vodafone_Huawei', '2026 Q2'),
        ('3_Ericsson', '2026 Q1'), ('3_Ericsson', '2026 Q2'),
    ])

    assert colours['Vodafone_Ericsson'] == '#E15759'
    assert colours['Vodafone_Huawei'] == '#9B1D20'
    assert colours['3_Ericsson'] == '#F28E2B'


def test_chart_colours_use_vendor_families_for_multi_operator_dimensions() -> None:
    keys = [('Vodafone', 'Ericsson'), ('Vodafone', 'Huawei'), ('3', 'Ericsson'), ('3', 'Huawei')]
    frame = pd.DataFrame({'__catalog_row_0': [], '__catalog_row_1': []})
    frame.attrs['catalogue_dimension_labels'] = {
        '__catalog_row_0': ('Operator',), '__catalog_row_1': ('Vendor',),
    }

    colours = _series_colours(keys, ['__catalog_row_0', '__catalog_row_1'], frame)

    assert colours[('Vodafone', 'Ericsson')] == '#2E8B57'
    assert colours[('3', 'Ericsson')] == '#0D5A34'
    assert colours[('Vodafone', 'Huawei')] == '#E15759'
    assert colours[('3', 'Huawei')] == '#A61E2B'

    line_colours = _series_colours(keys, ['__catalog_row_0', '__catalog_row_1'], frame, line_chart=True)
    assert line_colours == colours


def test_chart_colours_use_vendor_families_for_one_operator() -> None:
    keys = [('Vodafone', 'Ericsson'), ('Vodafone', 'Huawei'), ('Vodafone', 'Samsung'), ('Vodafone', 'NSN')]
    frame = pd.DataFrame({'__catalog_row_0': [], '__catalog_row_1': []})
    frame.attrs['catalogue_dimension_labels'] = {
        '__catalog_row_0': ('Operator',), '__catalog_row_1': ('Vendor',),
    }

    colours = _series_colours(keys, ['__catalog_row_0', '__catalog_row_1'], frame)

    assert colours == {
        ('Vodafone', 'Ericsson'): '#2E8B57', ('Vodafone', 'Huawei'): '#E15759',
        ('Vodafone', 'Samsung'): '#D9A514', ('Vodafone', 'NSN'): '#4E79A7',
    }


def test_chart_colours_detect_a_single_operator_from_composite_vendor_values() -> None:
    keys = [('3_Ericsson',), ('3_Huawei',), ('3_Nokia',)]
    frame = pd.DataFrame({'__catalog_column_0': []})
    frame.attrs['catalogue_dimension_labels'] = {'__catalog_column_0': ('Vendor',)}

    colours = _series_colours(keys, ['__catalog_column_0'], frame)

    assert colours[('3_Ericsson',)] == '#2E8B57'
    assert len(set(colours.values())) == 3


def test_catalogue_filter_contract_supports_not_in_and_not_contains() -> None:
    conditions = parse_catalog_filters('Session_Type NOT IN (WhatsApp, SMS); Vendor NOT CONTAINS (Mixed, Other); Campaign NOT CONTAINS legacy')
    assert [(item.column, item.operator, item.values) for item in conditions] == [
        ('Session_Type', 'NOT IN', ('WhatsApp', 'SMS')),
        ('Vendor', 'NOT CONTAINS', ('Mixed', 'Other')),
        ('Campaign', 'NOT CONTAINS', ('legacy',)),
    ]
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS) + '\n8,Quality,,Title and 1 column + Comments,Quality,CDR-Speech,LQ,Average Vertical Bars,Session_Type NOT IN (WhatsApp); Campaign NOT CONTAINS legacy,Operator,Campaign,,\n',
        'nsa',
    )[0]
    frame = pd.DataFrame({
        'Session_Type': ['VoLTE', 'WhatsApp', 'VoLTE'], 'Campaign': ['Q1', 'Q1', 'legacy-Q2'],
        'LQ': [3.2, 4.0, 3.8], 'Operator': ['EE', 'EE', 'O2'],
    })
    assert _apply_catalog_filters(frame, entry, False, 'LQ')['Operator'].tolist() == ['EE']


def test_not_contains_filter_excludes_each_comma_separated_term() -> None:
    entry = CatalogEntry(
        1, 'Quality', '', 'Title and 1 column', '', 'CDR-Speech', 'LQ', 'CDF Line',
        '', 'Vendor NOT CONTAINS (Mixed, Other)', 'Vendor', 'Campaign', 'Top',
    )
    frame = pd.DataFrame({
        'Vendor': ['Vodafone_Ericsson', 'Vodafone_Mixed Vendor', '3_Other Vendor'],
        'Campaign': ['2026 Q1'] * 3,
        'LQ': [3.5, 3.6, 3.7],
    })

    assert _apply_catalog_filters(frame, entry, False, 'LQ')['Vendor'].tolist() == ['Vodafone_Ericsson']


def test_catalogue_call_family_uses_documented_netcheck_session_values() -> None:
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS)
        + '\n8,Completed Call Ratio,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,"Call Family IN (VoLTE, MultiRAB, WhatsApp)",Call Family,Operator × Campaign,,\n',
        'nsa',
    )[0]
    frame = pd.DataFrame({
        'Session_Type': ['CALL', 'MultiRAB CALL', 'WhatsApp CALL'],
        'L1_Call_Mode_A': ['VoLTE', '', ''],
        'Operator': ['EE', 'EE', 'EE'],
        'Campaign': ['Q1', 'Q1', 'Q1'],
        'Call_Status': ['Completed', 'Completed', 'Completed'],
    })

    filtered = _apply_catalog_filters(frame, entry, False, 'Call_Status')
    grouped, primary, series = _apply_catalog_grouping(filtered, entry, False, 'Call_Status')

    assert grouped[primary].tolist() == ['VoLTE', 'MultiRAB', 'WhatsApp']
    assert grouped['__catalog_row_0'].tolist() == ['VoLTE', 'MultiRAB', 'WhatsApp']
    assert grouped['__catalog_column_0'].tolist() == ['EE', 'EE', 'EE']
    assert grouped['__catalog_column_1'].tolist() == ['Q1', 'Q1', 'Q1']

    with patch('src.modules.cdr_reporting._render_status_100_hierarchy') as hierarchy_renderer:
        hierarchy_renderer.return_value = BytesIO(b'nested-chart')
        chart = _render_status_100('Completed Call Ratio', grouped, primary, series)

    assert chart.getvalue() == b'nested-chart'
    assert hierarchy_renderer.call_args.args[2] == ['__catalog_row_0']
    assert hierarchy_renderer.call_args.args[3] == ['__catalog_column_0', '__catalog_column_1']


def test_status_chart_uses_nested_columns_without_a_row_grouping() -> None:
    entry = CatalogEntry(
        8, 'Completed Call Ratio', '', 'Title and 1 column + Comments', '', 'CDR-Voice',
        'Call_Status', '100% Stacked Vertical Bars', '', '', '', 'Operator × Campaign',
    )
    frame = pd.DataFrame({
        'Operator': ['Vodafone', 'Vodafone', 'O2', 'O2'],
        'Campaign': ['2025 Q4', '2026 Q1', '2025 Q4', '2026 Q1'],
        'Call_Status': ['Completed', 'Failed', 'Completed', 'Dropped'],
    })
    grouped, primary, series = _apply_catalog_grouping(frame, entry, False, 'Call_Status')

    with patch('src.modules.cdr_reporting._render_status_100_hierarchy') as hierarchy_renderer:
        hierarchy_renderer.return_value = BytesIO(b'nested-columns')
        chart = _render_status_100('Completed Call Ratio', grouped, primary, series)

    assert chart.getvalue() == b'nested-columns'
    assert hierarchy_renderer.call_args.args[2] == []
    assert hierarchy_renderer.call_args.args[3] == ['__catalog_column_0', '__catalog_column_1']


def test_hierarchical_grouping_keeps_campaign_bars_together_per_operator() -> None:
    frame = pd.DataFrame({
        '__catalog_column_0': ['Vodafone', 'O2', 'Vodafone', 'O2'],
        '__catalog_column_1': ['2025 Q4', '2025 Q4', '2026 Q1', '2026 Q1'],
    })

    keys = _hierarchical_unique_keys(frame, ['__catalog_column_0', '__catalog_column_1'])

    assert keys == [
        ('Vodafone', '2025 Q4'),
        ('Vodafone', '2026 Q1'),
        ('O2', '2025 Q4'),
        ('O2', '2026 Q1'),
    ]


def test_nsa_speech_catalogue_filters_produce_samples_and_use_latest_campaign() -> None:
    entries = load_catalog_csv(Path(__file__).parent / 'fixtures' / 'NSA Slide Template.csv', 'nsa')
    speech = pd.DataFrame({
        'sample': ['volte', 'multirab', 'whatsapp-old', 'whatsapp-latest', 'whatsapp-sa', 'o2-latest'],
        'Session_Type': ['CALL', 'MultiRAB CALL', 'WhatsApp CALL', 'WhatsApp CALL', 'WhatsApp CALL', 'WhatsApp CALL'],
        'L1_Call_Mode_A': ['VoLTE', 'VoLTE', 'VoIP', 'VoIP', 'VoIP', 'VoIP'],
        'Sample_RAT_A': [None, None, 'EN-DC', 'EN-DC', 'NR SA', 'EN-DC'],
        'Call_Status': ['Completed'] * 6,
        'Operator': ['Vodafone UK', '3', 'EE', 'EE', 'EE', 'O2 (UK)'],
        'Campaign': ['UK_Q3_2025', 'UK_Q3_2025', 'UK_Q3_2025', 'UK_Q4_2025', 'UK_Q4_2025', 'UK_Q4_2025'],
        'LQ': [3.8, 3.7, 3.9, 4.0, 4.1, 4.2],
    })
    nsa = classify_sessions(speech, 'nsa')
    filtered_by_entry = {
        (entry.slide, entry.chart_type, index): _apply_catalog_filters(nsa, entry, False, 'LQ')
        for index, entry in enumerate(entries)
        if entry.slide in {7, 8, 9}
    }

    assert all(not frame.empty for frame in filtered_by_entry.values())
    latest_whatsapp = [frame for (slide, _chart, _index), frame in filtered_by_entry.items() if slide == 8][2]
    assert latest_whatsapp['sample'].tolist() == ['whatsapp-latest']


def test_layout_chart_frames_are_always_ordered_by_visual_rows_then_columns() -> None:
    presentation = Presentation('assets/ppt-templates/Template_CDR_analysis.pptx')
    layout = _named_slide_layout(presentation, 'Title and 2 columns and 2 rows + Comments right')

    frames = _layout_chart_frames(layout)

    assert len(frames) == 4
    assert frames[0][0] < frames[1][0]
    assert frames[0][1] < frames[2][1]
    assert frames[2][0] < frames[3][0]


def test_failure_count_uses_row_and_column_hierarchies_without_flattening() -> None:
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS)
        + '\n9,Voice failures per Q/city,,Title and 1 column + Comments,Failures,CDR-Voice,Call_Status,Count Stacked Horizontal Bars,,Call Family × G Level 4,Operator × Campaign,Failed/Dropped,\n',
        'nsa',
    )[0]
    frame = pd.DataFrame({
        'Session_Type': ['VoLTE', 'VoLTE', 'MultiRAB CALL'],
        'G_Level_4': ['London', 'London', 'Belfast'],
        'Operator': ['EE', 'EE', '3'],
        'Campaign': ['Q2', 'Q2', 'Q3'],
        'Call_Status': ['Failed', 'Dropped', 'Failed'],
    })
    grouped, primary, series = _apply_catalog_grouping(frame, entry, False, 'Call_Status')

    with patch('src.modules.cdr_reporting._render_failure_count_hierarchy') as hierarchy_renderer:
        hierarchy_renderer.return_value = BytesIO(b'nested-failure-chart')
        chart = _render_failure_count('Voice failures per Q/city', grouped, primary, series)

    assert chart.getvalue() == b'nested-failure-chart'
    assert hierarchy_renderer.call_args.args[2] == ['__catalog_row_0', '__catalog_row_1']
    assert hierarchy_renderer.call_args.args[3] == ['__catalog_column_0', '__catalog_column_1']


def test_failure_count_keeps_zero_count_hierarchy_categories_from_all_filtered_rows() -> None:
    entry = parse_catalog_csv(
        ','.join(CATALOG_HEADERS)
        + '\n9,Voice failures,,Title and 1 column + Comments,Failures,CDR-Voice,Call_Status,Count Stacked Horizontal Bars,,Call Family,Operator × Campaign,Failed/Dropped,\n',
        'nsa',
    )[0]
    frame = pd.DataFrame({
        'Session_Type': ['VoLTE', 'VoLTE'], 'Operator': ['Vodafone', '3'],
        'Campaign': ['Q2', 'Q2'], 'Call_Status': ['Completed', 'Completed'],
    })
    grouped, primary, series = _apply_catalog_grouping(frame, entry, False, 'Call_Status')

    with patch('src.modules.cdr_reporting._render_failure_count_hierarchy') as hierarchy_renderer:
        hierarchy_renderer.return_value = BytesIO(b'zero-count-grid')
        chart = _render_failure_count('Voice failures', grouped, primary, series)

    assert chart.getvalue() == b'zero-count-grid'
    assert hierarchy_renderer.call_args.kwargs['comparison_frame'] is grouped
    assert hierarchy_renderer.call_args.args[1].empty


def test_failure_hierarchy_reserves_a_right_legend_lane() -> None:
    frame = pd.DataFrame({
        '__catalog_row_0': ['VoLTE'],
        '__catalog_column_0': ['Vodafone'],
        '__catalog_column_1': ['2026 Q2'],
        '__catalog_failure_state': ['Failed'],
    })

    with patch('src.modules.cdr_reporting._draw_chart_legend') as draw_legend:
        _render_failure_count_hierarchy(
            'Voice failures per Q/city', frame,
            ['__catalog_row_0'], ['__catalog_column_0', '__catalog_column_1'],
            legend_position='right',
        )

    assert draw_legend.call_args.kwargs['side_x'] == 1289


def test_failure_hierarchy_uses_dashed_child_boundaries_within_one_operator() -> None:
    frame = pd.DataFrame({
        '__catalog_row_0': ['VoLTE', 'VoLTE'],
        '__catalog_column_0': ['Vodafone', 'Vodafone'],
        '__catalog_column_1': ['2026 Q2', '2026 Q1'],
        '__catalog_failure_state': ['Failed', 'Dropped'],
    })

    with patch('src.modules.cdr_reporting._draw_dashed_vertical_line') as draw_dashed:
        _render_failure_count_hierarchy(
            'Voice failures per campaign', frame,
            ['__catalog_row_0'], ['__catalog_column_0', '__catalog_column_1'],
        )

    assert draw_dashed.call_count == 1


def test_failure_hierarchy_uses_dashed_child_boundaries_within_one_row_group() -> None:
    frame = pd.DataFrame({
        '__catalog_row_0': ['VoLTE', 'VoLTE', 'MultiRAB'],
        '__catalog_row_1': ['Belfast', 'Bristol', 'Belfast'],
        '__catalog_column_0': ['VF', 'VF', 'VF'],
        '__catalog_failure_state': ['Failed', 'Dropped', 'Failed'],
    })

    with patch('src.modules.cdr_reporting._draw_dashed_horizontal_line') as draw_dashed:
        _render_failure_count_hierarchy(
            'Voice failures per city', frame,
            ['__catalog_row_0', '__catalog_row_1'], ['__catalog_column_0'],
        )

    assert draw_dashed.call_count == 2


def test_nsa_catalogue_splits_template_screenshots_into_individual_charts() -> None:
    entries = load_catalog_csv(Path(__file__).parent / 'fixtures' / 'NSA Slide Template.csv', 'nsa')
    slide_ten = [entry for entry in entries if entry.slide == 10]
    slide_thirteen = [entry for entry in entries if entry.slide == 13]

    assert len(slide_ten) == 2
    assert {entry.layout for entry in slide_ten} == {'Title and 2 columns + Comments'}
    assert len(slide_thirteen) == 3
    assert {entry.layout for entry in slide_thirteen} == {'Title and 3 columns + Comments'}
    assert {slide: sum(entry.slide == slide for entry in entries) for slide in range(12, 17)} == {
        12: 2, 13: 3, 14: 3, 15: 4, 16: 2,
    }


def test_catalogue_uses_explicit_title_and_transition_slides() -> None:
    entries = load_catalog_csv(Path(__file__).parent / 'fixtures' / 'NSA Slide Template.csv', 'nsa')
    structural = [entry.chart_type for entry in entries if not entry.source_kind]
    title = next(entry for entry in entries if entry.slide == 1)
    conclusions = next(entry for entry in entries if entry.slide == 17)

    assert set(structural) == {'Title Slide', 'Transition Slide'}
    assert (title.chart_type, title.layout) == ('Title Slide', 'Title Page')
    assert (conclusions.chart_type, conclusions.layout) == ('Transition Slide', 'Title Only')


def test_catalogue_rows_use_matching_master_image_placeholders(tmp_path) -> None:
    catalogue = (
        ','.join(CATALOG_HEADERS)
        + '\n8,Completed Call Ratio,Voice quality,Title and 2 rows + Comments right,Status ratio,CDR-Voice,Call_Status,100% Stacked Vertical Bars,Call Family = VoLTE,Call Family,Operator × Campaign,Completed/Dropped/Failed,'
        + '\n8,Completed Call Ratio,Voice quality,Title and 2 rows + Comments right,Setup time,CDR-Voice,Call_Setup_Time,Average Vertical Bars,,Call Family = VoLTE,Call Family,Operator × Campaign\n'
    ).encode('utf-8')
    frames = {
        'data': pd.DataFrame(),
        'speech': pd.DataFrame(),
        'voice': pd.DataFrame({
            'Campaign': ['Q1'], 'Operator': ['EE'], 'Session_Type': ['VoLTE'],
            'Call_Status': ['Completed'], 'Call_Setup_Time': [1.2],
        }),
    }
    destination = tmp_path / 'catalogue-layout.pptx'

    render_cdr_report(
        destination,
        Path('assets/ppt-templates/Template_CDR_analysis.pptx'),
        frames,
        'nsa',
        False,
        parse_catalog_csv(catalogue, 'nsa'),
    )

    generated = Presentation(destination)
    assert len(generated.slides) == 1
    slide = generated.slides[0]
    assert slide.slide_layout.name == 'Title and 2 rows + Comments right'
    pictures = sorted((shape for shape in slide.shapes if hasattr(shape, 'image')), key=lambda shape: shape.top)
    assert len(pictures) >= 2
    assert pictures[0].top < pictures[1].top
    comments = next(shape for shape in slide.placeholders if shape.placeholder_format.idx == 10)
    layout_comments = next(shape for shape in slide.slide_layout.placeholders if shape.placeholder_format.idx == 10)
    assert (comments.left, comments.top, comments.width, comments.height) == (
        layout_comments.left, layout_comments.top, layout_comments.width, layout_comments.height,
    )
    title_shape = next(
        shape for shape in slide.shapes
        if getattr(shape, 'has_text_frame', False) and getattr(shape, 'is_placeholder', False)
        and shape.placeholder_format.type in {1, 3}
    )
    assert [paragraph.text for paragraph in title_shape.text_frame.paragraphs] == ['Completed Call Ratio', 'Voice quality']
    subtitle_paragraph = title_shape.text_frame.paragraphs[1]
    assert subtitle_paragraph.font.size.pt == 16
    assert subtitle_paragraph.font.color.rgb == RGBColor(36, 90, 150)
    assert not any(shape.name == 'catalogue-subtitle' for shape in slide.shapes)


def test_layout_only_template_builds_one_new_slide_per_catalogue_number(tmp_path) -> None:
    template = Path('assets/ppt-templates/Template_CDR_analysis.pptx')
    assert len(Presentation(template).slides) == 0
    catalogue = (
        ','.join(CATALOG_HEADERS)
        + '\n1,Quarterly report,NSA analysis,Title Page,,,,Title Slide,,,,,'
        + '\n2,Voice section,Seven cities,Title Only,,,,Transition Slide,,,,'
        + '\n8,Completed Call Ratio,,Title and 1 column + Comments,,CDR-Voice,Call_Status,100% Stacked Vertical Bars,,,Operator,Campaign\n'
    )
    destination = tmp_path / 'catalogue-built.pptx'

    render_cdr_report(
        destination,
        template,
        {'data': pd.DataFrame(), 'speech': pd.DataFrame(), 'voice': pd.DataFrame()},
        'nsa',
        False,
        parse_catalog_csv(catalogue, 'nsa'),
    )

    generated = Presentation(destination)
    assert len(generated.slides) == 3
    assert [slide.slide_layout.name for slide in generated.slides] == [
        'Title Page', 'Title Only', 'Title and 1 column + Comments',
    ]
    assert generated.slides[0].placeholders[0].text == 'Quarterly report'
    assert generated.slides[0].placeholders[1].text == 'NSA analysis'


def test_reporting_module_is_available_to_authenticated_users(client) -> None:
    response = client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    assert response.status_code == 303

    page = client.get('/reporting')

    assert page.status_code == 200
    assert 'NetCheck CDR Reports' in page.text
    assert 'Smart Orchestrator Logs Reports' in page.text
    assert 'data-reporting-module' in page.text
    assert 'data-reporting-module-panel="cdr"' in page.text
    assert 'data-reporting-module-panel="logs"' in page.text
    assert 'Generate PowerPoint Report' in page.text
    assert 'name="vodafone_mapping_dataset_id"' not in page.text
    assert 'name="three_mapping_dataset_id"' not in page.text
    assert '<option value="multivendor" disabled>Multivendor Comparison</option>' in page.text
    assert 'name="slides_templates"' in page.text
    assert 'value="nsa:NSA Slide Template"' in page.text
    assert 'data-report-job-form' in page.text
    assert 'data-report-multicampaign-dialog' in page.text
    assert 'Review selected campaigns' in page.text
    assert 'latest selected CDR for each type is preselected' in page.text
    assert 'campaignPeriod' in page.text
    assert 'uploadedRecency' in page.text
    assert 'latestCampaignOption' in page.text
    assert "select.dispatchEvent(new Event('change', {bubbles: true}))" in page.text
    assert 'Reports Jobs' in page.text
    assert 'Charts Jobs' in page.text
    assert 'data-report-job-stop' in page.text
    assert 'data-report-chart-job-stop' in page.text


def test_processing_report_and_chart_jobs_can_be_stopped_then_deleted(client) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    report_id = app_module.repository.create_report_job(
        report_type='netcheck_cdr', technology='nsa', scope='single',
        data_dataset_id=1, voice_dataset_id=2, speech_dataset_id=3,
        dataset_ids={'data': [1], 'voice': [2], 'speech': [3]},
        dataset_names={'data': ['Data'], 'voice': ['Voice'], 'speech': ['Speech']},
        slide_count=1, template_name='NSA Slide Template', output_file='stop-test.pptx',
        output_path=app_module.settings.output_dir / 'reports' / 'stop-test.pptx', created_by='admin',
    )
    app_module.repository.update_report_job(report_id, status='processing', progress=40)
    stopped_report = client.post(f'/reporting/jobs/{report_id}/stop')
    assert stopped_report.status_code == 200
    report = next(item for item in client.get('/api/reporting/jobs').json()['jobs'] if item['id'] == report_id)
    assert report['status'] == 'stopped'
    assert report['stop_url'] is None
    assert report['retry_url'] == f'/reporting/jobs/{report_id}/retry'
    assert client.post(report['delete_url']).status_code == 200

    chart_id = app_module.repository.create_report_chart_job(
        technology='nsa', scope='single', dataset_ids={'data': [1], 'voice': [2], 'speech': [3]},
        dataset_names={'data': ['Data'], 'voice': ['Voice'], 'speech': ['Speech']},
        template_name='NSA Slide Template', created_by='admin',
    )
    app_module.repository.update_report_chart_job(chart_id, status='processing', progress=40)
    stopped_chart = client.post(f'/reporting/chart-jobs/{chart_id}/stop')
    assert stopped_chart.status_code == 200
    chart = next(item for item in client.get('/api/reporting/chart-jobs').json()['jobs'] if item['id'] == chart_id)
    assert chart['status'] == 'stopped'
    assert chart['stop_url'] is None
    assert chart['retry_url'] == f'/reporting/chart-jobs/{chart_id}/retry'
    assert client.post(chart['delete_url']).status_code == 200


def test_chart_set_selector_excludes_published_but_processing_job(client) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    chart_set = app_module.persist_report_charts(
        'NSA Slide Template', 'single',
        [({'slide': 1, 'title': 'Completed chart', 'source': 'data', 'chart_type': 'Bar'}, b'PNG')],
        {'data': 1, 'voice': 1, 'speech': 1},
    )
    assert (app_module.report_charts_directory() / chart_set['generation']).is_dir()
    assert chart_set['generation'].startswith(
        datetime.strptime(chart_set['generated_at'], '%Y-%m-%d %H:%M:%S').strftime('%Y%m%d-%H%M%S')
    )
    job_id = app_module.repository.create_report_chart_job(
        technology='nsa', scope='single', dataset_ids={'data': [1], 'voice': [2], 'speech': [3]},
        dataset_names={'data': ['Data'], 'voice': ['Voice'], 'speech': ['Speech']},
        template_name='NSA Slide Template', created_by='admin',
    )
    app_module.repository.update_report_chart_job(job_id, status='processing', generation=chart_set['generation'])
    selector_value = f'value="standalone:{chart_set["generation"]}"'
    assert selector_value not in client.get('/reporting').text

    app_module.repository.update_report_chart_job(job_id, status='ready', progress=100, finished=True)
    assert selector_value in client.get('/reporting').text


def test_persisted_chart_set_keeps_template_order_when_rendered_by_cdr_source(client) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    chart_set = app_module.persist_report_charts(
        'NSA Slide Template', 'single',
        [
            ({'order': 2, 'slide': 3, 'title': 'Speech chart', 'source': 'speech', 'chart_type': 'Bar'}, b'SPEECH'),
            ({'order': 0, 'slide': 1, 'title': 'Data chart', 'source': 'data', 'chart_type': 'Bar'}, b'DATA'),
            ({'order': 1, 'slide': 2, 'title': 'Voice chart', 'source': 'voice', 'chart_type': 'Bar'}, b'VOICE'),
        ],
        {'data': 1, 'voice': 1, 'speech': 1},
    )

    manifest = json.loads((app_module.report_charts_directory() / chart_set['generation'] / 'manifest.json').read_text(encoding='utf-8'))
    assert [chart['title'] for chart in manifest['charts']] == ['Data chart', 'Voice chart', 'Speech chart']
    assert [chart['file'] for chart in manifest['charts']] == ['chart-001.png', 'chart-002.png', 'chart-003.png']


def test_retrying_a_failed_chart_job_reuses_its_row(client) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    job_id = app_module.repository.create_report_chart_job(
        technology='nsa', scope='single', dataset_ids={'data': [], 'voice': [], 'speech': []},
        dataset_names={}, template_name='NSA Slide Template', created_by='admin',
    )
    app_module.repository.update_report_chart_job(job_id, status='failed', progress=100, last_error='Synthetic failure', finished=True)
    before_ids = [row['id'] for row in app_module.repository.list_report_chart_jobs(limit=None)]

    response = client.post(f'/reporting/chart-jobs/{job_id}/retry')

    assert response.status_code == 202
    assert response.json()['job_id'] == job_id
    assert [row['id'] for row in app_module.repository.list_report_chart_jobs(limit=None)] == before_ids


def test_deleting_a_ready_chart_job_removes_its_chart_set(client) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    chart_set = app_module.persist_report_charts(
        'NSA Slide Template', 'single',
        [({'slide': 1, 'title': 'Chart', 'source': 'data', 'chart_type': 'Bar'}, b'PNG')],
        {'data': 1, 'voice': 1, 'speech': 1},
    )
    job_id = app_module.repository.create_report_chart_job(
        technology='nsa', scope='single', dataset_ids={'data': [], 'voice': [], 'speech': []},
        dataset_names={}, template_name='NSA Slide Template', created_by='admin',
    )
    app_module.repository.update_report_chart_job(
        job_id, status='ready', progress=100, generation=chart_set['generation'], finished=True,
    )

    deleted = client.post(f'/reporting/chart-jobs/{job_id}/delete')

    assert deleted.status_code == 200
    assert deleted.json()['generation'] == chart_set['generation']
    assert not (app_module.report_charts_directory() / chart_set['generation']).exists()


def test_reporting_multivendor_requires_a_previously_mapped_selected_cdr(client) -> None:
    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    uploads = [
        ('NetCheck_CDR_Data.csv', 'data', b'RAT,Operator,Mean_Data_Rate\nENDC,Vodafone UK,42\n'),
        ('NetCheck_CDR_Voice.csv', 'voice', b'RAT_A,Operator,Call_Duration\nENDC,Vodafone UK,60\n'),
        ('NetCheck_CDR_Speech.csv', 'speech', b'Sample_RAT_A,Operator,LQ\nENDC,Vodafone UK,3.8\n'),
    ]
    for filename, dataset_kind, content in uploads:
        response = client.post(
            '/dashboard/upload',
            data={'dataset_kinds': dataset_kind},
            files={'dataset_files': (filename, BytesIO(content), 'text/csv')},
        )
        assert response.status_code == 200

    page = client.get('/reporting')
    assert page.status_code == 200
    assert 'data-vendor-mapped="false"' in page.text
    report = client.post('/reporting/netcheck-cdr', data={
        'data_dataset_id': 1, 'voice_dataset_id': 2, 'speech_dataset_id': 3,
        'technology': 'nsa', 'report_scope': 'multivendor',
    })
    assert report.status_code == 400
    assert 'requires every selected Data, Voice and Speech CDR to have a Workspace Vendor mapping' in report.text


def test_netcheck_reporting_generates_template_backed_pptx(client) -> None:
    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    uploads = [
        ('NetCheck_CDR_Data.csv', b'RAT,Operator,Mean_Data_Rate,Test_Result\nENDC,Vodafone UK,42,Success\n', 'text/csv'),
        ('NetCheck_CDR_Voice.csv', b'RAT_A,Operator,Call_Status,Call_Duration\nENDC,Vodafone UK,Completed,60\n', 'text/csv'),
        ('NetCheck_CDR_Speech.csv', b'Sample_RAT_A,Operator,LQ\nENDC,Vodafone UK,3.8\n', 'text/csv'),
    ]
    for filename, content, media_type in uploads:
        response = client.post('/dashboard/upload', files={'dataset_files': (filename, BytesIO(content), media_type)})
        assert response.status_code == 200

    report = client.post('/reporting/netcheck-cdr', data={
        'data_dataset_id': 1,
        'voice_dataset_id': 2,
        'speech_dataset_id': 3,
        'technology': 'nsa',
        'report_scope': 'single',
        'slides_templates': 'nsa:NSA Slide Template',
    })

    assert report.status_code == 202
    job = wait_for_report_job(client, report.json()['job_id'])
    assert job['status'] == 'ready'
    assert job['slides'] == 17
    assert re.search(r'NetCheck_CDR_NSA_operator-comparison_\d{8}-\d{6}\.pptx', job['report_name'])
    download = client.get(job['download_url'])
    assert download.status_code == 200
    assert download.headers['content-type'].startswith('application/vnd.openxmlformats-officedocument.presentationml.presentation')
    assert download.content[:2] == b'PK'
    opened = client.get(job['open_url'])
    assert opened.status_code == 200
    assert opened.headers['content-disposition'].startswith('inline;')
    import src.DashboardAnalytic as app_module
    stale_file = app_module._report_job_directory(job['report_name']) / 'stale-output.txt'
    stale_file.write_text('remove me', encoding='utf-8')
    relaunched = client.post(job['retry_url'])
    assert relaunched.status_code == 202
    assert relaunched.json()['job_id'] == job['id']
    rerun = wait_for_report_job(client, job['id'])
    assert rerun['status'] == 'ready'
    assert not stale_file.exists()
    assert [item['id'] for item in client.get('/api/reporting/jobs').json()['jobs']] == [job['id']]
    deleted = client.post(job['delete_url'])
    assert deleted.status_code == 200
    assert client.get(job['download_url']).status_code == 404


def test_reporting_generates_template_chart_previews(client, monkeypatch) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    uploads = [
        ('NetCheck_CDR_Data.csv', 'data', b'RAT,Operator,Mean_Data_Rate\nENDC,Vodafone UK,42\n'),
        ('NetCheck_CDR_Voice.csv', 'voice', b'RAT_A,Operator,Call_Status\nENDC,Vodafone UK,Completed\n'),
        ('NetCheck_CDR_Speech.csv', 'speech', b'Sample_RAT_A,Operator,LQ\nENDC,Vodafone UK,3.8\n'),
    ]
    for filename, kind, content in uploads:
        response = client.post(
            '/dashboard/upload', data={'dataset_kinds': kind},
            files={'dataset_files': (filename, BytesIO(content), 'text/csv')},
        )
        assert response.status_code == 200

    rendered: list[tuple[str, bool]] = []

    def render_preview(frame, entry, *, multivendor=False):
        rendered.append((entry.cdr_source, multivendor))
        return b'PNG'

    monkeypatch.setattr(app_module, 'render_catalog_chart_preview', render_preview)
    response = client.post('/reporting/netcheck-cdr/charts', data={
        'data_dataset_id': 1, 'voice_dataset_id': 2, 'speech_dataset_id': 3,
        'technology': 'nsa', 'report_scope': 'single', 'slides_templates': 'nsa:NSA Slide Template',
    })

    assert response.status_code == 202
    job = wait_for_report_chart_job(client, response.json()['job_id'])
    assert job['status'] == 'ready'
    payload = client.get(job['open_url']).json()
    assert payload['template'] == 'NSA Slide Template'
    assert payload['scope'] == 'single'
    assert payload['dataset_counts'] == {'data': 1, 'voice': 1, 'speech': 1}
    assert re.fullmatch(r'\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2}', payload['generated_at'])
    assert job['date'] == payload['generated_at']
    assert payload['generation'] == datetime.strptime(payload['generated_at'], '%Y-%m-%d %H:%M:%S').strftime('%Y%m%d-%H%M%S')
    assert payload['charts']
    preview_context = client.get('/api/reporting/chart-preview/context', params={
        'source': 'standalone', 'identifier': payload['generation'], 'chart_index': 0,
    })
    assert preview_context.status_code == 200
    context_payload = preview_context.json()
    assert context_payload['dataset_ids_by_source'] == {'cdr-data': ['1'], 'cdr-voice': ['2'], 'cdr-speech': ['3']}
    source_key = context_payload['cdr_source'].lower()
    expected_id = {'cdr-data': '1', 'cdr-voice': '2', 'cdr-speech': '3'}[source_key]
    assert context_payload['dataset_ids'] == [expected_id]
    assert context_payload['datasets_by_source']['cdr-data'] == [{'value': '1', 'label': 'NetCheck_CDR_Data.csv'}]
    assert context_payload['datasets_by_source']['cdr-voice'] == [{'value': '2', 'label': 'NetCheck_CDR_Voice.csv'}]
    wrong_id = next(value for value in ('1', '2', '3') if value != expected_id)
    invalid_dataset_type = client.post('/api/reporting/chart-preview', json={
        'source': 'standalone', 'identifier': payload['generation'], 'chart_index': 0,
        'definition': {'cdr_source': context_payload['cdr_source'], 'dataset_ids': [wrong_id]},
    })
    assert invalid_dataset_type.status_code == 400
    image_url = payload['charts'][0]['image_url']
    assert re.match(r'/reporting/charts/\d{8}-\d{6}/chart-\d+\.png\?v=', image_url)
    assert client.get(image_url).content == b'PNG'
    second = client.post('/reporting/netcheck-cdr/charts', data={
        'data_dataset_id': 1, 'voice_dataset_id': 2, 'speech_dataset_id': 3,
        'technology': 'nsa', 'report_scope': 'single', 'slides_templates': 'nsa:NSA Slide Template',
    })
    assert second.status_code == 202
    second_job = wait_for_report_chart_job(client, second.json()['job_id'])
    assert second_job['status'] == 'ready'
    second_payload = client.get(second_job['open_url']).json()
    assert second_payload['generation'] != payload['generation']
    assert client.get(f"/api/reporting/chart-sets/{payload['generation']}").status_code == 200
    deleted = client.post(f"/reporting/chart-sets/{second_payload['generation']}/delete")
    assert deleted.status_code == 200
    assert [item['generation'] for item in deleted.json()['chart_sets']] == [payload['generation']]
    assert client.get(image_url).content == b'PNG'
    page = client.get('/reporting')
    assert image_url in page.text
    assert 'Operator Comparison' in page.text
    assert '(Data:1 | Voice:1 | Speech:1)' in page.text
    assert 'Delete Selected Charts Set' in page.text
    assert 'Delete All Charts Sets' in page.text
    assert 'data-report-chart-viewer' in page.text
    assert 'data-report-chart-viewer-canvas' in page.text
    assert 'data-report-chart-zoom-reset' in page.text
    assert rendered and all(multivendor is False for _, multivendor in rendered)
    orphaned_job = app_module.repository.create_report_chart_job(
        technology='nsa', scope='single', dataset_ids={}, dataset_names={},
        template_name='Interrupted Chart Set', created_by='admin',
    )
    orphaned_directory = app_module.report_charts_directory() / '.incomplete-chart-set'
    orphaned_directory.mkdir(parents=True)
    (orphaned_directory / 'partial.png').write_bytes(b'partial')
    cleared = client.post('/reporting/chart-sets/delete-all')
    assert cleared.status_code == 200
    assert cleared.json()['chart_sets'] == []
    assert cleared.json()['deleted_jobs'] == 2
    assert app_module.repository.get_report_chart_job(orphaned_job) is None
    assert app_module.repository.list_report_chart_jobs(limit=None) == []
    assert list(app_module.report_charts_directory().iterdir()) == []
    assert client.get(f"/api/reporting/chart-sets/{payload['generation']}").status_code == 404


def test_temporary_preview_accepts_dataset_ids_with_legacy_multiplication_separator(monkeypatch) -> None:
    import src.DashboardAnalytic as app_module

    monkeypatch.setattr(app_module.repository, 'list_datasets', lambda: [
        {'id': 2, 'status': 'ready', 'dataset_kind': 'voice'},
        {'id': 5, 'status': 'ready', 'dataset_kind': 'voice'},
    ])

    assert app_module._temporary_preview_dataset_ids({'dataset_ids': '2 × 5'}, {}, 'voice') == [2, 5]


def test_report_chart_generation_failures_return_json_and_are_logged(client, monkeypatch) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    for filename, kind, content in (
        ('NetCheck_CDR_Data.csv', 'data', b'RAT,Operator,Mean_Data_Rate\nENDC,Vodafone UK,42\n'),
        ('NetCheck_CDR_Voice.csv', 'voice', b'RAT_A,Operator,Call_Status\nENDC,Vodafone UK,Completed\n'),
        ('NetCheck_CDR_Speech.csv', 'speech', b'Sample_RAT_A,Operator,LQ\nENDC,Vodafone UK,3.8\n'),
    ):
        upload = client.post('/dashboard/upload', data={'dataset_kinds': kind}, files={'dataset_files': (filename, BytesIO(content), 'text/csv')})
        assert upload.status_code == 200

    def fail_render(*_args, **_kwargs):
        raise RuntimeError('Synthetic renderer failure')

    monkeypatch.setattr(app_module, 'render_catalog_chart_preview', fail_render)
    response = client.post('/reporting/netcheck-cdr/charts', data={
        'data_dataset_id': 1, 'voice_dataset_id': 2, 'speech_dataset_id': 3,
        'technology': 'nsa', 'report_scope': 'single', 'slides_templates': 'nsa:NSA Slide Template',
    })

    assert response.status_code == 202
    job = wait_for_report_chart_job(client, response.json()['job_id'])
    assert job['status'] == 'failed'
    assert job['error'] == 'Synthetic renderer failure'
    log = next(row for row in app_module.repository.list_logs() if row['action'] == 'generate_report_chart_set_failed')
    assert 'Synthetic renderer failure' in log['details']
    app_log = next(row for row in app_module.build_app_logs() if row['action'] == 'generate_report_chart_set_failed')
    assert app_log['log_type'] == 'Error'


def test_reporting_concatenates_multiple_campaign_cdrs_per_source(client, monkeypatch) -> None:
    import src.DashboardAnalytic as app_module

    client.post('/login', data={'username': 'admin', 'password': 'admin123'}, follow_redirects=False)
    uploads = [
        ('data-q1.csv', 'data', b'RAT_A,Campaign,Operator,Data_Q1\nENDC,2026 Q1,EE,10\n'),
        ('data-q2.csv', 'data', b'RAT_A,Campaign,Operator,Data_Q2\nENDC,2026 Q2,EE,20\n'),
        ('voice-q2.csv', 'voice', b'RAT_A,Campaign,Operator,Call_Status\nENDC,2026 Q2,EE,Completed\n'),
        ('speech-q2.csv', 'speech', b'RAT_A,Campaign,Operator,LQ\nENDC,2026 Q2,EE,3.8\n'),
    ]
    for filename, kind, content in uploads:
        response = client.post(
            '/dashboard/upload',
            data={'dataset_kinds': kind},
            files={'dataset_files': (filename, BytesIO(content), 'text/csv')},
        )
        assert response.status_code == 200

    captured: dict[str, pd.DataFrame] = {}

    def capture_report(destination, template, frames, technology, multivendor, catalog, **kwargs):
        frame_loader = kwargs.get('frame_loader')
        captured.update(frames or {kind: frame_loader(kind) for kind in ('data', 'voice', 'speech')})
        destination.write_bytes(b'PK')
        return destination

    monkeypatch.setattr(app_module, 'render_cdr_report', capture_report)
    payload = urlencode([
        ('data_dataset_id', '1'), ('data_dataset_id', '2'),
        ('voice_dataset_id', '3'), ('speech_dataset_id', '4'),
            ('technology', 'nsa'), ('report_scope', 'single'), ('slides_templates', 'nsa:NSA Slide Template'),
    ])
    response = client.post(
        '/reporting/netcheck-cdr',
        content=payload,
        headers={'content-type': 'application/x-www-form-urlencoded'},
    )

    assert response.status_code == 202
    job = wait_for_report_job(client, response.json()['job_id'])
    assert job['status'] == 'ready'
    assert captured['data']['Campaign'].tolist() == ['2026 Q1', '2026 Q2']
    # The shared CDR table and renderer materialise only fields required by
    # the chosen Slides Template; unrelated source metrics stay individual.
    assert not {'Data_Q1', 'Data_Q2'}.intersection(app_module.repository.list_reporting_row_columns('data'))
